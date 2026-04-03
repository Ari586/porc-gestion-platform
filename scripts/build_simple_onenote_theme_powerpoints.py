from __future__ import annotations

import argparse
import json
import re
import unicodedata
from collections import defaultdict
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


EXPORT_ROOT = Path("/Users/arielhavana/antigr/porc/mockups/onenote_mac_export/GeP_23-10_S-Notizbuch")
OUT_DIR = Path("/Users/arielhavana/antigr/porc/mockups/simple_onenote_theme_powerpoints")
MANIFEST_PATH = OUT_DIR / "manifest.json"

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(248, 250, 252)
TITLE = RGBColor(26, 43, 60)
TEXT = RGBColor(58, 76, 94)
ACCENT = RGBColor(45, 106, 129)
LINE = RGBColor(210, 220, 229)
MUTED = RGBColor(97, 117, 141)

NOISE_PREFIXES = (
    "OneNote type:",
    "Hierarchie:",
    "Kinder:",
    "Datum:",
    "Quelle:",
    "bearbeitet/Datum:",
    "erstellt/Datum:",
    "Seite ",
    "Modifiziert nach:",
    "abgerufen am",
)

NOISE_CONTAINS = (
    "qr-code scannen",
    "kopfhorer an",
    "kopfh rer an",
    "seite 1 von",
    "seite 2 von",
    "seite 3 von",
)


def normalize_key(value: str) -> str:
    value = unicodedata.normalize("NFKD", value).encode("ascii", "ignore").decode("ascii")
    value = value.lower()
    value = re.sub(r"[^a-z0-9]+", " ", value)
    return re.sub(r"\s+", " ", value).strip()


def canonical_bhf(name: str) -> str | None:
    match = re.match(r"^BHF_(\d+)([ab])?(?:_|$)", name, flags=re.IGNORECASE)
    if not match:
        return None
    number = int(match.group(1))
    if not 1 <= number <= 24:
        return None
    return f"{number:02d}"


def slugify(value: str) -> str:
    value = unicodedata.normalize("NFKD", value).encode("ascii", "ignore").decode("ascii")
    value = re.sub(r"[^A-Za-z0-9._-]+", "_", value).strip("._-")
    return value or "untitled"


def strip_leading_symbols(value: str) -> str:
    return re.sub(r"^[^A-Za-z0-9]+", "", value).strip()


def clean_label(value: str) -> str:
    value = value.replace(".pptx", "")
    value = value.replace(".docx", "")
    value = value.replace("_", " ")
    value = value.replace("  ", " ")
    value = strip_leading_symbols(value)
    value = re.sub(r"^\d+[a-z]?(?:\.\d+)*\s*", "", value, flags=re.IGNORECASE)
    value = re.sub(r"\s+", " ", value).strip(" -._")
    return value or "Untitled"


def clean_space(value: str) -> str:
    value = value.replace("\xa0", " ")
    value = value.replace("•", " ")
    value = value.replace("\u2028", " ")
    return re.sub(r"\s+", " ", value).strip()


def trim_line(value: str, limit: int = 108) -> str:
    value = clean_space(value)
    if len(value) <= limit:
        return value
    return f"{value[: limit - 3].rstrip()}..."


def read_heading(path: Path) -> str | None:
    try:
        with path.open("r", encoding="utf-8") as handle:
            for line in handle:
                if line.startswith("# "):
                    heading = line[2:].strip()
                    heading = heading.replace("*", "")
                    heading = strip_leading_symbols(heading)
                    return heading or None
    except OSError:
        return None
    return None


def should_skip_theme(label: str) -> bool:
    key = normalize_key(label)
    return any(token in key for token in ["ubersicht", "tagesfazit"])


def theme_sort_key(value: str) -> tuple[int, int, str]:
    value = value.replace("_", " ")
    match = re.match(r"^(\d+)(?:\.(\d+))?", value)
    if not match:
        return (999, 999, normalize_key(value))
    major = int(match.group(1))
    minor = int(match.group(2) or 0)
    return (major, minor, normalize_key(value))


def dedupe(values: Iterable[str]) -> list[str]:
    result: list[str] = []
    seen: set[str] = set()
    for value in values:
        key = normalize_key(value)
        if not key or key in seen:
            continue
        seen.add(key)
        result.append(value)
    return result


def is_noise_line(line: str) -> bool:
    lowered = normalize_key(line)
    if not lowered:
        return True
    if lowered == "kinder":
        return True
    if any(lowered.startswith(normalize_key(prefix)) for prefix in NOISE_PREFIXES):
        return True
    if any(fragment in lowered for fragment in NOISE_CONTAINS):
        return True
    if lowered.startswith("http"):
        return True
    if lowered.startswith("www"):
        return True
    if re.fullmatch(r"\d+", lowered):
        return True
    return False


def body_lines(md_path: Path, page_title: str, theme_label: str, limit: int = 4) -> list[str]:
    try:
        text = md_path.read_text(encoding="utf-8", errors="ignore")
    except OSError:
        return []

    lines: list[str] = []
    for raw in text.splitlines():
        line = clean_space(raw)
        if not line:
            continue
        if raw.startswith("# ") or raw.startswith("## "):
            continue
        line = re.sub(r"^[-*]\s*", "", line)
        line = clean_space(line)
        if not line or is_noise_line(line):
            continue
        if normalize_key(line) in {normalize_key(page_title), normalize_key(theme_label)}:
            continue
        if len(line.split()) < 3:
            continue
        lines.append(trim_line(line))
        if len(lines) >= limit:
            break
    return dedupe(lines)


def theme_page_bundle(theme_dir: Path, theme_label: str) -> dict[str, list[str]]:
    subtopics: list[str] = []
    highlights: list[str] = []
    for md_path in sorted(theme_dir.glob("*.md")):
        if md_path.name == "README.md":
            continue
        page_title = read_heading(md_path) or clean_label(md_path.stem)
        if normalize_key(page_title).startswith("untitled"):
            continue
        subtopics.append(page_title)
        page_highlights = body_lines(md_path, page_title, theme_label)
        if page_highlights:
            highlights.append(f"{page_title}: {page_highlights[0]}")
        else:
            highlights.append(page_title)
    return {
        "subtopics": dedupe(subtopics),
        "highlights": dedupe(highlights),
    }


def group_source_dirs() -> dict[str, list[Path]]:
    grouped: dict[str, list[Path]] = defaultdict(list)
    for path in sorted(EXPORT_ROOT.iterdir()):
        if not path.is_dir():
            continue
        code = canonical_bhf(path.name)
        if code:
            grouped[code].append(path)
    return grouped


def collect_bhf_themes(source_dirs: list[Path]) -> list[dict]:
    themes: dict[str, dict] = {}
    for source_dir in source_dirs:
        for child in sorted(source_dir.iterdir()):
            if not child.is_dir():
                continue
            label = clean_label(child.name)
            if should_skip_theme(label):
                continue
            key = normalize_key(label)
            entry = themes.setdefault(
                key,
                {
                    "label": label,
                    "sort_name": child.name,
                    "sources": [],
                    "theme_dirs": [],
                    "subtopics": [],
                    "highlights": [],
                },
            )
            entry["sources"].append(source_dir.name)
            entry["theme_dirs"].append(child)

    ordered = sorted(
        themes.values(),
        key=lambda item: (theme_sort_key(item["sort_name"]), normalize_key(item["label"])),
    )
    for item in ordered:
        subtopics: list[str] = []
        highlights: list[str] = []
        for theme_dir in item["theme_dirs"]:
            bundle = theme_page_bundle(theme_dir, item["label"])
            subtopics.extend(bundle["subtopics"])
            highlights.extend(bundle["highlights"])
        item["sources"] = sorted(set(item["sources"]))
        item["theme_dirs"] = [str(path) for path in item["theme_dirs"]]
        item["subtopics"] = dedupe(subtopics)[:18]
        item["highlights"] = dedupe(highlights)[:18]
    return ordered


def add_title(slide, text: str, top: float = 0.55, size: int = 28) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = TITLE


def add_subtitle(slide, text: str, top: float) -> None:
    box = slide.shapes.add_textbox(Inches(0.72), Inches(top), Inches(11.2), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(12)
    run.font.color.rgb = TEXT


def add_note(slide, text: str, top: float, size: int = 11) -> None:
    box = slide.shapes.add_textbox(Inches(0.72), Inches(top), Inches(11.3), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], top: float = 1.65, left: float = 0.9, width: float = 11.3, size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.1))
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"- {item}"
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = TEXT
        paragraph.space_after = Pt(9)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(7.02), Inches(11.5), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.color.rgb = TEXT


def add_rule(slide, top: float = 1.25) -> None:
    shape = slide.shapes.add_shape(1, Inches(0.72), Inches(top), Inches(11.6), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LINE
    shape.line.fill.background()


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def chunked(items: list[str], chunk_size: int) -> list[list[str]]:
    return [items[index : index + chunk_size] for index in range(0, len(items), chunk_size)] or [[]]


def learning_focus(theme: dict) -> list[str]:
    prompts = [
        f"Definiere '{theme['label']}' zuerst in 2 bis 3 Saetzen.",
        "Ordne danach die wichtigsten Unterthemen logisch und in eigenen Worten.",
        "Verbinde das Thema mit einem konkreten Beispiel aus der Pflegepraxis.",
        "Denke bei Antworten an Beobachtung, Ziele, Massnahmen und Evaluation.",
    ]
    if theme["subtopics"]:
        prompts.append(f"Sicher frei nennen koennen: {theme['subtopics'][0]}")
    return prompts[:5]


def build_theme_deck(code: str, theme: dict, deck_no: int, deck_total: int) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    prs.core_properties.title = theme["label"]
    prs.core_properties.subject = f"BHF {code}"
    prs.core_properties.author = "Codex"

    slide = new_slide(prs)
    add_title(slide, theme["label"])
    add_subtitle(slide, f"BHF {code} | Thema {deck_no} von {deck_total}", 1.15)
    add_rule(slide)
    info_lines = [
        f"{len(theme['subtopics'])} Unterthemen oder Materialien zusammengefuehrt",
        f"Quellen: {', '.join(theme['sources'][:4])}",
        "Erstellt aus dem lokalen OneNote-Mac-Export",
    ]
    add_bullets(slide, info_lines, top=1.95, left=1.0, width=10.8)
    add_footer(slide, "Codex | Themen-Deck")

    for page_no, chunk in enumerate(chunked(theme["subtopics"], 8), 1):
        slide = new_slide(prs)
        add_title(slide, f"{theme['label']} - Unterthemen ({page_no})")
        add_subtitle(slide, "Seiten, Materialien und Teilaspekte aus den OneNote-Notizen", 1.15)
        add_rule(slide)
        add_bullets(slide, chunk, top=1.75, left=1.0, width=10.8, size=18)
        add_footer(slide, f"BHF {code}")

    for page_no, chunk in enumerate(chunked(theme["highlights"], 6), 1):
        slide = new_slide(prs)
        add_title(slide, f"{theme['label']} - Kerngedanken ({page_no})")
        add_subtitle(slide, "Kurzfassung aus den lokal auslesbaren Notizseiten", 1.15)
        add_rule(slide)
        add_bullets(slide, chunk, top=1.75, left=1.0, width=10.8, size=17)
        add_note(slide, "Diese Stichpunkte sind bewusst knapp gehalten und eignen sich gut fuer Wiederholung oder Pruefungsvorbereitung.", 6.55)
        add_footer(slide, f"BHF {code}")

    slide = new_slide(prs)
    add_title(slide, f"{theme['label']} - Lernfokus")
    add_subtitle(slide, "So kannst du das Thema muendlich oder schriftlich sicher wiedergeben", 1.15)
    add_rule(slide)
    add_bullets(slide, learning_focus(theme), top=1.75, left=1.0, width=10.8, size=18)
    add_footer(slide, f"BHF {code}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    filename = OUT_DIR / f"BHF_{code}_{slugify(theme['label'])}.pptx"
    prs.save(filename)
    return filename


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Create one simple German PowerPoint per OneNote theme.")
    parser.add_argument("--bhf", nargs="*", help="Optional list of BHF numbers to build, e.g. 01 02 03")
    return parser.parse_args()


def wanted_codes(selected: list[str] | None) -> set[str] | None:
    if not selected:
        return None
    normalized = set()
    for raw in selected:
        raw = raw.strip()
        if not raw:
            continue
        normalized.add(f"{int(raw):02d}")
    return normalized


def main() -> None:
    args = parse_args()
    selected = wanted_codes(args.bhf)
    grouped = group_source_dirs()
    manifest: list[dict] = []
    for number in range(1, 25):
        code = f"{number:02d}"
        if selected and code not in selected:
            continue
        source_dirs = grouped.get(code, [])
        if not source_dirs:
            continue
        themes = collect_bhf_themes(source_dirs)
        for deck_no, theme in enumerate(themes, start=1):
            path = build_theme_deck(code, theme, deck_no, len(themes))
            manifest.append(
                {
                    "bhf": code,
                    "theme": theme["label"],
                    "deck": str(path),
                    "source_dirs": [str(item) for item in source_dirs],
                    "theme_dirs": theme["theme_dirs"],
                    "source_names": theme["sources"],
                    "subtopic_count": len(theme["subtopics"]),
                    "subtopics": theme["subtopics"],
                    "highlights": theme["highlights"],
                }
            )
            print(f"saved {path}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    MANIFEST_PATH.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"saved {MANIFEST_PATH}")


if __name__ == "__main__":
    main()
