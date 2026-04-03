from __future__ import annotations

import json
import re
import unicodedata
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


EXPORT_ROOT = Path("/Users/arielhavana/antigr/porc/mockups/onenote_mac_export/GeP_23-10_S-Notizbuch")
OUT_DIR = Path("/Users/arielhavana/antigr/porc/mockups/simple_onenote_powerpoints")
MANIFEST_PATH = OUT_DIR / "manifest.json"

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(248, 250, 252)
TITLE = RGBColor(26, 43, 60)
TEXT = RGBColor(58, 76, 94)
ACCENT = RGBColor(45, 106, 129)
LINE = RGBColor(210, 220, 229)


def normalize_key(value: str) -> str:
    value = unicodedata.normalize("NFKD", value).encode("ascii", "ignore").decode("ascii")
    value = value.lower()
    value = re.sub(r"[^a-z0-9]+", " ", value)
    return re.sub(r"\s+", " ", value).strip()


def canonical_bhf(name: str) -> str | None:
    match = re.match(r"^BHF_(\d+)([ab])?(?:_|$)", name, flags=re.IGNORECASE)
    if not match:
        return None
    number = int(match.group(1))
    if not 1 <= number <= 24:
        return None
    return f"{number:02d}"


def slugify(value: str) -> str:
    value = unicodedata.normalize("NFKD", value).encode("ascii", "ignore").decode("ascii")
    value = re.sub(r"[^A-Za-z0-9._-]+", "_", value).strip("._-")
    return value or "untitled"


def strip_leading_symbols(value: str) -> str:
    return re.sub(r"^[^A-Za-z0-9]+", "", value).strip()


def clean_label(value: str) -> str:
    value = value.replace(".pptx", "")
    value = value.replace(".docx", "")
    value = value.replace("_", " ")
    value = value.replace("  ", " ")
    value = strip_leading_symbols(value)
    value = re.sub(r"^\d+[a-z]?(?:\.\d+)*\s*", "", value, flags=re.IGNORECASE)
    value = re.sub(r"\s+", " ", value).strip(" -._")
    return value or "Untitled"


def read_heading(path: Path) -> str | None:
    try:
        with path.open("r", encoding="utf-8") as handle:
            for line in handle:
                if line.startswith("# "):
                    heading = line[2:].strip()
                    heading = strip_leading_symbols(heading)
                    return heading or None
    except OSError:
        return None
    return None


def should_skip_theme(label: str) -> bool:
    key = normalize_key(label)
    return any(token in key for token in ["ubersicht", "tagesfazit"])


def page_titles(theme_dir: Path, limit: int = 6) -> list[str]:
    titles: list[str] = []
    for md_path in sorted(theme_dir.glob("*.md")):
        if md_path.name == "README.md":
            continue
        heading = read_heading(md_path) or clean_label(md_path.stem)
        if normalize_key(heading).startswith("untitled"):
            continue
        titles.append(heading)
    deduped: list[str] = []
    seen: set[str] = set()
    for title in titles:
        key = normalize_key(title)
        if not key or key in seen:
            continue
        seen.add(key)
        deduped.append(title)
    return deduped[:limit]


def group_source_dirs() -> dict[str, list[Path]]:
    grouped: dict[str, list[Path]] = defaultdict(list)
    for path in sorted(EXPORT_ROOT.iterdir()):
        if not path.is_dir():
            continue
        code = canonical_bhf(path.name)
        if code:
            grouped[code].append(path)
    return grouped


def collect_bhf_themes(source_dirs: list[Path]) -> list[dict]:
    themes: dict[str, dict] = {}
    for source_dir in source_dirs:
        for child in sorted(source_dir.iterdir()):
            if not child.is_dir():
                continue
            label = clean_label(child.name)
            if should_skip_theme(label):
                continue
            key = normalize_key(label)
            entry = themes.setdefault(
                key,
                {
                    "label": label,
                    "sort_name": child.name,
                    "bullets": [],
                    "sources": [],
                },
            )
            entry["sources"].append(source_dir.name)
            for bullet in page_titles(child):
                if normalize_key(bullet) not in {normalize_key(item) for item in entry["bullets"]}:
                    entry["bullets"].append(bullet)
    ordered = sorted(
        themes.values(),
        key=lambda item: (theme_sort_key(item["sort_name"]), normalize_key(item["label"])),
    )
    for item in ordered:
        item["sources"] = sorted(set(item["sources"]))
        item["bullets"] = item["bullets"][:6]
    return ordered


def theme_sort_key(value: str) -> tuple[int, int, str]:
    value = value.replace("_", " ")
    match = re.match(r"^(\d+)(?:\.(\d+))?", value)
    if not match:
        return (999, 999, normalize_key(value))
    major = int(match.group(1))
    minor = int(match.group(2) or 0)
    return (major, minor, normalize_key(value))


def add_title(slide, text: str, top: float = 0.55, size: int = 28) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = TITLE


def add_subtitle(slide, text: str, top: float) -> None:
    box = slide.shapes.add_textbox(Inches(0.72), Inches(top), Inches(11.2), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(12)
    run.font.color.rgb = TEXT


def add_bullets(slide, items: list[str], top: float = 1.65, left: float = 0.9, width: float = 11.3) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.1))
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"- {item}"
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = TEXT
        paragraph.space_after = Pt(10)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.75), Inches(7.02), Inches(11.5), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.color.rgb = TEXT


def add_rule(slide, top: float = 1.25) -> None:
    shape = slide.shapes.add_shape(1, Inches(0.72), Inches(top), Inches(11.6), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LINE
    shape.line.fill.background()


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def overview_chunks(items: list[str], chunk_size: int = 10) -> list[list[str]]:
    return [items[index : index + chunk_size] for index in range(0, len(items), chunk_size)] or [[]]


def build_deck(code: str, themes: list[dict], source_dirs: list[Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide = new_slide(prs)
    add_title(slide, f"BHF {code}")
    add_subtitle(slide, "Simple PowerPoint aus dem Mac OneNote Export", 1.15)
    add_rule(slide)
    infos = [
        f"{len(themes)} Themen zusammengefuehrt",
        f"{len(source_dirs)} Quellordner zusammengefuehrt",
        "Einfaches Format fuer die schnelle Wiederholung",
    ]
    add_bullets(slide, infos, top=2.0, left=1.0, width=10.8)
    add_footer(slide, "Codex | OneNote-Mac-Export")

    overview_items = [theme["label"] for theme in themes]
    for page_no, chunk in enumerate(overview_chunks(overview_items), 1):
        slide = new_slide(prs)
        add_title(slide, f"BHF {code} - Themen ({page_no})")
        add_subtitle(slide, "Uebersicht der zusammengefuehrten Themen", 1.15)
        add_rule(slide)
        add_bullets(slide, chunk, top=1.75, left=1.0, width=10.8)
        add_footer(slide, "Themen")

    for theme in themes:
        slide = new_slide(prs)
        add_title(slide, theme["label"])
        add_subtitle(slide, f"Quellen: {', '.join(theme['sources'][:3])}", 1.15)
        add_rule(slide)
        bullets = theme["bullets"] or ["In OneNote vorhanden, aber ohne lokal auslesbare Seiteninhalte."]
        add_bullets(slide, bullets, top=1.75, left=1.0, width=10.8)
        add_footer(slide, f"BHF {code}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    filename = OUT_DIR / f"BHF_{code}_simple.pptx"
    prs.save(filename)
    return filename


def main() -> None:
    grouped = group_source_dirs()
    manifest: list[dict] = []
    for number in range(1, 25):
        code = f"{number:02d}"
        source_dirs = grouped.get(code, [])
        if not source_dirs:
            continue
        themes = collect_bhf_themes(source_dirs)
        path = build_deck(code, themes, source_dirs)
        manifest.append(
            {
                "bhf": code,
                "deck": str(path),
                "sources": [str(item) for item in source_dirs],
                "theme_count": len(themes),
                "themes": [item["label"] for item in themes],
            }
        )
        print(f"saved {path}")

    MANIFEST_PATH.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"saved {MANIFEST_PATH}")


if __name__ == "__main__":
    main()
