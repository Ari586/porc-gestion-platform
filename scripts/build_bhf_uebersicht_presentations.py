from __future__ import annotations

import json
import re
import sqlite3
from collections import defaultdict
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DB_PATH = Path(
    "/Users/arielhavana/Library/Containers/com.microsoft.onenote.mac/Data/Library/Application Support/"
    "Microsoft User Data/OneNote/15.0/FullTextSearchIndex/{F5FF1CE4-251B-4B4C-AC34-A9D31CD459EA}{30}.db"
)
OUT_DIR = Path("/Users/arielhavana/antigr/porc/mockups/exam_prep_uebersicht_clean")
JSON_PATH = OUT_DIR / "bhf_uebersicht_data.json"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


WHITE = rgb("FFFFFF")
INK = rgb("223247")
TEXT = rgb("32465B")
MUTED = rgb("61758D")
LINE = rgb("D5E1E9")
LIGHT = rgb("F8FBFD")

TONES = {
    "sky": rgb("E7F3FA"),
    "mint": rgb("E7F4EE"),
    "sand": rgb("FCF3DA"),
    "rose": rgb("F9E6EA"),
    "white": WHITE,
}

ACCENTS = [
    ("1E7A8C", "175E6D"),
    ("2F6F9F", "214F72"),
    ("7C5D9A", "604776"),
    ("4F8A53", "37673A"),
    ("AF6B1F", "884F12"),
    ("8B4E8E", "67386A"),
]

BHF_ORDER = [f"{number:02d}" for number in range(1, 25)]

SECTION_GROUPS = {
    "01": ["1a", "1b"],
    "02": ["2"],
    "03": ["3a", "3b"],
    "04": ["4"],
    "05": ["5"],
    "06": ["6"],
    "07": ["7"],
    "08": ["8"],
    "09": ["9"],
    "10": ["10"],
    "11": ["11"],
    "12": ["12"],
    "13": ["13"],
    "14": ["14"],
    "15": ["15"],
    "16": ["16"],
    "17": ["17"],
    "18": ["18"],
    "19": ["19"],
    "20": ["20"],
    "21": ["21"],
    "22": ["22"],
    "23": ["23"],
    "24": ["24"],
}

TOPIC_OVERRIDES = {
    "01": [
        "Kontakt treten und Orientierung",
        "Grundlagen der Kommunikation",
        "Privat- und Intimsphaere",
        "Rechte und Pflichten",
        "Pflegeprozess",
        "Informationsgespraeche und Dokumentation",
        "Berufsrolle und Reflexion",
        "System Familie und kollegiale Beratung",
    ],
    "02": [
        "Bewegung und Bewegungsfoerderung",
        "Pflege bei Bewegungseinschraenkung",
        "Sturzprophylaxe und Unfallverhuetung",
        "Kontraktur und Osteoporose",
        "Orientierung und Anleitung",
        "Persoenliche Gesunderhaltung",
        "Haftungsrecht",
    ],
    "03": [
        "Pflegealltag und Praxiseinsatz reflektieren",
        "Hygienisches Handeln",
        "Haut und Hautzustaende",
        "Koerperpflege und Intimsphaere",
        "Ernaehrung und Ausscheidung",
        "Pflegeprozess vertiefen",
        "Mangelernaehrung und Dehydratation",
        "Saeuglinge pflegen",
        "Pflegeethik und Kommunikation",
        "Kultursensibilitaet und Familienorientierung",
    ],
    "04": [
        "Einstieg, Anatomie und Physiologie",
        "Schwangerschaft und Geburt",
        "Versorgung des Neugeborenen",
        "Unterstuetzungssysteme",
    ],
    "05": [
        "Akute Wunden",
        "Organisatorische und oekonomische Rahmenbedingungen",
        "Akutes Schmerzmanagement",
        "Kinder im Krankenhaus",
        "Krankenhaushygiene",
        "Prae- und postoperative Pflege Erwachsener",
    ],
    "06": [
        "Vitalparameter und ihre Aussagekraft",
        "Erstmassnahmen in Notfallsituationen",
        "Arbeiten in der Notaufnahme",
        "Blutentnahme, Blutwerte und EKG",
        "Reanimation nach ERC",
        "Rechtliche Grundlagen der Notfallversorgung",
        "Erstmassnahmen bei akutem Abdomen, Luftnot, Schlaganfall, Verbrennung und Schock",
    ],
    "07": [
        "Arbeiten in der ambulanten Pflege",
        "Menschen mit Diabetes mellitus versorgen",
        "Menschen mit Hoer- und Seheinschraenkungen versorgen",
        "Menschen nach ambulanter OP versorgen",
        "Gesundheitsfoerderung und Praevention im haeuslichen Umfeld",
    ],
    "08": [
        "Henriette Schulz und multimorbide Patienten",
        "Kind mit Asthma",
        "Gesundheitsfoerderung, Empowerment und Adhaerenz",
        "Kollegiale Fallberatung",
        "Krankenhausfinanzierung",
    ],
    "09": [
        "Versorgung nach einem neurologischen Akutereignis",
        "Schlaganfall pflegerisch versorgen",
        "Pflegeprozess in der rehabilitativen Pflege",
        "Menschen mit minimalem Bewusstseinszustand",
        "Evidenzbasiert pflegen",
    ],
    "10": [
        "Organisation in der stationaeren Langzeitpflege",
        "Pflegeprozess in der Langzeitpflege",
        "Hygienisches Handeln und Pandemieplaene",
        "Existentielle Erfahrungen und Gewalt",
        "Gewalt erkennen und vermeiden",
        "Professionelles Selbstverstaendnis",
    ],
    "11": [
        "Diabetes Typ 1",
        "ADHS",
        "EACH-Charta",
        "Familien mit chronisch kranken Kindern",
        "Kindgerechte Kommunikation und Beteiligung",
    ],
    "12": [
        "Frau Bahde baut ab",
        "Betreuungskonzepte und rechtliche Vertretung",
        "Wohnformen und Versorgungswechsel",
        "Freiheitsentziehende Massnahmen",
        "Demenz und Differenzialdiagnostik",
    ],
    "13": [
        "Chronisch krank sein",
        "Morbus Parkinson",
        "Multiple Sklerose und Rheuma",
        "Trajektmodell",
        "Expertenstandards bei chronischen Schmerzen und Mobilitaet",
        "Autonomie unterstuetzen",
        "Integrierte Versorgungsprozesse",
        "Belastungen der Pflegenden",
    ],
    "14": [
        "Bedeutung der Diagnose und Handlungssituation",
        "Gutartige und boesartige Tumoren",
        "Diagnostik und Therapie",
        "Pflegerische Begleitung",
        "Kommunikation mit Betroffenen und Angehoerigen",
        "Pflege und eigenes Erleben",
        "Tod und Sterben",
    ],
    "15": [
        "Berufskrankheiten",
        "Burnout",
        "Mobbing",
        "Cybermobbing und Hate Speech",
        "Sucht und Suchtpraevention",
        "Betriebliche Gesundheitsfoerderung",
    ],
    "16": [
        "Geschichte der Psychiatrie",
        "Interventionen in der Psychiatrie",
        "Depression",
        "Schizophrenie",
        "Anorexia nervosa",
        "Deeskalation, Recovery und Empowerment",
    ],
    "17": [
        "Mit herausfordernden Pflegesituationen umgehen",
        "Infektionsschutz und Infektionskrankheiten",
        "Gesundheit erhalten und foerdern",
    ],
    "18": [
        "Frakturen allgemein",
        "Unfallgeschehen und Erstversorgung",
        "Arbeiten auf der ITS",
        "Polytrauma pflegen",
        "SHT und Frakturen pflegen",
        "Informierte Entscheidung treffen",
    ],
    "19": [
        "Nach Schlaganfall rehabilitativ pflegen",
        "Nach der Reha den Pflegeprozess steuern",
        "Mit Aphasikern richtig kommunizieren",
        "Aphasieformen: amnestisch, Broca, Wernicke, global",
        "Cerebralparese und Rehabilitation",
    ],
    "20": [
        "Anatomie, Physiologie und Erkrankungen des Urogenitalsystems",
        "Dialysepatienten pflegen und beraten",
        "Menschen mit Kontinenzproblemen pflegen",
        "Qualitaetssichernd handeln",
        "Verschiedene Perspektiven entwickeln",
    ],
    "21": [
        "Pflege eines Fruehgeborenen",
        "Schreibaby",
        "Pflege von Kindern mit Darmerkrankungen",
        "Pflege von Kindern mit paediatrischen Erkrankungen",
        "Friedemann",
        "Meningitis",
        "Akute Pyelonephritis",
        "Verdacht auf SHT",
        "BNS-Kraempfe",
    ],
    "22": [
        "Zukunftswerkstatt Pflege 2040",
        "Methode Zukunftswerkstatt",
        "Kritik- und Problemloesungsphase",
        "Utopiephase",
        "KI und Robotik",
        "Neue Versorgungsmodelle",
    ],
    "23": [
        "Koerperliche Misshandlung bei Kindern",
        "Vernachlaessigung bei Kindern",
        "Sexueller Missbrauch bei Kindern",
        "Muenchhausen-Stellvertreter-Syndrom",
        "Gewaltsituationen in der Pflege erkennen",
        "Konflikte erkennen und Loesungen ableiten",
        "Gewalt in der Pflege praeventiv begegnen",
    ],
    "24": [
        "Einfuehrung",
        "Falluebersicht",
        "Fall 1 Frau Weiss",
        "Fall 2 Frau Feldmann",
        "Fall 3 Herr Gaertner",
        "Fall 4 Lara",
        "Fall 5 Studie",
        "Fall 6 Frau Falk",
        "Fall 7 Leon",
        "Fall 8 Frau Wunsch",
        "Fall 9 Lena",
        "Fall 10 Nadine",
        "SIS",
    ],
}

KEYWORD_HINTS = {
    "diabetes": [
        "DURST: Definition, Typ/Unterform, Ursachen und Risikofaktoren klar nennen.",
        "Typische Zeichen, Akutkomplikationen und Pflegeschwerpunkte verbinden.",
    ],
    "adhs": [
        "Formen unterscheiden: unaufmerksam, hyperaktiv-impulsiv, kombiniert.",
        "Pflege: klare Struktur, kurze Anweisungen, Elternarbeit und Beobachtung.",
    ],
    "pneumonie": [
        "Definition, Risikofaktoren und Verlauf sicher erklaeren.",
        "Pneumonieprophylaxe, Beobachtung und Atemunterstuetzung begruenden.",
    ],
    "copd": [
        "Obstruktion, Exazerbation und Langzeitverlauf benennen.",
        "Atemerleichternde Positionen, Inhalation und Schulung mitdenken.",
    ],
    "herzinsuffizienz": [
        "Links-/Rechtsherzinsuffizienz und typische Symptome unterscheiden.",
        "Oedeme, Dyspnoe, Gewichtskontrolle und Medikamentenwirkung verknuepfen.",
    ],
    "asthma": [
        "Form, Ausloeser, Peak-Flow und Notfallverhalten nennen.",
        "Inhalation, Triggerkontrolle und Eltern-/Kinderschulung einbauen.",
    ],
    "dehydratation": [
        "Definition, Ursachen, Risikofaktoren und Warnzeichen nennen.",
        "Trinkmanagement, Beobachtung und Exsikkoseprophylaxe begruenden.",
    ],
    "malnutrition": [
        "Screening, Risikofaktoren und Konsequenzen unterscheiden.",
        "Ernaehrungsmanagement, PEG/PEJ und Aspirationsprophylaxe einordnen.",
    ],
    "parkinson": [
        "Definition, Leitsymptome und Verlaufsbesonderheiten nennen.",
        "Medikation, Mobilitaet, Schlucken und Alltagsstruktur mitdenken.",
    ],
    "multiple sklerose": [
        "Autoimmunprozess, Verlaufsformen und typische Symptome erklaeren.",
        "Fatigue, Schubbeobachtung und Ressourcenfoerderung einbauen.",
    ],
    "arthritis": [
        "Entzuendung, Gelenkveraenderungen und Funktionseinbussen erklaeren.",
        "Schmerz, Bewegung, Gelenkschutz und Therapieadharenz nennen.",
    ],
    "depression": [
        "Definition, Schweregrade, Risiken und Suizidalitaet mitdenken.",
        "Therapie, Beziehungsgestaltung und Aktivierungsangebote verbinden.",
    ],
    "schizophrenie": [
        "Positiv- und Negativsymptome klar unterscheiden.",
        "Beobachtung, Deeskalation, Psychoedukation und Medikation einordnen.",
    ],
    "anorexia": [
        "Definition, Ursachenmodell, Risiken und Komplikationen erklaeren.",
        "Therapie, Beziehung, Essbegleitung und Beobachtung begruenden.",
    ],
    "infektion": [
        "Uebertragungswege, Schutzmassnahmen und Verlaufsbeobachtung nennen.",
        "Hygiene, Isolation und Patientenedukation verknuepfen.",
    ],
    "schlaganfall": [
        "Definition, ischaemisch/haemorrhagisch und FAST-Zeichen nennen.",
        "Notfallkette, Reha, Lagerung und Beobachtung verbinden.",
    ],
    "aphasie": [
        "Aphasie von Dysarthrie abgrenzen und Formen benennen.",
        "Kommunikation anpassen, Ressourcen nutzen und Zeit geben.",
    ],
    "demenz": [
        "Definition, Formen, Verlauf und Alltagsfolgen nennen.",
        "Validierende Kommunikation, Orientierung und Milieugestaltung einbauen.",
    ],
    "gewalt": [
        "Formen von Gewalt erkennen: verbal, physisch, strukturell.",
        "Praevention, Deeskalation, Dokumentation und Teamreflexion nennen.",
    ],
    "schmerz": [
        "Akut vs. chronisch und Schmerzerfassung unterscheiden.",
        "Medikamentoese und nichtmedikamentoese Therapie mit Pflege verbinden.",
    ],
}

GENERIC_DISEASE_HINTS = [
    "DURST: Definition, Ursachen/Risiken, Symptome und Therapie/Pflege sortiert darstellen.",
    "Falls relevant: Formen, Typen, Grade oder Verlauf getrennt nennen.",
]

GENERIC_NON_DISEASE_HINTS = [
    "Definition in einem klaren Satz geben und direkt ein pflegerisches Beispiel anschliessen.",
    "Kommunikation, Sicherheit, Recht oder Ethik mit dem Thema verbinden.",
]

STOP_LINES = {
    "ich lerne etwas ueber:",
    "ich lerne etwas ueber",
    "ich reflektiere:",
    "ich reflektiere",
    "ich denke nach ueber:",
    "ich denke nach ueber",
    "ich kann",
}

OVERVIEW_TITLE_WORDS = (
    "uebersicht",
    "themenstruktur",
    "erwartet",
    "falluebersicht",
)


@dataclass
class PageRecord:
    rowid: int
    title: str
    parent_title: str | None
    grandparent_title: str | None
    greatgrandparent_title: str | None
    text_chunks: list[str] = field(default_factory=list)

    def lineage(self) -> list[str]:
        return [value for value in [self.title, self.parent_title, self.grandparent_title, self.greatgrandparent_title] if value]

    @property
    def full_text(self) -> str:
        return "\n".join(self.text_chunks)


def normalize_text(value: str) -> str:
    normalized = value.replace("Ü", "Ue").replace("ü", "ue")
    normalized = normalized.replace("Ä", "Ae").replace("ä", "ae")
    normalized = normalized.replace("Ö", "Oe").replace("ö", "oe")
    normalized = normalized.replace("ß", "ss")
    normalized = normalized.replace("’", "'").replace("“", '"').replace("”", '"')
    return normalized


def slugify(value: str) -> str:
    slug = normalize_text(value)
    slug = re.sub(r"[^A-Za-z0-9]+", "_", slug)
    slug = slug.strip("_")
    return slug or "bhf"


def bhf_key_from_title(value: str | None) -> str | None:
    if not value:
        return None
    match = re.match(r"^\s*BHF\s*0?(\d{1,2})([ab])?\b", value, re.IGNORECASE)
    if not match:
        return None
    number = int(match.group(1))
    suffix = (match.group(2) or "").lower()
    return f"{number}{suffix}"


def canonical_from_section(section_key: str | None) -> str | None:
    if not section_key:
        return None
    for canonical, aliases in SECTION_GROUPS.items():
        if section_key in aliases:
            return canonical
    return None


def select_canonical_bhf(page: PageRecord) -> str | None:
    for title in [page.parent_title, page.grandparent_title, page.greatgrandparent_title, page.title]:
        canonical = canonical_from_section(bhf_key_from_title(title))
        if canonical:
            return canonical
    return None


def connect_db() -> sqlite3.Connection:
    conn = sqlite3.connect(str(DB_PATH))
    conn.row_factory = sqlite3.Row
    return conn


def is_file_name(text: str) -> bool:
    lowered = text.strip().lower()
    return bool(re.search(r"\.(docx|pdf|pptx|xlsx|jpg|png)$", lowered))


def clean_text_chunk(text: str, title: str) -> str | None:
    text = text.strip()
    if not text:
        return None
    if text.strip() == title.strip():
        return None
    if is_file_name(text):
        return None
    text = text.replace("Computergenerierter Alternativtext:", "")
    text = text.replace("...", "\n")
    text = text.replace("\xa0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = text.strip()
    return text or None


def load_pages() -> list[PageRecord]:
    query = """
        SELECT
            e.rowid AS rowid,
            e.Title AS title,
            p.Title AS parent_title,
            gp.Title AS grandparent_title,
            ggp.Title AS greatgrandparent_title,
            pe.rowid AS pe_rowid,
            pe.Text AS text
        FROM Entities e
        JOIN PageElements pe ON pe.EntityRowId = e.rowid
        LEFT JOIN Entities p ON e.ParentGOID = p.GOID
        LEFT JOIN Entities gp ON p.ParentGOID = gp.GOID
        LEFT JOIN Entities ggp ON gp.ParentGOID = ggp.GOID
        ORDER BY e.rowid, pe.rowid
    """
    pages: dict[int, PageRecord] = {}
    with connect_db() as conn:
        for row in conn.execute(query):
            page = pages.setdefault(
                row["rowid"],
                PageRecord(
                    rowid=row["rowid"],
                    title=row["title"] or "",
                    parent_title=row["parent_title"],
                    grandparent_title=row["grandparent_title"],
                    greatgrandparent_title=row["greatgrandparent_title"],
                ),
            )
            cleaned = clean_text_chunk(row["text"] or "", page.title)
            if cleaned:
                page.text_chunks.append(cleaned)
    return list(pages.values())


def overview_score(page: PageRecord) -> int:
    score = 0
    title_l = normalize_text(page.title).lower()
    text_l = normalize_text(page.full_text).lower()
    if any(word in title_l for word in OVERVIEW_TITLE_WORDS):
        score += 5
    if title_l.startswith("0."):
        score += 1
    if title_l.startswith("3_"):
        score += 1
    if "lernfelduebersicht" in text_l or "graphische lernfelduebersicht" in text_l:
        score += 5
    if "themenstruktur" in text_l:
        score += 4
    if "ich lerne etwas ueber" in text_l:
        score += 3
    if "falluebersicht" in title_l or "falluebersicht" in text_l:
        score += 4
    if len(text_l) > 800:
        score += 2
    return score


def split_lines(text: str) -> list[str]:
    normalized = normalize_text(text)
    normalized = normalized.replace("", "\n• ")
    normalized = normalized.replace("▪", "\n• ")
    normalized = normalized.replace("•", "\n• ")
    normalized = normalized.replace("›", "\n- ")
    normalized = normalized.replace("=", "- ")
    normalized = normalized.replace(">", "- ")
    normalized = re.sub(r"\r\n?", "\n", normalized)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    raw_lines = [line.strip(" -\t") for line in normalized.split("\n")]
    merged: list[str] = []
    for line in raw_lines:
        line = line.strip()
        if not line:
            merged.append("")
            continue
        if merged and merged[-1] and merged[-1][-1] == "-":
            merged[-1] = f"{merged[-1][:-1]}{line}"
            continue
        merged.append(line)
    return merged


def is_boilerplate_line(line: str) -> bool:
    lowered = line.lower()
    if not line:
        return True
    if lowered in STOP_LINES:
        return True
    if lowered.startswith("seite ") or lowered.startswith("ersteller/datum") or lowered.startswith("bearbeitet/datum"):
        return True
    if lowered.startswith("zielgruppe") or lowered.startswith("setting") or lowered.startswith("kompetenzen"):
        return True
    if lowered.startswith("zeitlicher umfang"):
        return True
    if lowered.startswith("tk") or lowered.startswith("datum"):
        return True
    if re.match(r"^bhf\s*0?\d+[ab]?\b", lowered):
        return True
    return False


def is_topic_like_line(line: str) -> bool:
    lowered = line.lower().strip()
    if not lowered:
        return False
    if is_boilerplate_line(line):
        return False
    if lowered.startswith("•") or lowered.startswith("o ") or lowered.startswith("- "):
        return False
    if len(line) > 72:
        return False
    if re.match(r"^\d+$", line):
        return False
    if line.count(" ") < 1 and len(line) < 4:
        return False
    return True


def tidy_heading(value: str) -> str:
    value = value.strip(" -:;,.")
    value = re.sub(r"\s{2,}", " ", value)
    value = value.replace(" / ", " / ")
    return value


def tidy_bullet(value: str) -> str:
    value = value.strip(" -:;,.")
    value = re.sub(r"\s{2,}", " ", value)
    return value


def extract_topic_blocks(text: str) -> list[dict[str, list[str] | str]]:
    lines = split_lines(text)
    blocks: list[dict[str, list[str] | str]] = []
    marker_indices = [
        idx
        for idx, line in enumerate(lines)
        if line.lower() in STOP_LINES or "ich lerne etwas ueber" in line.lower() or "ich reflektiere" in line.lower()
    ]

    for marker_index in marker_indices:
        heading_lines: list[str] = []
        cursor = marker_index - 1
        while cursor >= 0 and len(heading_lines) < 3:
            candidate = lines[cursor].strip()
            if not candidate:
                if heading_lines:
                    break
                cursor -= 1
                continue
            if is_boilerplate_line(candidate):
                break
            if not is_topic_like_line(candidate):
                break
            heading_lines.append(candidate)
            cursor -= 1
        if not heading_lines:
            continue

        heading_candidates = [tidy_heading(line) for line in reversed(heading_lines)]
        bullets: list[str] = []
        cursor = marker_index + 1
        current_bullet = ""
        while cursor < len(lines):
            line = lines[cursor].strip()
            lowered = line.lower()
            if not line:
                if current_bullet:
                    bullets.append(tidy_bullet(current_bullet))
                    current_bullet = ""
                if bullets:
                    break
                cursor += 1
                continue
            if lowered in STOP_LINES or "ich lerne etwas ueber" in lowered or "ich reflektiere" in lowered:
                if current_bullet:
                    bullets.append(tidy_bullet(current_bullet))
                break
            if is_topic_like_line(line) and bullets:
                if current_bullet:
                    bullets.append(tidy_bullet(current_bullet))
                break
            if line.startswith("•") or line.startswith("o ") or line.startswith("- "):
                if current_bullet:
                    bullets.append(tidy_bullet(current_bullet))
                current_bullet = line[1:].strip()
            elif current_bullet:
                current_bullet = f"{current_bullet} {line}"
            elif not is_boilerplate_line(line):
                current_bullet = line
            cursor += 1
        if current_bullet:
            bullets.append(tidy_bullet(current_bullet))

        for heading in heading_candidates:
            if len(heading) < 4:
                continue
            blocks.append({"title": heading, "bullets": bullets[:6]})

    deduped: list[dict[str, list[str] | str]] = []
    seen: set[str] = set()
    for block in blocks:
        key = block["title"].lower()
        if key in seen:
            continue
        seen.add(key)
        deduped.append(block)
    return deduped


def unique(values: Iterable[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for value in values:
        if not value:
            continue
        norm = value.lower()
        if norm in seen:
            continue
        seen.add(norm)
        result.append(value)
    return result


def fallback_topics_from_titles(pages: list[PageRecord]) -> list[dict[str, list[str] | str]]:
    topics: list[dict[str, list[str] | str]] = []
    for page in pages:
        title = tidy_heading(page.title)
        if not title or is_file_name(title):
            continue
        if title.lower().startswith("0. ") or title.lower().startswith("3_"):
            title = re.sub(r"^\d+[\._ ]*", "", title).strip()
        title = title.replace("Fr.", "Frau").replace("Weiß", "Weiss").replace("Gärtner", "Gaertner")
        if len(title) < 3:
            continue
        bullets = [title]
        page_lines = [line for line in split_lines(page.full_text) if line and not is_boilerplate_line(line)]
        for line in page_lines[:3]:
            cleaned = tidy_bullet(line)
            if cleaned.lower() == title.lower():
                continue
            bullets.append(cleaned)
        topics.append({"title": title, "bullets": unique(bullets)[:4]})
    deduped: list[dict[str, list[str] | str]] = []
    seen: set[str] = set()
    for topic in topics:
        key = topic["title"].lower()
        if key in seen:
            continue
        seen.add(key)
        deduped.append(topic)
    return deduped


def manual_topics_for_bhf(bhf_no: str) -> list[dict[str, list[str] | str]]:
    manual_titles = TOPIC_OVERRIDES.get(bhf_no, [])
    return [{"title": title, "bullets": [title]} for title in manual_titles]


def topic_tokens(value: str) -> set[str]:
    normalized = normalize_text(value).lower()
    tokens = re.findall(r"[a-z0-9]+", normalized)
    stop = {
        "und",
        "der",
        "die",
        "das",
        "mit",
        "den",
        "dem",
        "des",
        "bei",
        "nach",
        "von",
        "eine",
        "einer",
        "eines",
        "im",
        "in",
        "auf",
        "fuer",
        "zum",
        "zur",
        "bzw",
        "bns",
        "pflege",
        "menschen",
        "bhf",
    }
    return {token for token in tokens if len(token) > 2 and token not in stop}


def related_lines_for_topic(title: str, extracted_topics: list[dict[str, list[str] | str]], section_titles: list[str]) -> list[str]:
    title_token_set = topic_tokens(title)
    related: list[str] = []
    for topic in extracted_topics:
        candidate_title = topic["title"]
        if topic_tokens(candidate_title) & title_token_set:
            related.append(candidate_title)
            related.extend(topic.get("bullets", []))
    for section_title in section_titles:
        if topic_tokens(section_title) & title_token_set:
            related.append(section_title)
    return trim_lines(unique([title] + related), 4)


def choose_source_pages(pages: list[PageRecord]) -> list[PageRecord]:
    ranked = sorted(pages, key=overview_score, reverse=True)
    selected: list[PageRecord] = []
    seen_texts: set[str] = set()
    for page in ranked:
        if overview_score(page) <= 0:
            continue
        text_key = normalize_text(page.full_text[:1000]).lower()
        if text_key and text_key in seen_texts:
            continue
        if text_key:
            seen_texts.add(text_key)
        selected.append(page)
        if len(selected) >= 3:
            break
    return selected


def build_bhf_sources() -> dict[str, dict]:
    pages = load_pages()
    grouped_pages: dict[str, list[PageRecord]] = defaultdict(list)
    for page in pages:
        canonical = select_canonical_bhf(page)
        if canonical:
            grouped_pages[canonical].append(page)

    result: dict[str, dict] = {}
    for bhf_no in BHF_ORDER:
        section_pages = grouped_pages.get(bhf_no, [])
        selected_pages = choose_source_pages(section_pages)
        extracted_topics: list[dict[str, list[str] | str]] = []
        for page in selected_pages:
            extracted_topics.extend(extract_topic_blocks(page.full_text))

        if not extracted_topics:
            extracted_topics = fallback_topics_from_titles(section_pages)

        manual_topics = manual_topics_for_bhf(bhf_no)
        if manual_topics:
            combined_topics = []
            for topic in manual_topics:
                title = tidy_heading(topic["title"])
                combined_topics.append(
                    {
                        "title": title,
                        "bullets": related_lines_for_topic(
                            title,
                            extracted_topics,
                            unique([page.title for page in section_pages if page.title]),
                        ),
                    }
                )
        else:
            combined_topics = []
            seen_titles: set[str] = set()
            for topic in extracted_topics:
                title = tidy_heading(topic["title"])
                if not title:
                    continue
                key = title.lower()
                if key in seen_titles:
                    continue
                seen_titles.add(key)
                combined_topics.append({"title": title, "bullets": unique(topic.get("bullets", []))[:6]})

        if not combined_topics:
            combined_topics = [{"title": f"BHF {bhf_no}", "bullets": [f"Keine auslesbaren Themen in der lokalen Uebersicht gefunden."]}]

        source_titles = unique([page.title for page in selected_pages]) or unique([page.title for page in section_pages[:5]])
        result[bhf_no] = {
            "bhf_no": bhf_no,
            "source_page_titles": source_titles,
            "source_page_count": len(selected_pages),
            "all_section_page_titles": unique([page.title for page in section_pages if page.title])[:30],
            "topics": combined_topics,
        }
    return result


def is_probable_disease(title: str, bullets: list[str]) -> bool:
    corpus = normalize_text(" ".join([title] + bullets)).lower()
    disease_signals = [
        "krankheitsbild",
        "diabetes",
        "copd",
        "asthma",
        "pneumonie",
        "herzinsuffizienz",
        "parkinson",
        "multiple sklerose",
        "arthritis",
        "depression",
        "schizophrenie",
        "anorexia",
        "schlaganfall",
        "aphasie",
        "demenz",
        "schmerz",
        "dehydratation",
        "malnutrition",
        "adhs",
        "infektionskrankheiten",
        "verbrennung",
        "luftnot",
    ]
    return any(signal in corpus for signal in disease_signals)


def exam_hints_for_topic(title: str, bullets: list[str]) -> list[str]:
    corpus = normalize_text(" ".join([title] + bullets)).lower()
    hints: list[str] = []
    for keyword, keyword_hints in KEYWORD_HINTS.items():
        if keyword in corpus:
            hints.extend(keyword_hints)
    if hints:
        return unique(hints)[:2]
    if is_probable_disease(title, bullets):
        return GENERIC_DISEASE_HINTS
    return GENERIC_NON_DISEASE_HINTS


def trim_lines(lines: Iterable[str], limit: int = 4) -> list[str]:
    cleaned: list[str] = []
    for line in lines:
        line = tidy_bullet(line)
        if not line:
            continue
        if len(line) > 108:
            line = f"{line[:105].rstrip()}..."
        cleaned.append(line)
        if len(cleaned) >= limit:
            break
    return cleaned


def deck_subtitle(deck_data: dict) -> str:
    topics = [topic["title"] for topic in deck_data["topics"][:4]]
    topic_preview = ", ".join(topics)
    return f"Aus OneNote-Uebersicht abgeleitet: {topic_preview}"


def set_fill(shape, color, transparency=0.0):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    fill.transparency = transparency


def set_line(shape, color, width=1.0, transparency=0.0):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width)
    line.transparency = transparency


def add_box(
    slide,
    x,
    y,
    w,
    h,
    text="",
    font_size=18,
    color=TEXT,
    bold=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_lines(
    slide,
    x,
    y,
    w,
    h,
    lines,
    font_size=15,
    color=TEXT,
    bold_first=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    line_spacing=1.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and idx == 0
        run.font.color.rgb = color
    return box


def add_background(slide, accent, accent_dark):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_fill(bg, LIGHT)
    set_line(bg, LIGHT, 0)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.24))
    set_fill(top, accent)
    set_line(top, accent, 0)
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(7.18), Inches(13.333), Inches(0.32))
    set_fill(bottom, rgb("EEF4F7"))
    set_line(bottom, rgb("EEF4F7"), 0)
    for x, y, size, tr in [(11.3, 0.55, 0.62, 0.8), (11.95, 0.95, 0.34, 0.84), (10.78, 1.0, 0.26, 0.86)]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        set_fill(circle, accent_dark, tr)
        set_line(circle, accent_dark, 0.5, 0.9)


def add_footer(slide, text, no, total, accent):
    add_box(slide, 0.55, 7.16, 10.8, 0.15, text, 9, MUTED)
    add_box(
        slide,
        12.0,
        7.12,
        0.8,
        0.18,
        f"{no}/{total}",
        10,
        accent,
        True,
        font_name="Aptos Display",
        align=PP_ALIGN.RIGHT,
    )


def add_tag(slide, text, accent):
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.48), Inches(2.15), Inches(0.34)
    )
    set_fill(tag, accent)
    set_line(tag, accent, 0.8)
    add_box(slide, 0.68, 0.51, 2.0, 0.24, text, 11, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_card(slide, x, y, w, h, title, body_lines, fill_color, accent, body_size=14):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, fill_color)
    set_line(shape, accent, 1.0, 0.25)
    add_box(slide, x + 0.14, y + 0.1, w - 0.28, 0.34, title, 16, accent, True, font_name="Aptos Display")
    add_lines(slide, x + 0.14, y + 0.46, w - 0.28, h - 0.56, body_lines, body_size, TEXT)


def chunk(items, size):
    return [items[index : index + size] for index in range(0, len(items), size)]


def add_title_slide(prs, deck, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.86), Inches(6.35), Inches(5.65)
    )
    set_fill(panel, deck["accent"])
    set_line(panel, deck["accent"], 0.8)
    add_box(slide, 0.95, 1.0, 2.2, 0.3, deck["no"], 13, WHITE, True)
    add_box(slide, 0.95, 1.42, 5.65, 1.55, deck["title"], 27, WHITE, True, font_name="Aptos Display")
    add_lines(slide, 0.95, 3.2, 5.4, 1.3, [deck["subtitle"]], 16, rgb("E7F3FA"))
    add_card(
        slide,
        7.35,
        1.0,
        5.15,
        1.55,
        "Quellen im OneNote",
        [f"• {line}" for line in deck["source_titles"][:4]] or ["• Lokale BHF-Uebersicht"],
        TONES["sky"],
        deck["accent_dark"],
    )
    add_card(
        slide,
        7.35,
        2.85,
        5.15,
        1.55,
        "Lernlogik",
        [
            "erst Thema und Einordnung sichern",
            "dann Unterpunkte/Faelle sortieren",
            "bei Krankheitslehre im DURST-Schema antworten",
        ],
        TONES["sand"],
        deck["accent_dark"],
    )
    add_card(
        slide,
        7.35,
        4.7,
        5.15,
        1.15,
        "Ziel",
        [
            "alle aus der Uebersicht sichtbaren Themen dieses BHF abdecken",
        ],
        TONES["mint"],
        deck["accent_dark"],
    )
    add_footer(slide, "OneNote -> BHF-Uebersicht -> Exam-PowerPoint", no, total, deck["accent"])


def add_topics_slide(prs, deck, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, "Themen", deck["accent"])
    add_box(slide, 0.72, 0.94, 9.5, 0.42, "Themen aus der OneNote-Uebersicht", 24, INK, True, font_name="Aptos Display")
    add_box(
        slide,
        0.75,
        1.3,
        10.8,
        0.25,
        "Extrahiert aus Uebersicht, Themenstruktur, Falluebersicht oder als Fallback aus den Abschnittsseiten.",
        12,
        MUTED,
    )
    topic_lines = [f"• {topic['title']}" for topic in deck["topics"]]
    left, right = chunk(topic_lines, (len(topic_lines) + 1) // 2), []
    if len(left) > 1:
        left, right = left[0], left[1]
    else:
        left, right = left[0], []
    left_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.8), Inches(5.9), Inches(4.95)
    )
    right_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.78), Inches(1.8), Inches(5.85), Inches(4.95)
    )
    set_fill(left_shape, TONES["sky"])
    set_fill(right_shape, TONES["white"])
    set_line(left_shape, deck["accent_dark"], 1.0, 0.25)
    set_line(right_shape, deck["accent_dark"], 1.0, 0.25)
    add_lines(slide, 0.94, 2.02, 5.38, 4.4, left, 15, TEXT)
    add_lines(slide, 7.0, 2.02, 5.28, 4.4, right, 15, TEXT)
    add_footer(slide, "Diese Themenliste ist die Basis fuer alle folgenden Lernkarten.", no, total, deck["accent"])


def add_theme_slide(prs, deck, theme_group, index, total_groups, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, "Themenkarten", deck["accent"])
    title = "Pruefungsrelevante Themenkarten"
    if total_groups > 1:
        title = f"{title} ({index}/{total_groups})"
    add_box(slide, 0.72, 0.94, 9.5, 0.42, title, 24, INK, True, font_name="Aptos Display")
    add_box(
        slide,
        0.75,
        1.28,
        11.0,
        0.25,
        "Jede Karte nimmt ein Thema direkt aus der Uebersicht auf und uebersetzt es in exam-taugliche Sprache.",
        12,
        MUTED,
    )
    positions = [
        (0.72, 1.72, 5.85, 2.18),
        (6.78, 1.72, 5.85, 2.18),
        (0.72, 4.02, 5.85, 2.18),
        (6.78, 4.02, 5.85, 2.18),
    ]
    tones = ["sky", "mint", "sand", "rose"]
    for pos, theme, tone in zip(positions, theme_group, tones):
        body_lines = [f"OneNote: {theme['bullets'][0]}"] if theme["bullets"] else []
        extra_lines = trim_lines(theme["bullets"][1:], 2)
        body_lines.extend([f"• {line}" for line in extra_lines])
        for hint in theme["exam_hints"]:
            body_lines.append(f"Pruefung: {hint}")
        add_card(slide, *pos, theme["title"], body_lines[:6], TONES[tone], deck["accent_dark"], body_size=13)
    add_footer(slide, "Bei Krankheitslehre zuerst D-U-R-S-T ordnen, dann am Fall anwenden.", no, total, deck["accent"])


def add_checklist_slide(prs, deck, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, "Checklist", deck["accent"])
    add_box(slide, 0.72, 0.94, 9.0, 0.42, "Pruefungscheck fuer dieses BHF", 24, INK, True, font_name="Aptos Display")
    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.78), Inches(5.9), Inches(4.95)
    )
    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.78), Inches(1.78), Inches(5.85), Inches(4.95)
    )
    set_fill(left, TONES["mint"])
    set_fill(right, TONES["sand"])
    set_line(left, deck["accent_dark"], 1.0, 0.25)
    set_line(right, deck["accent_dark"], 1.0, 0.25)
    add_box(slide, 0.95, 1.98, 5.2, 0.28, "Das solltest du frei sagen koennen", 17, deck["accent_dark"], True, font_name="Aptos Display")
    left_lines = [
        f"• {topic['title']}" for topic in deck["topics"][:6]
    ]
    add_lines(slide, 0.95, 2.3, 5.2, 4.1, left_lines, 15, TEXT)
    add_box(slide, 7.0, 1.98, 5.15, 0.28, "Letzte 48 Stunden vor der Pruefung", 17, deck["accent_dark"], True, font_name="Aptos Display")
    right_lines = [
        "• jeden Titel laut definieren und mit einem Fallbeispiel verbinden",
        "• bei Krankheitslehre Formen, Typen oder Grade extra nennen",
        "• Pflege, Sicherheit, Recht und Kommunikation immer mitdenken",
        "• in eigenen Worten statt auswendig gelernter Ketten antworten",
    ]
    add_lines(slide, 7.0, 2.3, 5.15, 4.1, right_lines, 15, TEXT)
    add_footer(slide, "Wiederhole diese Liste laut. Was du frei sagen kannst, sitzt in der Klausur.", no, total, deck["accent"])


def build_deck_payload(bhf_no: str, deck_data: dict) -> dict:
    accent, accent_dark = ACCENTS[(int(bhf_no) - 1) % len(ACCENTS)]
    topics = []
    for topic in deck_data["topics"]:
        bullets = trim_lines(topic.get("bullets", []), 3)
        topics.append(
            {
                "title": topic["title"],
                "bullets": bullets or [topic["title"]],
                "exam_hints": exam_hints_for_topic(topic["title"], bullets),
            }
        )
    return {
        "no": f"BHF {bhf_no}",
        "title": f"BHF {bhf_no} | Exam-Praep aus OneNote-Uebersicht",
        "subtitle": deck_subtitle(deck_data),
        "filename": f"BHF_{bhf_no}_OneNote_Uebersicht.pptx",
        "accent": rgb(accent),
        "accent_dark": rgb(accent_dark),
        "source_titles": deck_data["source_page_titles"],
        "topics": topics,
    }


def build_deck(deck: dict) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = deck["title"]
    prs.core_properties.subject = deck["subtitle"]
    prs.core_properties.author = "Codex"
    prs.core_properties.comments = "BHF-Deck auf Basis der lokalen OneNote-Uebersicht"

    theme_groups = chunk(deck["topics"], 4)
    total = 3 + len(theme_groups)
    slide_no = 1
    add_title_slide(prs, deck, slide_no, total)
    slide_no += 1
    add_topics_slide(prs, deck, slide_no, total)
    slide_no += 1
    for index, theme_group in enumerate(theme_groups, start=1):
        add_theme_slide(prs, deck, theme_group, index, len(theme_groups), slide_no, total)
        slide_no += 1
    add_checklist_slide(prs, deck, slide_no, total)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / deck["filename"]
    prs.save(str(path))
    return path


def build_all():
    bhf_sources = build_bhf_sources()
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    JSON_PATH.write_text(json.dumps(bhf_sources, indent=2, ensure_ascii=True), encoding="utf-8")
    paths: list[Path] = []
    for bhf_no in BHF_ORDER:
        deck_payload = build_deck_payload(bhf_no, bhf_sources[bhf_no])
        paths.append(build_deck(deck_payload))
    for path in paths:
        print(f"saved {path}")
    print(f"saved {JSON_PATH}")


if __name__ == "__main__":
    build_all()
