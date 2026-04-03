from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path("/Users/arielhavana/Downloads/KrampfanfallBNS-Krämpfe.pptx")
OUT = Path("/Users/arielhavana/antigr/porc/mockups/KrampfanfallBNS-Kraempfe-restyled.pptx")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


def section_theme(slide_no: int):
    if slide_no <= 2:
        return {
            "name": "Überblick",
            "bg": rgb("FFF8F2"),
            "accent": rgb("C86D53"),
            "soft": rgb("F6E2D7"),
        }
    if slide_no <= 10:
        return {
            "name": "Medizin",
            "bg": rgb("F5FBF8"),
            "accent": rgb("3F7C73"),
            "soft": rgb("DDEFEA"),
        }
    if slide_no <= 19:
        return {
            "name": "Pflege",
            "bg": rgb("FFFBF1"),
            "accent": rgb("B1883C"),
            "soft": rgb("F4E7C3"),
        }
    if slide_no <= 24:
        return {
            "name": "Familie",
            "bg": rgb("FBF7FD"),
            "accent": rgb("7A5F92"),
            "soft": rgb("E9DDF3"),
        }
    return {
        "name": "Abschluss",
        "bg": rgb("F4F8FB"),
        "accent": rgb("2F5872"),
        "soft": rgb("DCE8F1"),
    }


TEXT = rgb("283848")
MUTED = rgb("5E7285")
WHITE = rgb("FFFFFF")
TITLE_FONT = "Gill Sans"
BODY_FONT = "Trebuchet MS"


def apply_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_run_style(run, font_name, size=None, bold=None, color=None):
    run.font.name = font_name
    if size is not None:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def all_text_shapes(slide):
    items = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            items.append(shape)
    return items


def is_title_shape(shape, title_shape):
    if shape == title_shape:
        return True
    top_in = shape.top / 914400
    text = shape.text_frame.text.strip()
    return top_in < 1.45 and len(text) < 120


def restyle_text_shape(shape, title=False, accent=TEXT):
    tf = shape.text_frame
    tf.word_wrap = True
    if title:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in tf.paragraphs:
        if not paragraph.runs:
            run = paragraph.add_run()
            run.text = paragraph.text
        for run in paragraph.runs:
            if title:
                size = run.font.size or Pt(24)
                set_run_style(run, TITLE_FONT, size=size, bold=True, color=accent)
            else:
                size = run.font.size or Pt(16)
                set_run_style(run, BODY_FONT, size=size, bold=run.font.bold, color=TEXT)
        if title:
            paragraph.alignment = paragraph.alignment or PP_ALIGN.LEFT
        else:
            if paragraph.alignment is None:
                paragraph.alignment = PP_ALIGN.LEFT


def tint_shape(shape, accent, soft):
    if shape.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.TEXT_BOX):
        return
    try:
        fill = shape.fill
        line = shape.line
    except Exception:
        return

    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        line.color.rgb = accent
        line.transparency = 1.0
        return

    try:
        fill.solid()
        fill.fore_color.rgb = soft
        fill.transparency = 0.12
    except Exception:
        pass

    try:
        line.color.rgb = accent
        line.width = Pt(1.2)
        line.transparency = 0.15
    except Exception:
        pass


def add_section_tag(slide, label, accent):
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.7),
        Inches(0.28),
        Inches(1.95),
        Inches(0.34),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = accent
    tag.line.color.rgb = accent
    tf = tag.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if not tf.paragraphs[0].runs:
        run = tf.paragraphs[0].add_run()
    else:
        run = tf.paragraphs[0].runs[0]
    run.text = label.upper()
    run.font.name = BODY_FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_bottom_bar(slide, accent):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(6.98), Inches(2.2), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent


def build():
    prs = Presentation(str(SOURCE))
    prs.core_properties.title = "Krampfanfall BNS-Krämpfe - Restyled"
    prs.core_properties.comments = "Neu eingefärbte und neu typografierte Version"

    for idx, slide in enumerate(prs.slides, start=1):
        theme = section_theme(idx)
        apply_slide_background(slide, theme["bg"])
        text_shapes = all_text_shapes(slide)
        title_shape = min(text_shapes, key=lambda s: (s.top, s.left)) if text_shapes else None

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            tint_shape(shape, theme["accent"], theme["soft"])

        for shape in text_shapes:
            restyle_text_shape(shape, title=is_title_shape(shape, title_shape), accent=theme["accent"])

        add_section_tag(slide, theme["name"], theme["accent"])
        add_bottom_bar(slide, theme["accent"])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"saved {OUT}")


if __name__ == "__main__":
    build()
