from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("/Users/arielhavana/antigr/porc/mockups/BNS-West-Syndrom-Neuaufbau.pptx")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


NAVY = rgb("10233F")
DEEP = rgb("19345A")
INK = rgb("24344D")
TEAL = rgb("60B8B1")
SKY = rgb("DDEFF6")
POWDER = rgb("EEF5F8")
CORAL = rgb("F28A6C")
APRICOT = rgb("F8D6C9")
GOLD = rgb("F2C66D")
SAND = rgb("F7EEDC")
MINT = rgb("D8EFE6")
ROSE = rgb("F6DADF")
WHITE = rgb("FFFFFF")
SLATE = rgb("65748B")
GRAPHITE = rgb("344255")
LIGHT = rgb("F9FBFC")
LINE = rgb("BED0DB")


def set_fill(shape, color, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1.5, transparency: float = 0.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_box(slide, x, y, w, h, text="", font_size=18, color=INK, bold=False,
            font_name="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(5)
    tf.margin_bottom = Pt(5)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(slide, x, y, w, h, lines, font_size=18, color=INK, bold_first=False,
                  font_name="Aptos", leading=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(5)
    tf.margin_bottom = Pt(5)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = leading
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and idx == 0
        run.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, title, body_lines, fill_color, title_color=NAVY,
             body_color=INK, title_size=21, body_size=16, line_color=None):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(card, fill_color)
    set_line(card, line_color or fill_color, 1)
    add_box(slide, x + 0.16, y + 0.14, w - 0.32, 0.42, title, title_size, title_color, True,
            font_name="Aptos Display")
    add_multiline(slide, x + 0.16, y + 0.55, w - 0.32, h - 0.7, body_lines, body_size, body_color)
    return card


def add_tag(slide, x, y, w, h, text, fill_color, text_color=WHITE):
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(tag, fill_color)
    set_line(tag, fill_color, 1)
    add_box(slide, x + 0.04, y + 0.01, w - 0.08, h - 0.02, text, 12, text_color, True,
            font_name="Aptos", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return tag


def add_footer(slide, text, number, total):
    add_box(slide, 0.45, 7.0, 10.9, 0.22, text, 9, SLATE, False, font_name="Aptos")
    add_box(slide, 12.2, 6.95, 0.5, 0.25, f"{number}/{total}", 10, NAVY, True, font_name="Aptos")


def add_stat_card(slide, x, y, w, h, number, label, note, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, fill_color)
    set_line(shape, fill_color, 1)
    add_box(slide, x + 0.15, y + 0.12, w - 0.3, 0.38, number, 26, NAVY, True, font_name="Aptos Display")
    add_box(slide, x + 0.15, y + 0.58, w - 0.3, 0.26, label, 12, DEEP, True, font_name="Aptos")
    add_multiline(slide, x + 0.15, y + 0.95, w - 0.3, h - 1.05, note, 15, GRAPHITE)


def add_hex_cluster(slide, x, y, scale=1.0):
    specs = [
        (0.0, 0.0, 0.45, SKY),
        (0.34, 0.22, 0.45, WHITE),
        (0.68, 0.0, 0.45, SKY),
        (0.17, 0.41, 0.45, WHITE),
        (0.51, 0.41, 0.45, SKY),
    ]
    for ox, oy, size, color in specs:
        hex_shape = slide.shapes.add_shape(
            MSO_SHAPE.HEXAGON,
            Inches(x + ox * scale),
            Inches(y + oy * scale),
            Inches(size * scale),
            Inches(size * scale),
        )
        set_fill(hex_shape, color, transparency=0.12 if color == SKY else 0.55)
        set_line(hex_shape, SKY, 1.1)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Therapie der Blitz-Nick-Salaam-Epilepsie"
    prs.core_properties.subject = "West-Syndrom | Patientenleitlinie"
    prs.core_properties.author = "Codex"
    prs.core_properties.comments = "Neu aufgebaute Praesentation aus der Patientenleitlinie 2021"

    total = 9

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, LIGHT)
    set_line(bg, LIGHT, 0)
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.6), Inches(6.4), Inches(5.95))
    set_fill(panel, NAVY)
    set_line(panel, NAVY, 1)
    arc = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(4.75), Inches(1.6), Inches(1.6))
    set_fill(arc, TEAL, transparency=0.14)
    set_line(arc, TEAL, 0.8, transparency=0.3)
    add_hex_cluster(slide, 0.42, 0.4, 1.2)
    add_tag(slide, 0.88, 1.02, 2.45, 0.34, "PATIENTENLEITLINIE 2021", TEAL)
    add_box(slide, 0.88, 1.42, 5.55, 1.65, "Therapie der\nBlitz-Nick-Salaam-\nEpilepsie", 28, WHITE, True, font_name="Aptos Display")
    add_multiline(
        slide, 0.88, 3.48, 5.25, 1.3,
        ["West-Syndrom", "Fokus: fruehe Erkennung, schnelle Therapie,", "engmaschige Kontrolle und Elternbegleitung"],
        17, SKY
    )
    add_box(slide, 0.88, 5.45, 5.2, 0.5, "AWMF-Register 022/022 | Gesellschaft fuer Neuropaediatrie", 12, WHITE)
    add_card(
        slide, 7.45, 1.0, 5.05, 1.38, "1 | Diagnose schnell sichern",
        ["EEG bei Verdacht innerhalb weniger Tage", "mit Wach-, Schlaf- und moeglichst Video-Ableitung"],
        SKY, title_color=NAVY, body_color=INK
    )
    add_card(
        slide, 7.45, 2.64, 5.05, 1.38, "2 | Therapie sofort beginnen",
        ["Nach gesicherter Diagnose rasch wirksame Behandlung", "nicht erst Beobachtung ueber Wochen"],
        APRICOT, title_color=NAVY, body_color=INK
    )
    add_card(
        slide, 7.45, 4.28, 5.05, 1.38, "3 | Erfolg frueh pruefen",
        ["Kontrolle nach etwa 14 Tagen", "Ziel: Anfallsfreiheit und keine Hypsarrhythmie mehr"],
        SAND, title_color=NAVY, body_color=INK
    )
    add_footer(slide, "Quelle: Patientenleitlinie Therapie der BNS-Epilepsie, Abschnitte 1-2 und 6", 1, total)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, POWDER)
    set_line(bg, POWDER, 0)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.36), Inches(7.5))
    set_fill(stripe, TEAL)
    set_line(stripe, TEAL, 0)
    add_box(slide, 0.6, 0.45, 6.5, 0.55, "Kernbotschaften auf einen Blick", 24, NAVY, True, font_name="Aptos Display")
    add_box(slide, 0.62, 0.94, 8.6, 0.35, "Die Leitlinie betont Geschwindigkeit, Wirksamkeit und klare Elterninformation.", 13, GRAPHITE)
    add_stat_card(slide, 0.62, 1.45, 3.9, 1.45, "3-18", "Lebensmonat", ["typisches Auftreten", "nur ausnahmsweise frueher oder spaeter"], SKY)
    add_stat_card(slide, 4.72, 1.45, 3.9, 1.45, "1 von 2500", "Saeuglingen", ["seltene, aber folgenreiche Erkrankung", "fruehes Handeln ist entscheidend"], MINT)
    add_stat_card(slide, 8.82, 1.45, 3.9, 1.45, "3 Wochen", "kritische Zeitmarke", ["ein grosser Abstand bis Therapiebeginn", "verschlechtert die Entwicklungschance"], APRICOT)
    add_card(
        slide, 0.62, 3.35, 3.0, 1.45, "Verdacht ernst nehmen",
        ["• serielle, aehnliche Bewegungsablaeufe", "• oft kurz nach dem Aufwachen", "• Elternvideo ist sehr hilfreich"],
        WHITE, body_size=15, line_color=LINE
    )
    add_card(
        slide, 3.86, 3.35, 3.0, 1.45, "EEG macht die Diagnose fest",
        ["• schmerzlos, aber entscheidend", "• Schlaf-EEG mit Minuten nach dem Wecken", "• Hypsarrhythmie ist typisch"],
        WHITE, body_size=15, line_color=LINE
    )
    add_card(
        slide, 7.10, 3.35, 2.8, 1.45, "Erste Wahl",
        ["• Hormone: ACTH oder Prednisolon", "• haeufig kombiniert mit Vigabatrin"], WHITE, body_size=15, line_color=LINE
    )
    add_card(
        slide, 10.14, 3.35, 2.58, 1.45, "Sonderfall",
        ["• Tuberose Sklerose", "• oder Gruende gegen Hormone", "• dann Vigabatrin primaer"], WHITE, body_size=15, line_color=LINE
    )
    add_card(
        slide, 0.62, 5.08, 5.2, 1.45, "Eltern sind Teil der Therapie",
        ["Wiederholte Gespraeche in verstaendlicher Sprache, Besprechung von Aengsten,", "Anfallskalender, Befundkopien und psychosoziale Unterstuetzung gehoeren dazu."],
        ROSE, body_size=15
    )
    add_card(
        slide, 6.1, 5.08, 6.62, 1.45, "Therapieziel der Leitlinie",
        ["Rasche Anfallsfreiheit fuer BNS-Anfaelle und Sistieren der Hypsarrhythmie", "im Wach- und Schlaf-EEG als Voraussetzung fuer die bestmoegliche Entwicklung."],
        SAND, body_size=15
    )
    add_footer(slide, "Quelle: Zusammenfassung, Hintergrundwissen und Empfehlungen 1-7", 2, total)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, LIGHT)
    set_line(bg, LIGHT, 0)
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.45), Inches(12.0), Inches(0.72))
    set_fill(band, NAVY)
    set_line(band, NAVY, 0)
    add_box(slide, 0.92, 0.58, 7.0, 0.32, "BNS-Anfaelle erkennen", 25, WHITE, True, font_name="Aptos Display")
    add_card(
        slide, 0.72, 1.55, 4.1, 4.9, "Typisches Erscheinungsbild",
        [
            "Blitz: ploetzliches Zusammenzucken wie vom Blitz getroffen.",
            "Nick: Kopf und Rumpf beugen sich nach vorne.",
            "Salaam: gleichzeitig heben sich die Arme.",
            "",
            "Weitere Muster sind Streck- oder Mischformen.",
            "Manche Kinder schreien, andere wirken nur kurz veraendert."
        ],
        SKY, body_size=16
    )
    cluster = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.15), Inches(1.55), Inches(3.25), Inches(4.9))
    set_fill(cluster, WHITE)
    set_line(cluster, LINE, 1.4)
    add_box(slide, 5.35, 1.72, 2.6, 0.4, "So laeuft ein Cluster ab", 21, NAVY, True, font_name="Aptos Display")
    add_multiline(
        slide, 5.35, 2.24, 2.6, 0.9,
        ["Ein einzelner Spasmus dauert", "oft nur etwa 1 Sekunde."],
        18, GRAPHITE
    )
    for idx in range(6):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.45 + idx * 0.42), Inches(3.28), Inches(0.18), Inches(0.18))
        set_fill(dot, CORAL if idx in (1, 2, 3, 4) else TEAL)
        set_line(dot, WHITE, 0.2)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.45), Inches(3.9), Inches(2.2), Inches(0.08))
    set_fill(bar, LINE)
    set_line(bar, LINE, 0)
    add_multiline(
        slide, 5.35, 4.2, 2.7, 1.3,
        ["Cluster koennen bis zu 100 Spasmen umfassen", "und mehrere Minuten dauern.", "", "Haeufig in den ersten Minuten nach dem Aufwachen."],
        16, GRAPHITE
    )
    add_card(
        slide, 8.72, 1.55, 3.92, 1.45, "Wann besonders daran denken?",
        ["• bei neurologischer Vorerkrankung", "• bei Trisomie 21", "• bei immer gleichen Bewegungsserien"], MINT, body_size=15
    )
    add_card(
        slide, 8.72, 3.28, 3.92, 1.45, "Was verwechselt werden kann",
        ["• gutartige fruehkindliche Myoklonien", "• Sandifer-Syndrom", "• andere nicht-epileptische Bewegungen"], APRICOT, body_size=15
    )
    add_card(
        slide, 8.72, 5.0, 3.92, 1.45, "Praktischer Tipp",
        ["Ein Elternvideo ersetzt kein EEG, kann aber den Weg zur raschen Abklaerung deutlich verkuerzen."], SAND, body_size=15
    )
    add_footer(slide, "Quelle: Hintergrundwissen 3.2-3.5", 3, total)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, rgb("FFF8F3"))
    set_line(bg, rgb("FFF8F3"), 0)
    add_box(slide, 0.7, 0.48, 6.0, 0.5, "Ursachen und Entwicklung", 24, NAVY, True, font_name="Aptos Display")
    add_box(slide, 0.7, 0.95, 7.5, 0.3, "Die Prognose haengt stark von der Grunderkrankung und vom Zeitgewinn bis zur wirksamen Therapie ab.", 13, GRAPHITE)
    add_card(
        slide, 0.7, 1.45, 3.7, 1.52, "Vor der Geburt oder um die Geburt herum",
        ["• Fehlbildungen des Gehirns", "• Infektionen, Blutung, Sauerstoffmangel", "• andere fruehe Schaedigungen"], SKY, body_size=15
    )
    add_card(
        slide, 0.7, 3.18, 3.7, 1.52, "Genetisch oder metabolisch",
        ["• genetische Erkrankungen", "• Tuberose Sklerose", "• selten Stoffwechselstoerungen"], MINT, body_size=15
    )
    add_card(
        slide, 0.7, 4.91, 3.7, 1.52, "Nach der Geburt",
        ["• Hirnhautentzuendung", "• Hirninfarkt", "• Hirnverletzung", "• etwa ein Drittel bleibt ohne erkennbare Ursache"], APRICOT, body_size=15
    )
    main = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.55), Inches(7.9), Inches(4.95))
    set_fill(main, WHITE)
    set_line(main, LINE, 1.4)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.7), Inches(1.08), Inches(1.35), Inches(1.35))
    set_fill(circle, CORAL)
    set_line(circle, CORAL, 1)
    add_box(slide, 10.92, 1.37, 0.9, 0.45, "3\nWochen", 14, WHITE, True, font_name="Aptos", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_box(slide, 4.95, 1.78, 4.6, 0.42, "Entwicklungsrisiko durch die BNS-Epilepsie selbst", 21, NAVY, True, font_name="Aptos Display")
    add_multiline(
        slide, 4.95, 2.32, 7.2, 3.8,
        [
            "• Entwicklung kann sich verlangsamen, stehen bleiben oder ruecklaeufig werden.",
            "• Auch Blickkontakt und Verhalten koennen sich veraendern.",
            "• Stoerungen durch die BNS-Epilepsie koennen sich unter schneller, erfolgreicher Therapie zurueckbilden.",
            "• Ein schneller Start schuetzt nicht immer vollstaendig, verbessert aber die Chancen.",
            "",
            "Merksatz: Je kuerzer das Intervall zwischen Anfallsbeginn und wirksamer Therapie, desto besser die Aussichten."
        ],
        17, GRAPHITE
    )
    add_footer(slide, "Quelle: Hintergrundwissen 3.6-3.7", 4, total)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, rgb("F6FAFC"))
    set_line(bg, rgb("F6FAFC"), 0)
    add_box(slide, 0.7, 0.45, 5.0, 0.5, "Diagnostik in klaren Schritten", 24, NAVY, True, font_name="Aptos Display")
    add_box(slide, 0.72, 0.93, 8.5, 0.3, "Die Leitlinie verlangt keine lange Beobachtung, sondern rasche Sicherung oder Entkraeftung des Verdachts.", 13, GRAPHITE)
    steps = [
        ("1", "Video + Arztkontakt", SKY),
        ("2", "Anamnese + Entwicklung", MINT),
        ("3", "EEG Wach/Schlaf/Video", APRICOT),
        ("4", "Neurologischer Status", SAND),
        ("5", "MRT, Labor, Genetik nach Bedarf", ROSE),
    ]
    sx = 0.7
    for num, label, fill in steps:
        chevron = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(sx), Inches(1.75), Inches(2.38), Inches(1.02))
        set_fill(chevron, fill)
        set_line(chevron, WHITE, 1)
        add_box(slide, sx + 0.16, 1.89, 0.3, 0.3, num, 18, NAVY, True, font_name="Aptos Display")
        add_multiline(slide, sx + 0.46, 1.86, 1.56, 0.5, [label], 14, NAVY, True)
        sx += 2.28
    add_card(
        slide, 0.7, 3.2, 6.15, 2.72, "EEG: der Schluesselbefund",
        [
            "• zeichnet elektrische Aktivitaet des Gehirns auf",
            "• schmerzlos und unschaedlich",
            "• Hypsarrhythmie ist das typische chaotische Muster",
            "• Schlaf-EEG ist wichtig, weil die Auffaelligkeit dort erst sichtbar sein kann",
            "• nach dem Wecken weiter ableiten, weil Anfaelle oft dann auftreten"
        ],
        WHITE, body_size=16, line_color=LINE
    )
    add_card(
        slide, 7.15, 3.2, 5.55, 2.72, "Weitere Untersuchungen",
        [
            "• klinische und neurologische Untersuchung mit Entwicklungsbeurteilung",
            "• Haut- und Augenuntersuchung",
            "• fast immer MRT",
            "• Blut- und Urinwerte vor Therapie",
            "• genetische Tests, selten Stoffwechsel- oder Liquordiagnostik"
        ],
        WHITE, body_size=16, line_color=LINE
    )
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.6))
    set_fill(note, NAVY)
    set_line(note, NAVY, 0)
    add_box(slide, 0.95, 6.28, 11.3, 0.2, "Wichtig: Das EEG ist fuer die Bestaetigung der BNS-Epilepsie unbedingt erforderlich.", 15, WHITE, True)
    add_footer(slide, "Quelle: Diagnostik 4.1-4.3", 5, total)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, LIGHT)
    set_line(bg, LIGHT, 0)
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.48), Inches(12.0), Inches(0.86))
    set_fill(banner, DEEP)
    set_line(banner, DEEP, 0)
    add_box(slide, 0.95, 0.67, 9.6, 0.35, "Therapieziel und Behandlungsstart", 24, WHITE, True, font_name="Aptos Display")
    add_card(
        slide, 0.72, 1.65, 5.88, 3.45, "Erste Wahl laut Leitlinie",
        [
            "• Hormone: ACTH oder Prednisolon",
            "• oft kombiniert mit Vigabatrin",
            "• beste Erfolgschancen fuer eine rasche BNS-Anfallsfreiheit",
            "",
            "Ziel bleibt immer:",
            "Anfallsfreiheit plus Sistieren der Hypsarrhythmie im Wach- und Schlaf-EEG"
        ],
        SKY, body_size=17
    )
    add_card(
        slide, 6.82, 1.65, 5.78, 3.45, "Wann Vigabatrin primaer eingesetzt wird",
        [
            "• bei Tuberosem-Sklerose-Komplex",
            "• wenn Gruende gegen eine Hormontherapie sprechen",
            "",
            "Die Leitlinie stuft Vigabatrin in diesen Situationen als erste Wahl ein."
        ],
        APRICOT, body_size=17
    )
    lane = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(5.42), Inches(11.88), Inches(1.0))
    set_fill(lane, WHITE)
    set_line(lane, LINE, 1.4)
    for ox, label in [(1.0, "Tag 0"), (6.0, "Tag 14"), (10.3, "naechster Schritt")]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ox), Inches(5.68), Inches(0.38), Inches(0.38))
        set_fill(c, TEAL if label != "naechster Schritt" else CORAL)
        set_line(c, WHITE, 0.5)
    connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.35), Inches(5.83), Inches(8.9), Inches(0.08))
    set_fill(connector, LINE)
    set_line(connector, LINE, 0)
    add_box(slide, 1.0, 6.05, 1.0, 0.2, "Diagnose gesichert", 10, GRAPHITE, True)
    add_box(slide, 5.56, 6.05, 2.2, 0.2, "klinische + EEG-Evaluation", 10, GRAPHITE, True)
    add_box(slide, 9.9, 6.05, 2.0, 0.2, "wechseln/eskalieren, falls noetig", 10, GRAPHITE, True)
    add_footer(slide, "Quelle: Empfehlungen 1-8", 6, total)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, NAVY)
    set_line(bg, NAVY, 0)
    add_box(slide, 0.7, 0.48, 6.0, 0.5, "Wirksame Therapieschemata", 24, WHITE, True, font_name="Aptos Display")
    add_box(slide, 0.72, 0.95, 7.6, 0.3, "Die Leitlinie nennt diese Regime als wirksam und praxisrelevant.", 13, SKY)
    schemes = [
        (0.75, "Prednisolon", TEAL, ["40-60 mg/Tag per os", "2 Wochen Therapie", "+ 2 Wochen ausschleichen", "", "gut untersucht", "Tablettengabe moeglich"]),
        (4.42, "Depot-ACTH", CORAL, ["40-60 IE i.m.", "2 Wochen alle 2 Tage", "+ Beendigung ueber Prednisolon", "", "klassische Hormonoption", "kurzes Schema mit weniger Nebenwirkungen"]),
        (8.09, "Vigabatrin", GOLD, ["100-150 mg/kg/Tag", "3 Monate Therapie", "+ 1 Monat ausschleichen", "", "primaer bei tuberoeser Sklerose", "oder Kontraindikation gegen Hormone"]),
    ]
    for x, title, header_color, lines in schemes:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(3.25), Inches(4.85))
        set_fill(card, WHITE)
        set_line(card, WHITE, 0.8)
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.16), Inches(1.78), Inches(2.93), Inches(0.6))
        set_fill(header, header_color)
        set_line(header, header_color, 0)
        add_box(slide, x + 0.3, 1.93, 2.5, 0.2, title, 19, NAVY if header_color == GOLD else WHITE, True, font_name="Aptos Display")
        add_multiline(slide, x + 0.24, 2.6, 2.78, 3.45, lines, 17, INK)
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(6.6), Inches(10.15), Inches(0.42))
    set_fill(note, rgb("2A4D77"))
    set_line(note, rgb("2A4D77"), 0)
    add_box(slide, 0.95, 6.68, 9.6, 0.18, "Kurze Behandlungskonzepte verbessern die Vertraeglichkeit, ohne die Wirksamkeit aufzugeben.", 12, WHITE)
    add_footer(slide, "Quelle: Empfehlungen 6-8", 7, total)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, rgb("FFF9F1"))
    set_line(bg, rgb("FFF9F1"), 0)
    add_box(slide, 0.7, 0.48, 7.0, 0.5, "Wenn die erste Wahl nicht ausreicht", 24, NAVY, True, font_name="Aptos Display")
    add_card(
        slide, 0.72, 1.35, 5.95, 2.0, "Weitere Therapieoptionen",
        [
            "Ketogene Diaet, Sultiam, Topiramat, Valproat, Zonisamid oder Benzodiazepine",
            "",
            "Diese Optionen kommen nach unzureichender Wirkung der ersten Wahl in Betracht."
        ],
        WHITE, body_size=16, line_color=LINE
    )
    add_card(
        slide, 6.92, 1.35, 5.68, 2.0, "Frueh an Epilepsiechirurgie denken",
        [
            "vor allem bei fokalen Laesionen, fokalen Anfaellen, fokalen EEG-Befunden",
            "und nachgewiesener Therapieresistenz",
            "",
            "Abklaerung nur in erfahrenen Epilepsiezentren."
        ],
        WHITE, body_size=16, line_color=LINE
    )
    add_card(
        slide, 0.72, 3.72, 5.95, 2.42, "Nebenwirkungen von ACTH / Prednisolon",
        [
            "• Unruhe, Weinen, Schlafstoerung",
            "• Cushing-Zeichen und Gewichtszunahme",
            "• Infektanfaelligkeit, Bluthochdruck, Blutzuckeranstieg",
            "• bei Fieber sofort Aerztin/Aerztin kontaktieren"
        ],
        ROSE, body_size=16
    )
    add_card(
        slide, 6.92, 3.72, 5.68, 2.42, "Vigabatrin: worauf zu achten ist",
        [
            "• meist Schlaefrigkeit oder Reizbarkeit",
            "• Gesichtsfeldeinschraenkungen v. a. bei langer Behandlung",
            "• Risiko bei kurzer 3-6-monatiger Therapie wahrscheinlich sehr gering",
            "• bei unzureichender Wirkung konsequent wieder absetzen"
        ],
        MINT, body_size=16
    )
    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(6.35), Inches(11.88), Inches(0.52))
    set_fill(band, NAVY)
    set_line(band, NAVY, 0)
    add_box(slide, 0.94, 6.46, 11.2, 0.2, "Die Leitlinie betont: Das Risiko der Medikamente ist insgesamt kleiner als das Risiko einer unzureichend behandelten BNS-Epilepsie.", 12, WHITE)
    add_footer(slide, "Quelle: Empfehlungen 9-11 und Nebenwirkungen, Abschnitt 7", 8, total)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_fill(bg, POWDER)
    set_line(bg, POWDER, 0)
    add_box(slide, 0.7, 0.45, 7.0, 0.5, "Nachsorge und Unterstuetzung", 24, NAVY, True, font_name="Aptos Display")
    add_card(
        slide, 0.72, 1.28, 3.72, 5.25, "Nach der Akuttherapie",
        [
            "• regelmaessige EEG- und klinische Kontrollen",
            "• Entwicklungsdiagnostik im Alter von 18 Monaten",
            "• erneute standardisierte Testung vor Schuleintritt",
            "• bei Rueckfall gleiche oder andere Therapien moeglich"
        ],
        WHITE, body_size=16, line_color=LINE
    )
    add_card(
        slide, 4.8, 1.28, 3.72, 5.25, "Was Familien konkret hilft",
        [
            "• eigene Akte mit Befunden, Arztbriefen und Beobachtungen",
            "• Anfallskalender und Nebenwirkungsdokumentation",
            "• wiederholte Aufklaerung in verstaendlicher Sprache",
            "• psychologische Begleitung und Sozialdienst"
        ],
        WHITE, body_size=16, line_color=LINE
    )
    add_card(
        slide, 8.88, 1.28, 3.72, 5.25, "Nuetzliche Adressen",
        [
            "epilepsie bundes-elternverband",
            "www.epilepsie-elternverband.de",
            "",
            "Gesellschaft fuer Neuropaediatrie",
            "gesellschaft-fuer-neuropaediatrie.org",
            "",
            "Weitere Hilfe: Epilepsie-Liga, Deutsche Epilepsievereinigung, epikurier"
        ],
        WHITE, body_size=15, line_color=LINE
    )
    strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(6.55), Inches(11.88), Inches(0.42))
    set_fill(strip, TEAL)
    set_line(strip, TEAL, 0)
    add_box(slide, 0.95, 6.63, 11.0, 0.18, "Internetrecherche kann helfen, sollte aber immer mit dem Behandlungsteam eingeordnet werden.", 12, WHITE)
    add_footer(slide, "Quelle: Abschnitt 8-9 der Patientenleitlinie", 9, total)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"saved {OUT_PATH}")


if __name__ == "__main__":
    build()
