from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("/Users/arielhavana/antigr/porc/mockups/BNS-Kraempfe-und-Therapie-Neu.pptx")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


NAVY = rgb("122847")
DEEP = rgb("1E3F63")
TEXT = rgb("23364D")
MUTED = rgb("61758D")
WHITE = rgb("FFFFFF")
BG = rgb("F8FBFD")
LINE = rgb("C9D9E4")
BLUE = rgb("DCEFFC")
MINT = rgb("DDF2E7")
APRICOT = rgb("F9DED4")
SUN = rgb("F7ECC8")
ROSE = rgb("F8E0E6")
LILAC = rgb("E7DDF6")
TEAL = rgb("68B8B2")
CORAL = rgb("F28B6A")
GOLD = rgb("EABF62")


def set_fill(shape, color, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1.3, transparency: float = 0.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_box(
    slide,
    x,
    y,
    w,
    h,
    text="",
    font_size=18,
    color=TEXT,
    bold=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_lines(
    slide,
    x,
    y,
    w,
    h,
    lines,
    font_size=16,
    color=TEXT,
    bold_first=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    line_spacing=1.12,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and idx == 0
        run.font.color.rgb = color
    return box


def add_card(
    slide,
    x,
    y,
    w,
    h,
    title,
    body_lines,
    fill_color,
    body_size=15,
    title_size=20,
    title_color=NAVY,
    body_color=TEXT,
    border=LINE,
):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(card, fill_color)
    set_line(card, border, 1.1)
    add_box(
        slide,
        x + 0.16,
        y + 0.12,
        w - 0.32,
        0.36,
        title,
        title_size,
        title_color,
        True,
        font_name="Aptos Display",
    )
    add_lines(slide, x + 0.16, y + 0.52, w - 0.32, h - 0.66, body_lines, body_size, body_color)
    return card


def add_tag(slide, x, y, w, h, text, fill_color, text_color=WHITE):
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(tag, fill_color)
    set_line(tag, fill_color, 0.8)
    add_box(
        slide,
        x + 0.03,
        y + 0.01,
        w - 0.06,
        h - 0.02,
        text,
        11,
        text_color,
        True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_stat(slide, x, y, w, h, big, small, fill_color):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(box, fill_color)
    set_line(box, fill_color, 1)
    add_box(slide, x + 0.14, y + 0.1, w - 0.28, 0.34, big, 24, NAVY, True, font_name="Aptos Display")
    add_lines(slide, x + 0.14, y + 0.52, w - 0.28, h - 0.6, [small], 12, DEEP, True)


def add_hex_pattern(slide, x, y, scale=1.0):
    specs = [
        (0.00, 0.00, BLUE, 0.25),
        (0.30, 0.18, WHITE, 0.55),
        (0.60, 0.00, BLUE, 0.25),
        (0.15, 0.38, WHITE, 0.55),
        (0.45, 0.38, BLUE, 0.25),
        (0.75, 0.38, WHITE, 0.55),
    ]
    for ox, oy, fill, tr in specs:
        hexagon = slide.shapes.add_shape(
            MSO_SHAPE.HEXAGON,
            Inches(x + ox * scale),
            Inches(y + oy * scale),
            Inches(0.42 * scale),
            Inches(0.42 * scale),
        )
        set_fill(hexagon, fill, transparency=tr)
        set_line(hexagon, rgb("A8C6D7"), 0.9, 0.2)


def add_footer(slide, text, number, total):
    add_box(slide, 0.55, 7.02, 10.8, 0.18, text, 9, MUTED)
    add_box(slide, 12.18, 6.98, 0.55, 0.22, f"{number}/{total}", 10, NAVY, True, align=PP_ALIGN.RIGHT)


def add_cycle_core(slide, cx, cy):
    outer = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 1.15), Inches(cy - 1.15), Inches(2.3), Inches(2.3))
    set_fill(outer, BLUE, transparency=0.6)
    set_line(outer, rgb("9FC7DF"), 2.0, 0.0)
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.72), Inches(cy - 0.72), Inches(1.44), Inches(1.44))
    set_fill(inner, WHITE)
    set_line(inner, WHITE, 0.5)
    add_box(slide, cx - 0.45, cy - 0.34, 0.9, 0.22, "BNS", 20, NAVY, True, font_name="Aptos Display", align=PP_ALIGN.CENTER)
    add_lines(
        slide,
        cx - 0.62,
        cy - 0.05,
        1.24,
        0.48,
        ["Serien / Cluster", "kurze Spasmen", "kritisch beim Aufwachen"],
        10,
        TEXT,
        align=PP_ALIGN.CENTER,
    )
    positions = [
        (cx - 1.55, cy - 0.15, 310),
        (cx - 0.22, cy - 1.52, 40),
        (cx + 1.2, cy - 0.05, 130),
        (cx - 0.2, cy + 1.2, 220),
    ]
    for x, y, rotation in positions:
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.65), Inches(0.34))
        set_fill(arrow, rgb("8FC6E6"))
        set_line(arrow, rgb("7BB5D8"), 0.8)
        arrow.rotation = rotation


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_fill(bg, color)
    set_line(bg, color, 0)
    add_hex_pattern(slide, 0.2, 0.18, 1.4)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "BNS-Kraempfe im Saeuglingsalter"
    prs.core_properties.subject = "Neue kombinierte PowerPoint aus Leitlinie und Pflegefall"
    prs.core_properties.author = "Codex"
    prs.core_properties.comments = "Neuaufbau aus zwei Quelldecks"

    total = 8

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    side = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.58), Inches(6.15), Inches(5.9))
    set_fill(side, NAVY)
    set_line(side, NAVY, 0.8)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.6), Inches(4.7), Inches(1.55), Inches(1.55))
    set_fill(glow, TEAL, transparency=0.22)
    set_line(glow, TEAL, 0.5, 0.6)
    add_tag(slide, 0.92, 0.98, 2.5, 0.34, "NEUE KOMBINATION AUS 2 QUELLDECKS", TEAL)
    add_box(slide, 0.92, 1.42, 5.25, 1.45, "BNS-Krämpfe im\nSäuglingsalter", 30, WHITE, True, font_name="Aptos Display")
    add_lines(
        slide,
        0.92,
        3.15,
        5.2,
        1.15,
        ["West-Syndrom zwischen Leitlinie,", "Fallsituation und pflegerischer Betreuung"],
        18,
        rgb("E4F2F9"),
        False,
    )
    add_lines(
        slide,
        0.92,
        4.65,
        5.2,
        0.9,
        ["Neue Dramaturgie: Erkennen, Diagnostik, Therapie,", "Pflege, Familie und Nachsorge in einem Deck."],
        13,
        WHITE,
    )
    add_card(
        slide,
        7.25,
        0.95,
        5.1,
        1.18,
        "Quellenbasis",
        [
            "1. BNS-Krämpfe im Säuglingsalter: Fallsituation, ATL, Familienanalyse",
            "2. Therapie der Blitz-Nick-Salaam-Epilepsie: Leitlinie und Medikamente",
        ],
        BLUE,
        body_size=13,
    )
    add_stat(slide, 7.25, 2.42, 1.52, 1.15, "3-18", "typischer Lebensmonat", MINT)
    add_stat(slide, 8.98, 2.42, 1.52, 1.15, "1/2500", "betroffene Säuglinge", APRICOT)
    add_stat(slide, 10.71, 2.42, 1.52, 1.15, "< 3 Wo", "kritisch bis Therapiebeginn", SUN)
    add_card(
        slide,
        7.25,
        4.0,
        5.1,
        2.25,
        "Was im neuen Deck zusammenkommt",
        [
            "• typische BNS-Anfälle und Risikokonstellationen",
            "• EEG, Hypsarrhythmie und weitere Diagnostik",
            "• erste Wahl der Therapie und Evaluation nach 14 Tagen",
            "• Beobachtung, Sicherheit, Elternanleitung und psychosoziale Unterstützung",
        ],
        WHITE,
        body_size=15,
        border=LINE,
    )
    add_footer(slide, "Neue Präsentation auf Basis der beiden Quelldecks", 1, total)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_box(slide, 1.1, 0.18, 11.2, 0.48, "BNS-Anfälle im Überblick", 26, NAVY, True, font_name="Aptos Display", align=PP_ALIGN.CENTER)
    add_card(
        slide,
        3.8,
        0.78,
        5.65,
        1.18,
        "Definition",
        [
            "Epileptische Spasmen im Säuglings- und Kleinkindalter, meist vor dem 2. Lebensjahr.",
            "Typisch sind Blitz-, Nick- und Salaam-Bewegungen in Serien.",
        ],
        BLUE,
        body_size=14,
    )
    add_card(
        slide,
        0.55,
        2.02,
        3.28,
        1.8,
        "Ursachen & Risiken",
        [
            "• Fehlbildungen, Infektion, Blutung, Sauerstoffmangel",
            "• genetische Erkrankungen, Tuberöse Sklerose",
            "• Trisomie 21 und neurologische Vorerkrankungen",
            "• etwa ein Drittel ohne klare Ursache",
        ],
        BLUE,
        body_size=14,
    )
    add_card(
        slide,
        9.52,
        2.02,
        3.28,
        1.8,
        "Symptome & Anfälle",
        [
            "• einzelne Spasmen oft nur etwa 1 Sekunde",
            "• häufig in Serien mit vielen Spasmen",
            "• besonders oft nach dem Aufwachen",
            "• Schreien, Weinen oder Erschöpfung möglich",
        ],
        BLUE,
        body_size=14,
    )
    add_cycle_core(slide, 6.67, 3.9)
    add_card(
        slide,
        0.55,
        4.65,
        3.28,
        1.72,
        "Diagnostik",
        [
            "• EEG mit Wach- und Schlafableitung",
            "• Video durch Eltern sehr hilfreich",
            "• neurologische Untersuchung, MRT, Labor, Genetik",
        ],
        MINT,
        body_size=14,
    )
    add_card(
        slide,
        4.15,
        5.34,
        5.02,
        1.1,
        "Therapie",
        [
            "• Ziel: Anfallsfreiheit + hypsarrhythmiefreies EEG",
            "• erste Wahl: Hormone, oft mit Vigabatrin kombiniert",
            "• Kontrolle nach ca. 14 Tagen",
        ],
        LILAC,
        body_size=14,
    )
    add_card(
        slide,
        9.52,
        4.65,
        3.28,
        1.72,
        "Pflege, Familie & Verlauf",
        [
            "• engmaschige Beobachtung und sichere Umgebung",
            "• Elternanleitung und emotionale Begleitung",
            "• Prognose hängt stark vom frühen Therapiebeginn ab",
        ],
        APRICOT,
        body_size=14,
    )
    add_footer(slide, "Verdichtete Übersichtsfolie aus beiden Quelldecks", 2, total)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_box(slide, 0.7, 0.42, 7.5, 0.42, "Fallsituation: Marvin, 6 Monate", 24, NAVY, True, font_name="Aptos Display")
    add_box(slide, 0.72, 0.88, 7.5, 0.28, "Die Pflegepräsentation bringt die klinische Situation in einen konkreten Versorgungsalltag.", 13, MUTED)
    add_card(
        slide,
        0.72,
        1.42,
        4.05,
        2.4,
        "Klinisches Bild",
        [
            "• typische BNS-Krämpfe in Serien",
            "• Häufung beim Aufwachen und Einschlafen",
            "• Risiko für Beeinträchtigung der Hirnentwicklung",
            "• schnelle Anfallsbeobachtung entscheidet über weitere Schritte",
        ],
        BLUE,
        body_size=16,
    )
    add_card(
        slide,
        4.98,
        1.42,
        3.95,
        2.4,
        "Akute Aufgaben im Setting",
        [
            "• Atmung, Hautfarbe und Vigilanz überwachen",
            "• EEG und weitere Diagnostik vorbereiten",
            "• ärztlich angeordnete Medikamente mitbeobachten",
            "• Belastung für Eltern früh erkennen",
        ],
        MINT,
        body_size=16,
    )
    add_card(
        slide,
        9.15,
        1.42,
        3.45,
        2.4,
        "Warum der Fall wichtig ist",
        [
            "Er übersetzt die Leitlinie in konkrete Pflegehandlungen:",
            "",
            "Beobachtung",
            "Sicherheit",
            "Dokumentation",
            "Anleitung",
        ],
        APRICOT,
        body_size=16,
    )
    ribbon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(4.18), Inches(11.88), Inches(0.52))
    set_fill(ribbon, NAVY)
    set_line(ribbon, NAVY, 0)
    add_box(slide, 0.96, 4.28, 11.0, 0.2, "Zentrale pflegerische Leitfrage: Was muss ich während, zwischen und nach den Krampfserien sicher erfassen und sofort weitergeben?", 13, WHITE)
    add_card(
        slide,
        0.72,
        5.02,
        3.7,
        1.45,
        "Beobachtung",
        [
            "Anfallsdauer, Häufigkeit, Auslöser und Verhalten vor / nach dem Spasmus dokumentieren."
        ],
        WHITE,
        body_size=15,
    )
    add_card(
        slide,
        4.58,
        5.02,
        3.7,
        1.45,
        "Kommunikation",
        [
            "Eltern aktiv nach Veränderungen fragen und Beobachtungen in klare Worte übersetzen."
        ],
        WHITE,
        body_size=15,
    )
    add_card(
        slide,
        8.44,
        5.02,
        4.16,
        1.45,
        "Pflegepriorität",
        [
            "Reizarme, sichere Umgebung und ruhige Orientierung für Kind und Familie schaffen."
        ],
        WHITE,
        body_size=15,
    )
    add_footer(slide, "Quelle: BNS-Krämpfe im Säuglingsalter, Fallsituation und Beobachtungsfokus", 3, total)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_box(slide, 0.7, 0.42, 7.2, 0.42, "Erkennen und diagnostisch absichern", 24, NAVY, True, font_name="Aptos Display")
    add_card(
        slide,
        0.72,
        1.22,
        3.82,
        5.25,
        "Verdacht auf BNS-Epilepsie",
        [
            "• serielle, gleich oder ähnlich aussehende Bewegungsabläufe",
            "• kurze Abstände zwischen den Spasmen",
            "• häufig erste Minuten nach dem Aufwachen",
            "• bei neurologischer Vorerkrankung oder Trisomie 21 besonders aufmerksam sein",
            "",
            "Merke:",
            "Ein Handyvideo kann die Abklärung deutlich beschleunigen.",
        ],
        BLUE,
        body_size=16,
    )
    add_card(
        slide,
        4.76,
        1.22,
        3.92,
        5.25,
        "EEG als Schlüsseluntersuchung",
        [
            "• Diagnose soll innerhalb weniger Tage gesichert werden",
            "• Wach- und Schlaf-EEG, ideal mit Video",
            "• Hypsarrhythmie ist das typische chaotische Muster",
            "• auch nach dem Wecken noch einige Minuten ableiten",
            "",
            "Das EEG ist schmerzlos, aber für die Bestätigung unverzichtbar.",
        ],
        MINT,
        body_size=16,
    )
    add_card(
        slide,
        8.9,
        1.22,
        3.7,
        5.25,
        "Weitere Diagnostik",
        [
            "• klinische und neurologische Untersuchung",
            "• Entwicklungsbeurteilung",
            "• MRT des Gehirns",
            "• Blut- und Urinuntersuchungen",
            "• genetische Diagnostik",
            "• selten Stoffwechsel- oder Liquoruntersuchung",
            "",
            "Differenzialdiagnosen: benigne Myoklonien, Sandifer-Syndrom und andere Bewegungsstörungen.",
        ],
        SUN,
        body_size=15,
    )
    add_footer(slide, "Quelle: Verdachtszeichen, EEG und weitere Diagnostik aus der Leitlinienpräsentation", 4, total)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.68), Inches(0.42), Inches(12.0), Inches(0.78))
    set_fill(header, DEEP)
    set_line(header, DEEP, 0)
    add_box(slide, 0.94, 0.62, 9.2, 0.24, "Therapiepfad nach der Patientenleitlinie", 24, WHITE, True, font_name="Aptos Display")
    add_card(
        slide,
        0.72,
        1.55,
        4.0,
        2.0,
        "Therapieziel",
        [
            "• rasche Anfallsfreiheit für BNS-Anfälle",
            "• Sistieren der Hypsarrhythmie im Wach- und Schlaf-EEG",
            "• bestmögliche Entwicklungsprognose",
        ],
        BLUE,
        body_size=16,
    )
    add_card(
        slide,
        4.92,
        1.55,
        4.0,
        2.0,
        "Erste Wahl",
        [
            "• ACTH oder Prednisolon",
            "• häufig in Kombination mit Vigabatrin",
            "• bei Tuberöser Sklerose oder Kontraindikation gegen Hormone: Vigabatrin primär",
        ],
        LILAC,
        body_size=16,
    )
    add_card(
        slide,
        9.12,
        1.55,
        3.48,
        2.0,
        "Beurteilung nach ca. 14 Tagen",
        [
            "• klinisch: 48 Stunden ohne BNS-Anfälle",
            "• EEG: hypsarrhythmiefrei",
            "• nur weniger Anfälle reicht nicht aus",
        ],
        APRICOT,
        body_size=16,
    )
    lane = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(3.95), Inches(11.88), Inches(1.02))
    set_fill(lane, WHITE)
    set_line(lane, LINE, 1.1)
    points = [
        (1.05, "Tag 0", "Diagnose sichern\nund sofort starten", TEAL),
        (4.2, "1-2 Wochen", "wirksame Erstlinientherapie", TEAL),
        (7.35, "Tag 14", "klinische + EEG-Evaluation", GOLD),
        (10.5, "danach", "wechseln / erweitern,\nfalls nötig", CORAL),
    ]
    for x, top, bottom, color in points:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(4.18), Inches(0.38), Inches(0.38))
        set_fill(dot, color)
        set_line(dot, WHITE, 0.5)
        add_box(slide, x - 0.1, 4.58, 0.6, 0.16, top, 10, NAVY, True, align=PP_ALIGN.CENTER)
        add_lines(slide, x - 0.55, 4.84, 1.45, 0.36, bottom.split("\n"), 11, MUTED, align=PP_ALIGN.CENTER)
    connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.38), Inches(4.35), Inches(9.55), Inches(0.06))
    set_fill(connector, LINE)
    set_line(connector, LINE, 0)
    add_card(
        slide,
        0.72,
        5.42,
        5.8,
        1.18,
        "Wenn die erste Wahl nicht wirkt",
        [
            "Ketogene Diät, Sultiam, Topiramat, Valproat, Zonisamid oder Benzodiazepine kommen als weitere Optionen infrage."
        ],
        MINT,
        body_size=15,
    )
    add_card(
        slide,
        6.76,
        5.42,
        5.84,
        1.18,
        "Früh an Epilepsiechirurgie denken",
        [
            "Besonders bei fokalen Läsionen und therapierefraktärem Verlauf soll früh ein erfahrenes Epilepsiezentrum eingebunden werden."
        ],
        SUN,
        body_size=15,
    )
    add_footer(slide, "Quelle: Therapieziele, Erstlinientherapie und Evaluation nach 14 Tagen", 5, total)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_box(slide, 0.7, 0.42, 7.5, 0.42, "Pflegerische Betreuung im Akut- und Zwischenintervall", 24, NAVY, True, font_name="Aptos Display")
    add_card(
        slide,
        0.72,
        1.22,
        3.8,
        2.35,
        "1. Beobachtung & Dokumentation",
        [
            "• Art, Dauer und Häufigkeit der Krampfserien genau festhalten",
            "• Zusammenhang mit Wach-Schlaf-Rhythmus dokumentieren",
            "• Veränderungen der Reaktion des Kindes beschreiben",
        ],
        BLUE,
        body_size=16,
    )
    add_card(
        slide,
        4.78,
        1.22,
        3.8,
        2.35,
        "2. Sicherheit & Vitalfunktionen",
        [
            "• sichere Lagerung und reizangepasste Umgebung",
            "• Atmung, Hautfarbe und Vigilanz überwachen",
            "• Verletzungsprophylaxe und ruhiges Vorgehen im Anfall",
        ],
        MINT,
        body_size=16,
    )
    add_card(
        slide,
        8.84,
        1.22,
        3.76,
        2.35,
        "3. Essen, Trinken, Schlafen",
        [
            "• Trinkschwäche nach Anfällen beachten",
            "• Zeiten und Belastung anpassen",
            "• Schlafverhalten als Auslöser und Verlaufshinweis erfassen",
        ],
        APRICOT,
        body_size=16,
    )
    add_card(
        slide,
        0.72,
        3.92,
        5.9,
        2.2,
        "4. Medikamentenmanagement & Nebenwirkungen",
        [
            "• Medikamente pünktlich und korrekt verabreichen",
            "• auf Unruhe, Schläfrigkeit, Reizbarkeit, Infektzeichen und Blutdruckprobleme achten",
            "• bei Fieber oder auffälligen Nebenwirkungen früh ärztlich rückmelden",
        ],
        ROSE,
        body_size=16,
    )
    add_card(
        slide,
        6.86,
        3.92,
        5.74,
        2.2,
        "5. Elternanleitung",
        [
            "• Verhalten im Anfall erklären und üben",
            "• Sorgen ernst nehmen und in einfache, klare Informationen übersetzen",
            "• Notfallsituationen, Dokumentation und nächste Schritte gemeinsam besprechen",
        ],
        SUN,
        body_size=16,
    )
    add_footer(slide, "Quelle: ATL-Folien, pflegerische Maßnahmen und Überwachung aus beiden Quelldecks", 6, total)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_box(slide, 0.7, 0.42, 8.0, 0.42, "Familie und interprofessionelle Unterstützung", 24, NAVY, True, font_name="Aptos Display")
    add_card(
        slide,
        0.72,
        1.22,
        4.15,
        5.18,
        "Familienanalyse nach Friedemann",
        [
            "Systemerhaltung:",
            "Der Alltag soll trotz Erkrankung irgendwie weiterlaufen.",
            "",
            "Systemveränderung:",
            "Neue Informationen, neue Rollen und neue Handlungsweisen werden nötig.",
            "",
            "Pflegeauftrag:",
            "Orientierung geben, Unsicherheit abbauen und das familiäre Gleichgewicht stabilisieren.",
        ],
        BLUE,
        body_size=16,
    )
    add_card(
        slide,
        5.08,
        1.22,
        3.6,
        2.45,
        "Was Eltern brauchen",
        [
            "• wiederholte Gespräche in verständlicher Sprache",
            "• genug Raum für Fragen und Ängste",
            "• Befundkopien, schriftliche Informationen und klare Ansprechpartner",
        ],
        MINT,
        body_size=16,
    )
    add_card(
        slide,
        8.92,
        1.22,
        3.68,
        2.45,
        "Hilfreiche Berufsgruppen",
        [
            "• Kinderneuropädiatrie / Epileptologie",
            "• Psychologie",
            "• Sozialdienst / Sozialpädagogik",
            "• Frühförderung und Beratungsangebote",
        ],
        APRICOT,
        body_size=16,
    )
    add_card(
        slide,
        5.08,
        3.98,
        7.52,
        2.42,
        "Praktische Unterstützung im Verlauf",
        [
            "• eigene Akte mit Arztbriefen, Befunden und Beobachtungen anlegen",
            "• Kontakt zu Selbsthilfe und Elternnetzwerken vermitteln",
            "• psychosoziale Belastung aktiv ansprechen, nicht nur medizinische Fakten",
            "• interprofessionell abstimmen: Medizin, Pflege, Entwicklung und Familie gehören zusammen",
        ],
        SUN,
        body_size=16,
    )
    add_footer(slide, "Quelle: Familienanalyse, Elterninformation und Unterstützungsangebote", 7, total)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    closing = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.52), Inches(12.0), Inches(0.76))
    set_fill(closing, NAVY)
    set_line(closing, NAVY, 0)
    add_box(slide, 0.96, 0.71, 10.8, 0.22, "Verlauf, Prognose und Nachsorge", 24, WHITE, True, font_name="Aptos Display")
    add_card(
        slide,
        0.72,
        1.55,
        3.82,
        4.85,
        "Entwicklung",
        [
            "• BNS-Epilepsie kann Entwicklung verlangsamen oder rückläufig machen",
            "• Blickkontakt und Verhalten können sich verändern",
            "• schnelle, erfolgreiche Therapie verbessert die Chancen",
            "• vollständig vorhersehbar ist die Prognose trotzdem nicht",
        ],
        BLUE,
        body_size=16,
    )
    add_card(
        slide,
        4.76,
        1.55,
        3.82,
        4.85,
        "Nachkontrollen",
        [
            "• regelmäßige EEG- und klinische Kontrollen",
            "• standardisierte Entwicklungstests mit etwa 18 Monaten",
            "• weitere Testung vor Schuleintritt",
            "• bei Rückfall erneute Therapieentscheidung",
        ],
        MINT,
        body_size=16,
    )
    add_card(
        slide,
        8.8,
        1.55,
        3.8,
        4.85,
        "Schlussbotschaft",
        [
            "Frühes Erkennen, rasche Diagnostik und eine zügig wirksame Therapie sind der medizinische Kern.",
            "",
            "Genauso wichtig bleiben Beobachtung, Sicherheit und die kontinuierliche Begleitung der Familie.",
        ],
        APRICOT,
        body_size=16,
    )
    add_tag(slide, 4.7, 6.55, 3.95, 0.34, "früh erkennen  |  rasch behandeln  |  Familie mitnehmen", TEAL)
    add_footer(slide, "Kombinierte Abschlussfolie aus Leitlinie, Pflegefall und Verlaufsperspektive", 8, total)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"saved {OUT_PATH}")


if __name__ == "__main__":
    build()
