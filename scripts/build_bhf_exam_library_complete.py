from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("/Users/arielhavana/antigr/porc/mockups/exam_prep_complete")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


def L(text: str):
    return [line.strip() for line in text.strip().splitlines() if line.strip()]


WHITE = rgb("FFFFFF")
INK = rgb("223247")
TEXT = rgb("33465A")
MUTED = rgb("667A8F")
LINE = rgb("D6E0E8")
LIGHT = rgb("F8FBFD")

TONES = {
    "sky": rgb("E7F3FA"),
    "mint": rgb("E7F4EE"),
    "sand": rgb("FCF3DA"),
    "rose": rgb("F9E6EA"),
    "white": WHITE,
}


DECKS = [
    {
        "no": "01",
        "title": "BHF 01 | Lernen, Kommunikation und Berufsrolle",
        "subtitle": "BHF 1a + 1b: Einstieg in Pflege, Rechte, Beobachtung und Teamarbeit",
        "filename": "BHF_01_Lernen_Kommunikation_Berufsrolle.pptx",
        "accent": rgb("1E7A8C"),
        "accent_dark": rgb("175E6D"),
        "topics": L(
            """
            Lernen lernen und Operatoren kennen
            Datenschutz und Schweigepflicht
            Rechte und Pflichten
            Kommunikation und Informationsgespraeche
            Berufsbild und Geschichte der Pflege
            Pflegeprozess
            Wahrnehmen und Beobachten
            Im Pflegealltag mitwirken
            Kollegiale Fallberatung
            Emotion, Empathie und System Familie
            """
        ),
        "focus_cards": [
            {
                "title": "Lernen und Operatoren",
                "tone": "sky",
                "body": L(
                    """
                    Verstehen, erklaeren, begruenden und anwenden klar unterscheiden
                    In der Pruefung erst die Frageart erkennen, dann antworten
                    """
                ),
            },
            {
                "title": "Rechte und Schweigepflicht",
                "tone": "mint",
                "body": L(
                    """
                    Datenschutz, Privatsphaere und Einwilligung sind Basis der Pflege
                    Alles nur weitergeben, was pflegerisch notwendig und erlaubt ist
                    """
                ),
            },
            {
                "title": "Kommunikation und Empathie",
                "tone": "sand",
                "body": L(
                    """
                    aktiv zuhoeren, klar formulieren, Gefuehle benennen
                    empathisch sein ohne die professionelle Rolle zu verlieren
                    """
                ),
            },
            {
                "title": "Beobachtung und Pflegeprozess",
                "tone": "rose",
                "body": L(
                    """
                    erst objektiv beschreiben, dann interpretieren
                    Assessment, Problem, Ziel, Massnahme, Evaluation sauber trennen
                    """
                ),
            },
        ],
        "blocks_title": "Pruefungslogik und Anwendung",
        "blocks_subtitle": "Diese vier Punkte bringen in schriftlicher und muendlicher Pruefung besonders viele sichere Punkte.",
        "blocks": [
            {
                "title": "Pflegeprozess",
                "tone": "sky",
                "body": L(
                    """
                    Definition: systematisches Vorgehen von Assessment bis Evaluation
                    Kern: Daten sammeln, Probleme erkennen, Ziele setzen
                    Praxis: immer Problem + Ziel + Massnahme + Evaluation nennen
                    """
                ),
            },
            {
                "title": "Wahrnehmen vs. Interpretieren",
                "tone": "mint",
                "body": L(
                    """
                    Wahrnehmen: ich sehe, hoere, messe etwas
                    Interpretieren: ich ordne es fachlich ein
                    Merksatz: in der Doku erst Beobachtung, dann Bewertung
                    """
                ),
            },
            {
                "title": "Kollegiale Fallberatung",
                "tone": "sand",
                "body": L(
                    """
                    Ziel: gemeinsam Loesungen finden ohne Schuldzuweisung
                    Struktur: Fall schildern, Fragen klaeren, Optionen sammeln
                    Nutzen: Sicherheit, Reflexion und Teamlernen
                    """
                ),
            },
            {
                "title": "System Familie",
                "tone": "rose",
                "body": L(
                    """
                    Pflege betrifft nie nur eine Person, sondern oft das Umfeld mit
                    Familie kann Ressource, Belastung oder beides sein
                    Kommunikation immer systemisch und respektvoll denken
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Pflegeprozess in eigenen Worten erklaeren
            Schweigepflicht, Datenschutz und Einwilligung unterscheiden
            Beobachtung von Interpretation trennen
            Empathie professionell beschreiben
            Kollegiale Fallberatung einfach darstellen
            """
        ),
        "check_right": L(
            """
            In Antworten zuerst definieren, dann Beispiel geben
            Rechte, Sicherheit und Wuerde fast immer mitdenken
            Kurze, klare Saetze schlagen Fachwortketten
            Wenn unsicher: Was ist beobachtbar, was ist pflegerisch relevant?
            """
        ),
    },
    {
        "no": "02",
        "title": "BHF 02 | Bewegung, Prophylaxen und Haftung",
        "subtitle": "Beweglichkeit, Sturz, Kontraktur, Osteoporose, Anleitung und Haftungsrecht",
        "filename": "BHF_02_Bewegung_Prophylaxen_Haftung.pptx",
        "accent": rgb("2F6F9F"),
        "accent_dark": rgb("214F72"),
        "topics": L(
            """
            Beweglichkeit unterstuetzen
            Sturzprophylaxe
            Erste Hilfe bei Sturz
            Kontraktur
            Osteoporose
            Orientierung unterstuetzen
            Anleiten und beraten
            Persoenliche Gesundheit erhalten
            Haftungsrecht
            """
        ),
        "focus_cards": [
            {
                "title": "Mobilitaet foerdern",
                "tone": "sky",
                "body": L(
                    """
                    Beweglichkeit erhaelt Selbststaendigkeit und verhindert Komplikationen
                    Transfer, Hilfsmittel, Schmerz und Sicherheit immer zusammen denken
                    """
                ),
            },
            {
                "title": "Sturz und Erste Hilfe",
                "tone": "mint",
                "body": L(
                    """
                    nach Sturz zuerst Verletzung, Bewusstsein und Schmerzen pruefen
                    sichere Umgebung und ruhige Kommunikation sind Pflicht
                    """
                ),
            },
            {
                "title": "Anleitung und Beratung",
                "tone": "sand",
                "body": L(
                    """
                    nur erklaeren reicht nicht, Patient muss es nachvollziehen koennen
                    Sprache, Tempo und Vorerfahrung anpassen
                    """
                ),
            },
            {
                "title": "Haftung und Verantwortung",
                "tone": "rose",
                "body": L(
                    """
                    wer Risiken erkennt, muss angemessen handeln und dokumentieren
                    Unterlassen kann haftungsrelevant sein
                    """
                ),
            },
        ],
        "blocks_title": "Krankheitslehre und Transfer",
        "blocks_subtitle": "Bei BHF 2 geht es darum, Risiken zu erkennen und daraus sichere Pflege abzuleiten.",
        "blocks": [
            {
                "title": "Osteoporose",
                "tone": "sky",
                "body": L(
                    """
                    D: verminderte Knochendichte mit erhoehter Frakturgefahr
                    F/T: primaer oder sekundaer, oft lange unbemerkt
                    U + R: Alter, Bewegungsmangel, Mangelernahrung, Kortison
                    S: Rueckenschmerz, Groessenverlust, Frakturen
                    T/P: Bewegung, Kalzium/Vitamin D, Sturzschutz, Schmerz- und Mobilitaetsmanagement
                    """
                ),
            },
            {
                "title": "Sturzprophylaxe",
                "tone": "mint",
                "body": L(
                    """
                    Ziel: Sturzrisiko systematisch reduzieren
                    Risiken: Schwindel, Medikamente, Sehprobleme, Unsicherheit
                    Pflege: Umgebung sichern, Hilfsmittel, Anleitung, Rufsystem, Doku
                    """
                ),
            },
            {
                "title": "Kontrakturprophylaxe",
                "tone": "sand",
                "body": L(
                    """
                    Definition: Verkuerzung von Muskeln und Bewegungseinschraenkung
                    Risiken: Immobilitaet, Schmerzen, Schonhaltung
                    Pflege: Mobilisieren, lagern, anleiten, Schmerzen ernst nehmen
                    """
                ),
            },
            {
                "title": "Orientierung und Sicherheit",
                "tone": "rose",
                "body": L(
                    """
                    bei Verwirrung oder Unsicherheit Umgebung klar strukturieren
                    Namen, Tageszeit, Ort und Handlung immer mitbenennen
                    Sicherheit geht vor Tempo
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Sturzprophylaxe mit Beispiel erklaeren
            Osteoporose kurz im DURST-Schema darstellen
            Kontraktur und Immobilitaet verknuepfen
            Anleitung patientengerecht beschreiben
            Haftung in einem Pflegefall begruenden
            """
        ),
        "check_right": L(
            """
            Immer Risiko + Beobachtung + konkrete Massnahme nennen
            Bei Sturz zuerst sichern, dann mobilisieren
            Beratung ohne Verstaendnissicherung ist unvollstaendig
            Doku schuetzt Patient und Pflegekraft
            """
        ),
    },
    {
        "no": "03",
        "title": "BHF 03 | Hygiene, Haut, Ernaehrung und Pflegeethik",
        "subtitle": "BHF 3a + 3b: hygienisches Handeln, Hautbeobachtung, Koerperpflege und Ethik",
        "filename": "BHF_03_Hygiene_Haut_Ernaehrung_Ethik.pptx",
        "accent": rgb("22806D"),
        "accent_dark": rgb("196252"),
        "topics": L(
            """
            Im Pflegealltag handeln
            Grundlagen hygienischen Handelns
            Haut und Hautbeobachtung
            Dekubitus
            Koerperpflege
            Mundpflege
            Interaktion bei koerpernaher Pflege
            Grundlagen von Ernaehrung und Nahrungsaufnahme
            Recht und Pflegeethik
            Familie und Kultursensibilitaet
            """
        ),
        "focus_cards": [
            {
                "title": "Hygiene als Sicherheitsprinzip",
                "tone": "sky",
                "body": L(
                    """
                    Haende, Material und Arbeitsflaechen konsequent mitdenken
                    Hygiene ist Infektionsschutz und Qualitaet zugleich
                    """
                ),
            },
            {
                "title": "Haut und Dekubitus",
                "tone": "mint",
                "body": L(
                    """
                    Hautbeobachtung ist ein fruehes Warnsystem
                    Druck, Feuchtigkeit, Mangelernaehrung und Immobilitaet sind Schluesselrisiken
                    """
                ),
            },
            {
                "title": "Koerpernahe Pflege",
                "tone": "sand",
                "body": L(
                    """
                    Intimsphaere schuetzen, Schritte erklaeren, Mitmachen foerdern
                    Koerperpflege ist zugleich Beziehung und Assessment
                    """
                ),
            },
            {
                "title": "Ernaehrung und Ethik",
                "tone": "rose",
                "body": L(
                    """
                    Essen ist nicht nur Nahrungsaufnahme, sondern auch Biografie und Kultur
                    Ethik wird sichtbar bei Abhaengigkeit, Scham und Selbstbestimmung
                    """
                ),
            },
        ],
        "blocks_title": "Dekubitus, Ernaehrung und Pflegebeziehung",
        "blocks_subtitle": "Hier werden oft praktische Alltagssituationen mit Beobachtung und Begruendung abgefragt.",
        "blocks": [
            {
                "title": "Dekubitus",
                "tone": "sky",
                "body": L(
                    """
                    D: lokale Gewebeschaedigung durch Druck oder Druck + Scherung
                    Grade: Stadium 1 bis 4 je nach Tiefe des Schadens
                    U + R: Immobilitaet, Feuchtigkeit, Mangelernaehrung, Sensibilitaetsstoerung
                    S: Roetung, Schmerz, Hautdefekt, Nekrose
                    T/P: Druckentlastung, Lagerung, Hautpflege, Ernaehrung, Wundmanagement
                    """
                ),
            },
            {
                "title": "Mangelernaehrung",
                "tone": "mint",
                "body": L(
                    """
                    Definition: unzureichende Energie- oder Nahrstoffzufuhr
                    Risiken: Alter, Krankheit, Kau-/Schluckprobleme, Appetitmangel
                    Pflege: Essprotokoll, Lieblingsspeisen, Hilfen anbieten, Gewicht beobachten
                    """
                ),
            },
            {
                "title": "Interaktion in Intimnaehe",
                "tone": "sand",
                "body": L(
                    """
                    vorher anklaeren, Einwilligung holen, Grenzen respektieren
                    ruhige Sprache und Sichtschutz geben Sicherheit
                    Beruehrung immer professionell und transparent gestalten
                    """
                ),
            },
            {
                "title": "Pflegeethik im Alltag",
                "tone": "rose",
                "body": L(
                    """
                    Kernfragen: Was ist gut? Was will der Patient? Was ist sicher?
                    Autonomie, Fuersorge und Nicht-Schaden abwaegen
                    Kultursensibel handeln ohne zu stereotypisieren
                    """
                ),
            },
        ],
        "check_left": L(
            """
            5 Momente der Haendedesinfektion kennen
            Dekubitus im DURST-Schema erklaeren
            Koerperpflege als Assessment begruenden
            Mangelernaehrung frueh erkennen
            Ethik und Kultursensibilitaet einfach erklaeren
            """
        ),
        "check_right": L(
            """
            Bei Haut immer Farbe, Temperatur, Feuchtigkeit und Druckstellen nennen
            Bei Pflegehandlungen immer Intimsphaere und Einwilligung mitdenken
            Ethik wird konkret im Einzelfall, nicht nur in Definitionen
            """
        ),
    },
    {
        "no": "04",
        "title": "BHF 04 | Schwangerschaft, Geburt und Neugeborene",
        "subtitle": "Anatomie, Geburt, Versorgung des Neugeborenen und familienorientierte Pflege",
        "filename": "BHF_04_Schwangerschaft_Geburt_Neugeborene.pptx",
        "accent": rgb("D06D54"),
        "accent_dark": rgb("A25340"),
        "topics": L(
            """
            Einstieg und Anatomie
            Schwangerschaft und Geburt
            Versorgung des Neugeborenen
            Unterstuetzungssysteme
            Familie und Sicherheit in der fruehen Phase
            """
        ),
        "focus_cards": [
            {
                "title": "Anatomie verstehen",
                "tone": "sky",
                "body": L(
                    """
                    Koerperkenntnis ist Grundlage fuer Beobachtung und Erklaerung
                    normale Veraenderungen muessen von Warnzeichen getrennt werden
                    """
                ),
            },
            {
                "title": "Schwangerschaft beobachten",
                "tone": "mint",
                "body": L(
                    """
                    Wohlbefinden von Mutter und Kind sichern
                    Schmerzen, Blutung, Oedeme und RR-Anstieg immer ernst nehmen
                    """
                ),
            },
            {
                "title": "Neugeborenes versorgen",
                "tone": "sand",
                "body": L(
                    """
                    Atmung, Temperatur, Hautfarbe, Trinkverhalten beobachten
                    Elternanleitung ist Teil der Pflege
                    """
                ),
            },
            {
                "title": "Unterstuetzungssysteme",
                "tone": "rose",
                "body": L(
                    """
                    Hebamme, Familie, Klinik, Nachsorge und soziale Ressourcen kennen
                    frueh Hilfen vermitteln statt spaet Krisen verwalten
                    """
                ),
            },
        ],
        "blocks_title": "Pruefungstransfer fuer die fruehe Familienphase",
        "blocks_subtitle": "In diesem BHF zaehlen vor allem sichere Beobachtung, ruhige Anleitung und Red-Flag-Denken.",
        "blocks": [
            {
                "title": "Schwangerschaft",
                "tone": "sky",
                "body": L(
                    """
                    Normale Veraenderungen: Uebelkeit, Muedigkeit, Gewichtszunahme
                    Red Flags: Blutung, starke Schmerzen, Oedeme + RR-Anstieg
                    Pflege: informieren, beobachten, weiterleiten
                    """
                ),
            },
            {
                "title": "Geburt",
                "tone": "mint",
                "body": L(
                    """
                    Mutter und Kind gleichzeitig im Blick behalten
                    Kommunikation ruhig und klar halten
                    Sicherheit, Schmerzbeobachtung und Begleitung sind zentral
                    """
                ),
            },
            {
                "title": "Neugeborenes",
                "tone": "sand",
                "body": L(
                    """
                    Atmung, Temperatur und Hautfarbe zuerst
                    Trinkverhalten, Ausscheidung und Bindung mitbeobachten
                    Eltern anleiten statt nur uebernehmen
                    """
                ),
            },
            {
                "title": "Familienpflege",
                "tone": "rose",
                "body": L(
                    """
                    Eltern sind Partner der Versorgung
                    Unsicherheit normalisieren und Handlungssicherheit geben
                    Ressourcen und Unterstuetzungsangebote sichtbar machen
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Wichtige Warnzeichen in der Schwangerschaft nennen
            Versorgung eines Neugeborenen erklaeren
            Elternanleitung als Pflegeaufgabe begruenden
            Unterstuetzungssysteme aufzählen
            """
        ),
        "check_right": L(
            """
            Immer Mutter und Kind gemeinsam denken
            Beobachtung vor Bewertung
            Bei Red Flags frueh weiterleiten
            """
        ),
    },
    {
        "no": "05",
        "title": "BHF 05 | Krankenhaus, OP und Schmerz",
        "subtitle": "Aufnahme, perioperative Pflege, Schmerz, Wunden, Hygiene und Arzneimittel",
        "filename": "BHF_05_Krankenhaus_OP_Schmerz.pptx",
        "accent": rgb("8C6C31"),
        "accent_dark": rgb("674F24"),
        "topics": L(
            """
            Pflegeprozess bei Aufnahme
            Praeoperative Pflege
            Postoperative Pflege
            Thromboseprophylaxe
            Management des akuten Schmerzes
            Akute Wunden
            Krankenhaushygiene
            Organisation und Oekonomie
            Kinder im Krankenhaus
            Ambulante OP
            Arzneimittelsicherheit
            """
        ),
        "focus_cards": [
            {
                "title": "Aufnahme und Vorbereitung",
                "tone": "sky",
                "body": L(
                    """
                    Informationen, Risiken und Ressourcen frueh erfassen
                    OP-Vorbereitung immer sicher, nachvollziehbar und dokumentiert
                    """
                ),
            },
            {
                "title": "Post-OP beobachten",
                "tone": "mint",
                "body": L(
                    """
                    Atmung, Kreislauf, Schmerzen, Uebelkeit, Blutung und Mobilitaet
                    fruehe Komplikationen aktiv suchen
                    """
                ),
            },
            {
                "title": "Schmerz und Wunde",
                "tone": "sand",
                "body": L(
                    """
                    Schmerz ist subjektiv und trotzdem messbar
                    Wundbeobachtung gehoert zur Sicherheitsarbeit
                    """
                ),
            },
            {
                "title": "Kinder und ambulante OP",
                "tone": "rose",
                "body": L(
                    """
                    kurze Wege, gute Aufklaerung und Elternarbeit sind entscheidend
                    Angst und Orientierung immer mitdenken
                    """
                ),
            },
        ],
        "blocks_title": "Kernthemen fuer die OP-Pflege",
        "blocks_subtitle": "Bei BHF 5 wird haeufig abgefragt, was vor und nach einer OP pflegerisch priorisiert wird.",
        "blocks": [
            {
                "title": "Thromboseprophylaxe",
                "tone": "sky",
                "body": L(
                    """
                    Ziel: Venenthrombose und Embolie verhindern
                    Risiken: Immobilitaet, OP, Alter, Adipositas, Krebs
                    Pflege: Mobilisation, Kompression, Fluessigkeit, Verordnung beachten
                    """
                ),
            },
            {
                "title": "Akuter Schmerz",
                "tone": "mint",
                "body": L(
                    """
                    D: ploetzlich einsetzender Schmerz mit Schutz- und Warnfunktion
                    Assessment: Ort, Staerke, Qualität, Ausloeser, Wirkung
                    T/P: Analgetika, Lagerung, Ruhe, Evaluation, Schmerzskala nutzen
                    """
                ),
            },
            {
                "title": "Akute Wunde",
                "tone": "sand",
                "body": L(
                    """
                    Beobachte: Roetung, Schwellung, Sekret, Schmerz, Geruch
                    Aseptisches Arbeiten ist Pflicht
                    Pflege: Verband, Doku, Infektionszeichen und Heilungsverlauf
                    """
                ),
            },
            {
                "title": "Arzneimittelsicherheit",
                "tone": "rose",
                "body": L(
                    """
                    6-R-Regel + Wirkung + Nebenwirkung
                    OP-Umfeld braucht klare Uebergaben und Kontrollen
                    bei Unsicherheit lieber stoppen und nachfragen
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Prae- und postoperative Pflege unterscheiden
            Thromboseprophylaxe konkret erklaeren
            Schmerzassessment beschreiben
            Wundbeobachtung strukturieren
            Arzneimittelsicherheit begruenden
            """
        ),
        "check_right": L(
            """
            Post-OP immer an Blutung, Atmung, Kreislauf und Schmerz denken
            Hygiene und Doku sind Sicherheitsinstrumente
            Kinderpflege braucht Elternarbeit und Angstabbau
            """
        ),
    },
    {
        "no": "06",
        "title": "BHF 06 | Notfaelle, Schock und Blut",
        "subtitle": "Vitalzeichen, Akutversorgung, Schock, Anaphylaxie, COPD, Schlaganfall, Blut und Injektionen",
        "filename": "BHF_06_Notfaelle_Schock_Blut.pptx",
        "accent": rgb("C05D43"),
        "accent_dark": rgb("934734"),
        "topics": L(
            """
            Vitalzeichen messen
            Kardiogener Notfall
            Menschen im Schock pflegen
            Notfaelle in Pflege und Reha
            Akutes Abdomen
            Anaphylaxie
            Exazerbierte COPD
            Schlaganfall
            Verbrennung
            Aufgaben des Blutes
            Erythrozyten, Leukozyten, Thrombozyten
            Injektionen
            Rechtliche Grundlagen der Notfallversorgung
            """
        ),
        "focus_cards": [
            {
                "title": "Vitalzeichen und Trends",
                "tone": "sky",
                "body": L(
                    """
                    RR, Puls, AF, SpO2, Temp und Bewusstsein immer zusammen lesen
                    Trends sind oft wichtiger als Einzelwerte
                    """
                ),
            },
            {
                "title": "ABCDE und Schockdenken",
                "tone": "mint",
                "body": L(
                    """
                    Akut lebensbedrohliches zuerst sichern
                    Atemweg, Atmung und Kreislauf priorisieren
                    """
                ),
            },
            {
                "title": "Blut und Injektionen",
                "tone": "sand",
                "body": L(
                    """
                    Blutbestandteile einfach, aber sicher erklaeren koennen
                    Injektionen immer aseptisch, korrekt und beobachtend
                    """
                ),
            },
            {
                "title": "Recht in Notfaellen",
                "tone": "rose",
                "body": L(
                    """
                    Hilfe holen, dokumentieren, Notkompetenz und Standards beachten
                    in Akutsituationen klar kommunizieren
                    """
                ),
            },
        ],
        "blocks_title": "DURST-Kurzformat fuer zentrale Notfaelle",
        "blocks_subtitle": "Wenn Krankheitslehre gefragt wird, antworte strukturiert und prioritaetsorientiert.",
        "blocks": [
            {
                "title": "Anaphylaxie",
                "tone": "sky",
                "body": L(
                    """
                    D: akute schwere allergische Sofortreaktion
                    Grade: Haut bis Kreislaufstillstand
                    U + R: Nahrungsmittel, Medikamente, Insekten, Asthma
                    S: Urtikaria, Atemnot, Hypotonie, Angst
                    T/P: Notruf, Lagerung, O2, Adrenalin, Monitoring, beruhigen
                    """
                ),
            },
            {
                "title": "COPD-Exazerbation",
                "tone": "mint",
                "body": L(
                    """
                    D: akute Verschlechterung der COPD-Symptomatik
                    Formen: infektiös oder nicht-infektioes getriggert
                    U + R: Rauchen, Infekt, Belastung, unzureichende Therapie
                    S: Atemnot, Husten, Sputum, Tachypnoe
                    T/P: O2 nach Standard, Inhalation, Lagerung, Atemerleichterung, Beobachtung
                    """
                ),
            },
            {
                "title": "Schlaganfall",
                "tone": "sand",
                "body": L(
                    """
                    D: akute Durchblutungsstoerung oder Blutung im Gehirn
                    Typen: ischaemisch oder haemorrhagisch
                    U + R: Hypertonie, Vorhofflimmern, Rauchen, Diabetes
                    S: FAST, Hemiparese, Sprachstoerung, Sehverlust
                    T/P: Stroke Unit, Aspirationsschutz, Monitoring, schnelle Eskalation
                    """
                ),
            },
            {
                "title": "Schock",
                "tone": "rose",
                "body": L(
                    """
                    D: kritische Minderdurchblutung mit Organunterversorgung
                    Formen: hypovolaem, kardiogen, distributiv, obstruktiv
                    S: Tachykardie, Hypotonie, Kaltschweissigkeit, Unruhe
                    Pflege: Hilfe holen, engmaschig beobachten, Lagerung, Waerme, Ruhe
                    """
                ),
            },
        ],
        "check_left": L(
            """
            ABCDE kurz erklaeren
            Anaphylaxie, Schock und Schlaganfall im DURST-Schema nennen
            Blutbestandteile einfach erklaeren
            Injektionen sicher beschreiben
            """
        ),
        "check_right": L(
            """
            Akut geht immer vor vollstaendig
            Im Zweifel zuerst Lebensgefahr benennen
            Pflege heisst hier: beobachten, sichern, eskalieren, dokumentieren
            """
        ),
    },
    {
        "no": "07",
        "title": "BHF 07 | Sinne, ambulante Versorgung und chronische Wunden",
        "subtitle": "Seh-/Hoerbeeintraechtigung, Diabetes im Alltag, Wundversorgung und ambulante Pflege",
        "filename": "BHF_07_Sinne_Ambulante_Versorgung_Wunden.pptx",
        "accent": rgb("4B8B61"),
        "accent_dark": rgb("386749"),
        "topics": L(
            """
            Bei Seh- und Hoereinschraenkungen pflegen
            Menschen mit Diabetes pflegen
            Chronische Wunden versorgen
            Ambulante Versorgung gestalten
            Ambulante OP unterstuetzen
            Gesundheitsfoerderung
            Rechtliche Grundlagen
            """
        ),
        "focus_cards": [
            {
                "title": "Sinnesbeeintraechtigung",
                "tone": "sky",
                "body": L(
                    """
                    Orientierung, Sicherheit und Kommunikation anpassen
                    Hilfsmittel und Umgebung aktiv nutzen
                    """
                ),
            },
            {
                "title": "Chronische Versorgung",
                "tone": "mint",
                "body": L(
                    """
                    ambulante Pflege ist planend, beratend und ressourcenorientiert
                    Familie und Alltag gehoeren mit in den Pflegeplan
                    """
                ),
            },
            {
                "title": "Wunden und Diabetes",
                "tone": "sand",
                "body": L(
                    """
                    Stoffwechsel und Wundheilung haengen eng zusammen
                    Beobachtung und Anleitung sind Schluessel
                    """
                ),
            },
            {
                "title": "Recht und Gesundheit",
                "tone": "rose",
                "body": L(
                    """
                    ambulant braucht klare Absprachen, Datenschutz und Verordnungssicherheit
                    Gesundheitsfoerderung muss alltagstauglich sein
                    """
                ),
            },
        ],
        "blocks_title": "Pruefungstransfer in der ambulanten Pflege",
        "blocks_subtitle": "BHF 7 verbindet Alltag, Chronizitaet und patientenzentrierte Kommunikation.",
        "blocks": [
            {
                "title": "Chronische Wunde",
                "tone": "sky",
                "body": L(
                    """
                    D: Wunde ohne regelrechte Heilung ueber laengere Zeit
                    Typen: z. B. Dekubitus, Ulcus cruris, diabetisches Fusssyndrom
                    U + R: Durchblutungsstoerung, Druck, Diabetes, Mangelernaehrung
                    S: langes Bestehen, Exsudat, Schmerz, Infektzeichen
                    T/P: Ursache behandeln, aseptisch versorgen, Druck entlasten, Anleitung
                    """
                ),
            },
            {
                "title": "Diabetes im Alltag",
                "tone": "mint",
                "body": L(
                    """
                    Pflege zielt auf Selbstmanagement und Beobachtung
                    Hypo- und Hyperglykämie frueh erkennen
                    Fuesse, Haut, Essen und Bewegung mitdenken
                    """
                ),
            },
            {
                "title": "Sehen und Hoeren",
                "tone": "sand",
                "body": L(
                    """
                    Ansprache vor Beruehrung, Umfeld erklaeren, Stolperfallen beseitigen
                    Sprache klar, Blickkontakt und Hilfsmittel konsequent nutzen
                    """
                ),
            },
            {
                "title": "Ambulante Versorgung",
                "tone": "rose",
                "body": L(
                    """
                    Pflege findet im Lebensumfeld des Patienten statt
                    Ziel: Selbststaendigkeit erhalten, Risiken vermeiden, Angehoerige anleiten
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Chronische Wunde im DURST-Schema erklaeren
            Pflege bei Seh- und Hoerbeeintraechtigung beschreiben
            Besonderheiten der ambulanten Versorgung nennen
            Gesundheitsfoerderung praxisnah erklaeren
            """
        ),
        "check_right": L(
            """
            Im ambulanten Bereich sind Anleitung und Umfeld besonders wichtig
            Sicherheit und Selbststaendigkeit immer zusammen denken
            Wunden nie isoliert, sondern ursachenbezogen betrachten
            """
        ),
    },
    {
        "no": "08",
        "title": "BHF 08 | Kinder mit Asthma, Fieber und Beratungsbedarf",
        "subtitle": "Kinderpflege, Fall Henriette Schulz, Fieber, kollegiale Beratung und Gesundheitsfoerderung",
        "filename": "BHF_08_Kinder_Asthma_Fieber_Beratung.pptx",
        "accent": rgb("5D7DB8"),
        "accent_dark": rgb("46608C"),
        "topics": L(
            """
            Wiederholung Anatomie und Physiologie
            Ein Kind mit Asthma pflegen
            Fall Henriette Schulz
            Menschen mit Fieber pflegen
            Kollegiale Beratung
            Gesundheitsfoerderung
            Krankenhausfinanzierung
            """
        ),
        "focus_cards": [
            {
                "title": "Paediatrische Pflege",
                "tone": "sky",
                "body": L(
                    """
                    Kind und Eltern immer gemeinsam im Blick
                    Angstabbau und kindgerechte Erklaerung sind zentral
                    """
                ),
            },
            {
                "title": "Asthma und Fieber",
                "tone": "mint",
                "body": L(
                    """
                    Atmung, Allgemeinzustand und Fluessigkeit eng beobachten
                    Symptome koennen sich bei Kindern rasch verschlechtern
                    """
                ),
            },
            {
                "title": "Falldenken",
                "tone": "sand",
                "body": L(
                    """
                    Henriette Schulz zeigt, wie Beobachtung und Beratung zusammenwirken
                    erst Situation ordnen, dann Pflege planen
                    """
                ),
            },
            {
                "title": "Team und System",
                "tone": "rose",
                "body": L(
                    """
                    Kollegiale Beratung und Ressourcen im System nutzen
                    Gesundheit und Organisation gehoeren mit zur Versorgung
                    """
                ),
            },
        ],
        "blocks_title": "Krankheitslehre und Falllogik",
        "blocks_subtitle": "Asthma und Fieber sind klassische Prüfungsthemen, weil hier Beobachtung, Elternarbeit und Eskalation zusammenspielen.",
        "blocks": [
            {
                "title": "Asthma bronchiale",
                "tone": "sky",
                "body": L(
                    """
                    D: chronisch entzuendliche Atemwegserkrankung mit reversibler Obstruktion
                    Formen: allergisch, nicht-allergisch, gemischt
                    U + R: Allergene, Infekte, Belastung, Rauch
                    S: Giemen, Husten, Atemnot, verlängerte Exspiration
                    T/P: Inhalation, Trigger meiden, Atembeobachtung, Eltern anleiten
                    """
                ),
            },
            {
                "title": "Fieber",
                "tone": "mint",
                "body": L(
                    """
                    D: erhoehte Koerpertemperatur als Symptom
                    Ursachen: meist Infektion, auch Entzuendung oder andere Trigger
                    Risiken: Exsikkose, Krampfneigung, Leistungseinbruch
                    Pflege: Temperatur, Fluessigkeit, Verhalten und Kreislauf beobachten
                    """
                ),
            },
            {
                "title": "Kollegiale Beratung",
                "tone": "sand",
                "body": L(
                    """
                    bei komplexen Kinderfaellen gemeinsam handeln statt allein improvisieren
                    Fragen klaeren, Optionen sammeln, Verantwortung sauber verteilen
                    """
                ),
            },
            {
                "title": "Gesundheitsfoerderung",
                "tone": "rose",
                "body": L(
                    """
                    Elternkompetenz staerken, Alltagsstrategien vermitteln
                    nicht nur Krankheit behandeln, sondern Gesundheit foerdern
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Asthma im DURST-Schema erklaeren
            Fieber als Symptom einordnen
            Elternarbeit in der Kinderpflege begruenden
            Kollegiale Beratung beschreiben
            """
        ),
        "check_right": L(
            """
            Bei Kindern veraendern sich Zustaende oft schnell
            Beobachtung und Anleitung gehoeren zusammen
            Fallbeispiele immer strukturiert beantworten
            """
        ),
    },
    {
        "no": "09",
        "title": "BHF 09 | Neurologie, Wachkoma und Rehabilitation",
        "subtitle": "Neurophysiologie, Schlaganfall beim Kind, neurologische Pflege und Wachkoma",
        "filename": "BHF_09_Neurologie_Wachkoma_Rehabilitation.pptx",
        "accent": rgb("4A69A7"),
        "accent_dark": rgb("374F7D"),
        "topics": L(
            """
            Neurophysiologie
            Schlaganfall beim Kind
            Neurologische Pflege
            Rehabilitative Pflege
            Ethik und Wachkoma
            Evidenz
            Neuropsychologische Stoerungen
            """
        ),
        "focus_cards": [
            {
                "title": "Neuro-Basics",
                "tone": "sky",
                "body": L(
                    """
                    Funktion des Nervensystems verstehen, um Symptome einzuordnen
                    kleine Ausfaelle koennen grosse Bedeutung haben
                    """
                ),
            },
            {
                "title": "Neurologische Beobachtung",
                "tone": "mint",
                "body": L(
                    """
                    Bewusstsein, Pupillen, Motorik, Sprache, Schlucken, Verhalten
                    Veraenderungen frueh merken und sofort weitergeben
                    """
                ),
            },
            {
                "title": "Reha und Funktion",
                "tone": "sand",
                "body": L(
                    """
                    alltagsnahe Ziele, kleine Schritte und Wiederholung
                    Ressourcen sichtbar machen und trainieren
                    """
                ),
            },
            {
                "title": "Ethik im Wachkoma",
                "tone": "rose",
                "body": L(
                    """
                    Würde, Kommunikation, Basale Stimulation und Teamabstimmung
                    ethische Entscheidungen nie isoliert treffen
                    """
                ),
            },
        ],
        "blocks_title": "Neurologische Krankheitslehre im Kurzformat",
        "blocks_subtitle": "Fuer die Pruefung zaehlt, neurologische Symptome mit Pflegefolgen zu verbinden.",
        "blocks": [
            {
                "title": "Schlaganfall",
                "tone": "sky",
                "body": L(
                    """
                    D: akute cerebrovaskulaere Stoerung
                    Typen: ischaemisch oder haemorrhagisch
                    S: Paresen, Sprachstoerung, Bewusstseinsaenderung, Schluckstoerung
                    Pflege: Monitoring, Aspirationsschutz, Fruehmobilisation, Reha-Ziele
                    """
                ),
            },
            {
                "title": "Wachkoma",
                "tone": "mint",
                "body": L(
                    """
                    D: schwere Bewusstseinsstoerung mit Wachphasen ohne klare Kommunikation
                    Beobachtung: Reize, Reaktionen, Tonus, Atmung, Lagerung
                    Pflege: basale Stimulation, Dekubitusprophylaxe, Mundpflege, Teamarbeit
                    """
                ),
            },
            {
                "title": "Neurologische Pflege",
                "tone": "sand",
                "body": L(
                    """
                    sicher lagern, schlucken beachten, Kommunikation anpassen
                    Wahrnehmung, Mobilitaet und Selbstversorgung foerdern
                    """
                ),
            },
            {
                "title": "Evidenz in der Pflege",
                "tone": "rose",
                "body": L(
                    """
                    Entscheidungen sollen nicht nur Gewohnheit, sondern auch Wissen folgen
                    Leitlinien, Standards und Beobachtungen zusammendenken
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Neurologische Beobachtung erklaeren
            Schlaganfall und Reha verknuepfen
            Wachkoma pflegerisch beschreiben
            Evidenz in der Pflege einfach erklaeren
            """
        ),
        "check_right": L(
            """
            In Neuro-Faellen haben kleine Veraenderungen Gewicht
            Schlucken und Kommunikation nie vergessen
            Reha bedeutet alltagsnahe Funktion, nicht nur Uebung
            """
        ),
    },
    {
        "no": "10",
        "title": "BHF 10 | Langzeitpflege, Assessment und Gewalt erkennen",
        "subtitle": "Organisation, Finanzierung, Biografiearbeit, Pflegediagnostik, Hygiene und Gewalt in der LP",
        "filename": "BHF_10_Langzeitpflege_Assessment_Gewalt.pptx",
        "accent": rgb("7A8B39"),
        "accent_dark": rgb("5A672B"),
        "topics": L(
            """
            Organisation in der Langzeitpflege
            Finanzierung von Betreuung
            Biografieorientiert pflegen
            Pflegeprozess in der Langzeitpflege
            Assessment und Pflegediagnostik
            Hygienisch sicher handeln
            Gewaltausuebungen erkennen
            """
        ),
        "focus_cards": [
            {
                "title": "Langzeitpflege verstehen",
                "tone": "sky",
                "body": L(
                    """
                    Tagesstruktur, Beziehung und Ressourcen stehen im Vordergrund
                    Langfristige Stabilitaet ist wichtiger als schnelle Einzelaktionen
                    """
                ),
            },
            {
                "title": "Biografie und Assessment",
                "tone": "mint",
                "body": L(
                    """
                    Biografie hilft Verhalten zu verstehen und Pflege anzupassen
                    Assessment macht Probleme, Risiken und Ressourcen sichtbar
                    """
                ),
            },
            {
                "title": "Hygiene in der LP",
                "tone": "sand",
                "body": L(
                    """
                    auch in Routinen bleibt Hygiene ein Sicherheitskern
                    Wiederholung darf nicht zu Nachlaessigkeit fuehren
                    """
                ),
            },
            {
                "title": "Gewalt erkennen",
                "tone": "rose",
                "body": L(
                    """
                    Gewalt kann offen, subtil oder strukturell sein
                    Ueberforderung und Routineblindheit sind echte Risiken
                    """
                ),
            },
        ],
        "blocks_title": "Pruefungstransfer fuer den Langzeitbereich",
        "blocks_subtitle": "In BHF 10 zaehlt, Versorgung langfristig zu denken und trotzdem Risiken frueh zu sehen.",
        "blocks": [
            {
                "title": "Biografieorientierung",
                "tone": "sky",
                "body": L(
                    """
                    Vergangenheit erklaert oft aktuelle Reaktionen
                    Pflege wird besser, wenn Gewohnheiten und Werte bekannt sind
                    Ziel: personzentriert statt schematisch handeln
                    """
                ),
            },
            {
                "title": "Pflegediagnostik",
                "tone": "mint",
                "body": L(
                    """
                    aus Assessment-Daten werden pflegerische Probleme und Ziele abgeleitet
                    Diagnosen muessen beobachtbar und pflegerisch bearbeitbar sein
                    """
                ),
            },
            {
                "title": "Gewalt in der Pflege",
                "tone": "sand",
                "body": L(
                    """
                    Formen: physisch, psychisch, strukturell, verbal
                    Warnzeichen: Demuetigung, Zeitdruck, unerklaerte Zwangssituationen
                    Reaktion: ansprechen, sichern, dokumentieren, reflektieren
                    """
                ),
            },
            {
                "title": "Organisation und Finanzierung",
                "tone": "rose",
                "body": L(
                    """
                    Rahmenbedingungen beeinflussen Pflegequalitaet
                    trotzdem bleibt die personzentrierte Versorgung Kernauftrag
                    Ressourcen muessen geplant und begruendet werden
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Biografieorientierung erklaeren
            Assessment und Pflegediagnostik verbinden
            Gewalt in der Langzeitpflege erkennen
            Hygiene im Routinebereich begruenden
            """
        ),
        "check_right": L(
            """
            Langzeitpflege ist Beziehung + Struktur + Beobachtung
            Gewalt beginnt oft schon vor dem offensichtlichen Vorfall
            Gute Pflege braucht auch gute Organisation
            """
        ),
    },
    {
        "no": "11",
        "title": "BHF 11 | Diabetes Typ 1, ADHS und Kinderrechte",
        "subtitle": "Paediatrische Krankheitslehre und familienzentrierte Versorgung",
        "filename": "BHF_11_Diabetes_ADHS_Kinderrechte.pptx",
        "accent": rgb("2B8C90"),
        "accent_dark": rgb("20676A"),
        "topics": L(
            """
            Diabetes Typ 1
            ADHS
            EACH-Charta
            kindgerechte Kommunikation und Beteiligung
            familienzentrierte Versorgung
            """
        ),
        "focus_cards": [
            {
                "title": "Kinderrechte im KH",
                "tone": "sky",
                "body": L(
                    """
                    Kinder brauchen Information, Schutz und Beteiligung
                    Eltern sind Partner, nicht Besucher zweiter Ordnung
                    """
                ),
            },
            {
                "title": "Diabetes im Alltag",
                "tone": "mint",
                "body": L(
                    """
                    Blutzucker, Essen, Bewegung und Insulin muessen zusammen gedacht werden
                    Eltern und Kind brauchen Sicherheit im Handling
                    """
                ),
            },
            {
                "title": "ADHS verstehen",
                "tone": "sand",
                "body": L(
                    """
                    Verhalten ist nicht bloss Ungehorsam
                    Struktur, klare Reize und kurze Anweisungen helfen
                    """
                ),
            },
            {
                "title": "Kommunikation",
                "tone": "rose",
                "body": L(
                    """
                    altersgerecht, konkret und beruhigend sprechen
                    Angst und Scham aktiv abbauen
                    """
                ),
            },
        ],
        "blocks_title": "Krankheitslehre fuer die Kinderpflege",
        "blocks_subtitle": "Hier ist das DURST-Schema besonders hilfreich, weil Elternanleitung und Alltagstransfer sofort mitgedacht werden muessen.",
        "blocks": [
            {
                "title": "Diabetes Typ 1",
                "tone": "sky",
                "body": L(
                    """
                    D: Autoimmunerkrankung mit absolutem Insulinmangel
                    Typen: klar von Typ 2 abgrenzen
                    U + R: Autoimmunneigung, familiäre Belastung
                    S: Polydipsie, Polyurie, Gewichtsverlust, Muedigkeit
                    T/P: Insulin, BZ-Kontrolle, Hypo-/Hyperglykaemie erkennen, Eltern schulen
                    """
                ),
            },
            {
                "title": "ADHS",
                "tone": "mint",
                "body": L(
                    """
                    D: Aufmerksamkeitsstoerung mit Impulsivitaet und/oder Hyperaktivitaet
                    Formen: unaufmerksam, hyperaktiv-impulsiv, kombiniert
                    S: Ablenkbarkeit, Unruhe, Impulsdurchbrueche
                    T/P: Struktur, kurze Auftraege, positives Feedback, Elternanleitung
                    """
                ),
            },
            {
                "title": "EACH-Charta",
                "tone": "sand",
                "body": L(
                    """
                    Kindgerechte Versorgung, Elternnaehe und Beteiligung
                    Angst, Trennung und Informationsdefizite aktiv vermeiden
                    Rechte muessen im Alltag sichtbar sein
                    """
                ),
            },
            {
                "title": "Familienzentrierung",
                "tone": "rose",
                "body": L(
                    """
                    Eltern sind Mitversorgende und wichtige Beobachter
                    Schulung und Entlastung gehoeren zur professionellen Pflege
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Diabetes Typ 1 im DURST-Schema erklaeren
            ADHS kurz und nicht wertend beschreiben
            EACH-Charta in einfachen Worten erklaeren
            Elternarbeit begruenden
            """
        ),
        "check_right": L(
            """
            In Kinderpflege nie nur die Krankheit sehen
            Kind + Eltern + Alltag gehoeren zusammen
            Sicherheit entsteht durch klare, ruhige Anleitung
            """
        ),
    },
    {
        "no": "12",
        "title": "BHF 12 | Fall Bahde, Wohnformen und freiheitsentziehende Massnahmen",
        "subtitle": "Verschlechterung erkennen, Betreuung organisieren und ethisch sauber handeln",
        "filename": "BHF_12_Fall_Bahde_Wohnformen_FEM.pptx",
        "accent": rgb("8C5AA3"),
        "accent_dark": rgb("69447A"),
        "topics": L(
            """
            Frau Bahde baut ab
            Frau Bahde braucht mehr Betreuung
            Frau Bahde zieht in eine Wohnform
            Frau Bahde wird fixiert
            Checkliste und Fallstruktur
            """
        ),
        "focus_cards": [
            {
                "title": "Verschlechterung erkennen",
                "tone": "sky",
                "body": L(
                    """
                    Funktionsverlust, Unruhe, Desorientierung und Sicherheitsrisiken systematisch beobachten
                    Fallverlauf ernst nehmen, nicht nur Einzelereignis
                    """
                ),
            },
            {
                "title": "Versorgung anpassen",
                "tone": "mint",
                "body": L(
                    """
                    Wenn der Bedarf steigt, muessen Setting und Pflegeplanung mitwachsen
                    Wohnform, Hilfsmittel und Teamentscheidungen prüfen
                    """
                ),
            },
            {
                "title": "FEM kritisch denken",
                "tone": "sand",
                "body": L(
                    """
                    freiheitsentziehende Massnahmen sind immer schwerwiegende Eingriffe
                    Alternativen muessen zuerst geprueft werden
                    """
                ),
            },
            {
                "title": "Ethik und Kommunikation",
                "tone": "rose",
                "body": L(
                    """
                    Würde, Sicherheit und Selbstbestimmung abwaegen
                    Angehoerige und Team transparent einbinden
                    """
                ),
            },
        ],
        "blocks_title": "Falllogik fuer die Pruefung",
        "blocks_subtitle": "Auch wenn nicht jede Diagnose benannt ist, verlangt BHF 12 strukturiertes Fallmanagement und ethische Begruendung.",
        "blocks": [
            {
                "title": "Fallanalyse",
                "tone": "sky",
                "body": L(
                    """
                    Was hat sich veraendert?
                    Welche Risiken sind akut?
                    Welche Ressourcen bestehen noch?
                    """
                ),
            },
            {
                "title": "Wohnformen und Uebergaenge",
                "tone": "mint",
                "body": L(
                    """
                    passende Wohnform orientiert sich an Pflegebedarf und Sicherheit
                    Uebergaenge muessen geplant und begleitet werden
                    """
                ),
            },
            {
                "title": "Freiheitsentziehende Massnahmen",
                "tone": "sand",
                "body": L(
                    """
                    nur ultima ratio und rechtlich/ethisch streng gebunden
                    Alternativen: Beobachtung, Umgebung anpassen, Deeskalation, Begleitung
                    Doku und Teamentscheidung sind Pflicht
                    """
                ),
            },
            {
                "title": "Pflegeplanung",
                "tone": "rose",
                "body": L(
                    """
                    Sicherheit, Orientierung, Mobilitaet und Beziehung priorisieren
                    Ziele realistisch und alltagsnah formulieren
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Fall Bahde strukturiert darstellen
            Wohnformwahl begruenden
            FEM als ethisch-rechtliches Thema erklaeren
            Alternativen zu Fixierung nennen
            """
        ),
        "check_right": L(
            """
            Bei Unsicherheit zuerst Risiko und Ziel benennen
            Sicherheit ja, aber nicht automatisch Zwang
            Kommunikation mit Team und Angehoerigen ist zentral
            """
        ),
    },
    {
        "no": "13",
        "title": "BHF 13 | Chronische Erkrankungen und Autonomie",
        "subtitle": "Parkinson, MS, Rheuma, Trajektmodell, Expertenstandard und integrierte Versorgung",
        "filename": "BHF_13_Chronische_Erkrankungen_Autonomie.pptx",
        "accent": rgb("607C35"),
        "accent_dark": rgb("485D28"),
        "topics": L(
            """
            Einstieg chronisch krank
            Parkinson
            MS und Rheuma
            Trajektmodell
            Expertenstandard
            Autonomie unterstuetzen
            Integrierte Versorgung
            Belastung im Umgang mit chronischer Krankheit
            """
        ),
        "focus_cards": [
            {
                "title": "Chronizitaet verstehen",
                "tone": "sky",
                "body": L(
                    """
                    Verlauf, Krisen und Anpassung ueber lange Zeit denken
                    Pflege begleitet Alltag, nicht nur Symptome
                    """
                ),
            },
            {
                "title": "Autonomie und Beratung",
                "tone": "mint",
                "body": L(
                    """
                    Betroffene sollen moeglichst selbst handeln koennen
                    Anleitung, Motivation und realistisches Planen sind zentral
                    """
                ),
            },
            {
                "title": "Expertenstandard",
                "tone": "sand",
                "body": L(
                    """
                    Standards begruenden professionelles Handeln
                    sie machen Pflege verlässlich und pruefbar
                    """
                ),
            },
            {
                "title": "Belastungssystem",
                "tone": "rose",
                "body": L(
                    """
                    chronische Krankheit belastet Betroffene und Umfeld
                    Familie, Arbeit und Lebensqualitaet immer mitdenken
                    """
                ),
            },
        ],
        "blocks_title": "DURST-Kurzformat fuer chronische Krankheitslehre",
        "blocks_subtitle": "Bei chronischen Erkrankungen in der Pruefung immer Verlauf, Alltag und Adhaerenz mitdenken.",
        "blocks": [
            {
                "title": "Parkinson",
                "tone": "sky",
                "body": L(
                    """
                    D: neurodegeneratives Syndrom mit Dopaminmangel
                    Typen: idiopathisch, atypisch, sekundaer
                    S: Bradykinese, Tremor, Rigor, Instabilitaet
                    T/P: Medikamente pünktlich, Sturzprophylaxe, Zeit geben, Obstipation beachten
                    """
                ),
            },
            {
                "title": "Multiple Sklerose",
                "tone": "mint",
                "body": L(
                    """
                    D: chronisch entzuendliche Erkrankung des ZNS
                    Verlaeufe: schubfoermig, progredient
                    S: Sensibilitaetsstoerungen, Fatigue, Paresen, Sehstörungen
                    T/P: Energie einteilen, Symptome beobachten, Mobilitaet und Beratung
                    """
                ),
            },
            {
                "title": "Rheuma",
                "tone": "sand",
                "body": L(
                    """
                    D: entzuendliche oder degenerative Gelenkserkrankung
                    S: Schmerz, Morgensteifigkeit, Bewegungseinschraenkung
                    T/P: Medikamente, Bewegung im passenden Mass, Gelenkschutz, Schmerzmanagement
                    """
                ),
            },
            {
                "title": "Trajektmodell",
                "tone": "rose",
                "body": L(
                    """
                    beschreibt Verlauf, Krisen und Arbeit einer chronischen Krankheit
                    hilft, Belastung und Pflegebedarf ueber Zeit zu verstehen
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Chronische Krankheit als Verlauf erklaeren
            Parkinson, MS und Rheuma knapp darstellen
            Trajektmodell einfach beschreiben
            Autonomiefoerderung begruenden
            """
        ),
        "check_right": L(
            """
            Nicht nur Krankheit, sondern Alltag und Ressourcen beschreiben
            Bei Chronik ist Beratung fast immer pruefungsrelevant
            Medikamentenmanagement und Selbststaendigkeit zusammen denken
            """
        ),
    },
    {
        "no": "14",
        "title": "BHF 14 | Diagnostik, Kommunikation und Tod und Sterben",
        "subtitle": "Therapie, Pflege, Kommunikation, eigenes Erleben und palliative Grundhaltung",
        "filename": "BHF_14_Diagnostik_Kommunikation_Tod_Sterben.pptx",
        "accent": rgb("7C5F97"),
        "accent_dark": rgb("5E486F"),
        "topics": L(
            """
            Diagnostik und Therapie
            Pflege
            Kommunikation
            Pflege und eigenes Erleben
            Tod und Sterben
            """
        ),
        "focus_cards": [
            {
                "title": "Diagnostik und Therapie",
                "tone": "sky",
                "body": L(
                    """
                    Patient muss verstehen, was untersucht und behandelt wird
                    Pflege beobachtet Wirkung, Belastung und Nebenwirkungen
                    """
                ),
            },
            {
                "title": "Kommunikation in Belastung",
                "tone": "mint",
                "body": L(
                    """
                    schlechte Nachrichten brauchen Klarheit, Ruhe und Echtheit
                    zuhören ist oft wichtiger als sofort loesen
                    """
                ),
            },
            {
                "title": "Eigenes Erleben",
                "tone": "sand",
                "body": L(
                    """
                    Mitgefuehl ja, Ueberidentifikation nein
                    Reflexion schuetzt Professionalitaet und Team
                    """
                ),
            },
            {
                "title": "Sterben begleiten",
                "tone": "rose",
                "body": L(
                    """
                    Wuerde, Symptomlinderung und Beziehungsarbeit stehen im Zentrum
                    Anhoerige mitdenken und Sicherheit geben
                    """
                ),
            },
        ],
        "blocks_title": "Pruefungsfokus im palliativen Denken",
        "blocks_subtitle": "BHF 14 verlangt vor allem Haltung, Kommunikation und saubere Beobachtung bei schwerer Krankheit.",
        "blocks": [
            {
                "title": "Kommunikation",
                "tone": "sky",
                "body": L(
                    """
                    ehrlich, klar und dosiert sprechen
                    Pausen zulassen und Gefuehle benennen
                    keine falschen Versprechen machen
                    """
                ),
            },
            {
                "title": "Pflege bei Belastung",
                "tone": "mint",
                "body": L(
                    """
                    Symptome beobachten: Schmerz, Atemnot, Uebelkeit, Angst
                    kleine Massnahmen koennen grosse Erleichterung bringen
                    """
                ),
            },
            {
                "title": "Eigenes Erleben",
                "tone": "sand",
                "body": L(
                    """
                    belastende Situationen im Team reflektieren
                    Grenzen erkennen und Unterstuetzung annehmen
                    """
                ),
            },
            {
                "title": "Tod und Sterben",
                "tone": "rose",
                "body": L(
                    """
                    Wuerde, Symptomkontrolle und Beziehung stehen im Vordergrund
                    Palliative Haltung heisst nicht aufgeben, sondern Ziele veraendern
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Kommunikation in schweren Situationen erklaeren
            palliative Pflegegrundsaetze nennen
            eigenes Erleben professionell einordnen
            Tod und Sterben respektvoll beschreiben
            """
        ),
        "check_right": L(
            """
            In diesem BHF ist Haltung genauso wichtig wie Wissen
            Symptomlinderung und Wuerde immer nennen
            Ehrliche, ruhige Kommunikation gibt Sicherheit
            """
        ),
    },
    {
        "no": "15",
        "title": "BHF 15 | Berufskrankheiten, Burnout und Suchtpraevention",
        "subtitle": "Mobbing, Cybermobbing, betriebliche Gesundheit und psychosoziale Risiken",
        "filename": "BHF_15_Berufskrankheiten_Burnout_Sucht.pptx",
        "accent": rgb("B56D2D"),
        "accent_dark": rgb("884F22"),
        "topics": L(
            """
            Berufskrankheiten
            Burnout
            Mobbing
            Cybermobbing und Hate Speech
            Sucht und Suchtpraevention
            Betriebliche Gesundheitsfoerderung
            Bewertungskriterien und Arbeitsorganisation
            """
        ),
        "focus_cards": [
            {
                "title": "Arbeitsgesundheit",
                "tone": "sky",
                "body": L(
                    """
                    Gesundheit im Beruf ist Voraussetzung fuer gute Pflege
                    Belastungen frueh erkennen und nicht normalisieren
                    """
                ),
            },
            {
                "title": "Psychosoziale Risiken",
                "tone": "mint",
                "body": L(
                    """
                    Mobbing, Ueberforderung und Sucht gefaehrden Menschen und Teams
                    Praevention braucht offene Kultur und klare Regeln
                    """
                ),
            },
            {
                "title": "Burnout und Grenzen",
                "tone": "sand",
                "body": L(
                    """
                    Erschoepfung ist ein Warnsignal, kein Leistungsbeweis
                    Pausen, Teamunterstuetzung und Selbstschutz sind professionell
                    """
                ),
            },
            {
                "title": "Betrieblich denken",
                "tone": "rose",
                "body": L(
                    """
                    Gesundheitsfoerderung ist Aufgabe von Person und Organisation
                    gute Strukturen senken Risiken
                    """
                ),
            },
        ],
        "blocks_title": "Krankheitslehre und Praevention",
        "blocks_subtitle": "In der Pruefung reicht es nicht, Belastungen zu nennen. Du musst auch Praevention und Hilfen ableiten koennen.",
        "blocks": [
            {
                "title": "Burnout",
                "tone": "sky",
                "body": L(
                    """
                    D: arbeitsbezogene Erschoepfung mit Distanzierung und Leistungsabfall
                    Risiken: Dauerstress, Personalmangel, fehlende Erholung
                    S: Muedigkeit, Zynismus, Konzentrationsverlust
                    T/P: Belastungen benennen, Hilfe holen, Strukturen und Selbstschutz verbessern
                    """
                ),
            },
            {
                "title": "Suchtstoerung",
                "tone": "mint",
                "body": L(
                    """
                    D: schädlicher oder abhaengiger Konsum mit Kontrollverlust
                    Risiken: Stress, psychische Belastung, Verfuegbarkeit
                    S: Craving, Vernachlaessigung, Toleranzentwicklung
                    T/P: fruehes Ansprechen, Beratungswege, Praevention, Schutz des Teams
                    """
                ),
            },
            {
                "title": "Mobbing / Cybermobbing",
                "tone": "sand",
                "body": L(
                    """
                    wiederholte Abwertung oder Ausgrenzung, auch digital
                    Folgen: Angst, Rueckzug, Krankheit, Teamverlust
                    Reaktion: ansprechen, dokumentieren, Strukturen und Schutz nutzen
                    """
                ),
            },
            {
                "title": "Betriebliche Praevention",
                "tone": "rose",
                "body": L(
                    """
                    Arbeitsorganisation, Fuehrung, Pausen und Kultur sind Gesundheitsfaktoren
                    Praevention ist wirksamer als spaete Krisenintervention
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Burnout und Sucht kurz einordnen
            Mobbing von Konflikt unterscheiden
            Betriebliche Gesundheitsfoerderung erklaeren
            Praeventionsmassnahmen nennen
            """
        ),
        "check_right": L(
            """
            Psychosoziale Risiken sind echte Pflege- und Patientensicherheitsrisiken
            Praevention braucht Person + Team + Organisation
            Belastung benennen ist ein Zeichen von Professionalitaet
            """
        ),
    },
    {
        "no": "16",
        "title": "BHF 16 | Psychiatrie, Depression und Anorexie",
        "subtitle": "Geschichte der Psychiatrie, Interventionen, Depressionen und Anorexia nervosa",
        "filename": "BHF_16_Psychiatrie_Depression_Anorexie.pptx",
        "accent": rgb("7A5A9B"),
        "accent_dark": rgb("5C4476"),
        "topics": L(
            """
            Das weisse Rauschen
            Geschichte der Psychiatrie
            Interventionen in der Psychiatrie
            Anorexia nervosa
            Depressionen
            """
        ),
        "focus_cards": [
            {
                "title": "Psychiatrische Grundhaltung",
                "tone": "sky",
                "body": L(
                    """
                    Sicherheit, Beziehung und Struktur bilden das Fundament
                    Menschen nicht auf ihre Diagnose reduzieren
                    """
                ),
            },
            {
                "title": "Interventionen",
                "tone": "mint",
                "body": L(
                    """
                    ruhig, klar, wertschätzend und nicht konfrontativ
                    Beobachtung und Beziehungsarbeit sind aktive Pflege
                    """
                ),
            },
            {
                "title": "Depression",
                "tone": "sand",
                "body": L(
                    """
                    ernst nehmen, Aktivierung dosieren, Suizidalitaet mitdenken
                    kleine Schritte und verlässliche Beziehung helfen
                    """
                ),
            },
            {
                "title": "Anorexie",
                "tone": "rose",
                "body": L(
                    """
                    Koerperbild, Kontrolle und somatische Gefahr gleichzeitig sehen
                    konsequent und nicht moralisierend pflegen
                    """
                ),
            },
        ],
        "blocks_title": "DURST fuer psychiatrische Krankheitslehre",
        "blocks_subtitle": "In diesem BHF sind Beziehungsgestaltung und somatische Risiken gleich wichtig.",
        "blocks": [
            {
                "title": "Depression",
                "tone": "sky",
                "body": L(
                    """
                    D: affektive Stoerung mit gedrueckter Stimmung und Antriebsmangel
                    Formen: leicht, mittel, schwer, rezidivierend
                    S: Interessenverlust, Schlafstoerung, Schuld, Suizidgedanken
                    T/P: Therapie begleiten, Tagesstruktur, Suizidalitaet ansprechen, Sicherheit
                    """
                ),
            },
            {
                "title": "Anorexia nervosa",
                "tone": "mint",
                "body": L(
                    """
                    D: Essstoerung mit absichtlich herbeigefuehrtem Untergewicht
                    Formen: restriktiv oder binge/purging
                    S: Untergewicht, Bradykardie, Frieren, verzerrtes Koerperbild
                    T/P: medizinische Ueberwachung, Essstruktur, Beziehung, klare Grenzen
                    """
                ),
            },
            {
                "title": "Psychiatrische Intervention",
                "tone": "sand",
                "body": L(
                    """
                    Praesenz, Orientierung, Reizarme Umgebung und klare Sprache
                    Eskalation vermeiden und Sicherheit im Team halten
                    """
                ),
            },
            {
                "title": "Geschichte der Psychiatrie",
                "tone": "rose",
                "body": L(
                    """
                    zeigt den Wandel von Verwahrung zu patientenzentrierter Behandlung
                    wichtig fuer Haltung, Menschenbild und Kritikfaehigkeit
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Depression und Anorexie im DURST-Schema erklaeren
            Suizidalitaet professionell einordnen
            psychiatrische Grundhaltung nennen
            """
        ),
        "check_right": L(
            """
            Psychiatrie braucht Beziehung UND Struktur
            Sicherheit geht vor Diskussion
            Somatische Risiken nie unterschätzen
            """
        ),
    },
    {
        "no": "17",
        "title": "BHF 17 | Herausfordernde Pflegesituationen und Infektionsschutz",
        "subtitle": "Infektionskrankheiten, Praevention und Umgang mit schwierigen Situationen",
        "filename": "BHF_17_Infektionsschutz_Herausfordernde_Situationen.pptx",
        "accent": rgb("2B8F76"),
        "accent_dark": rgb("206A58"),
        "topics": L(
            """
            Mit herausfordernden Pflegesituationen umgehen
            Infektionsschutz
            Infektionskrankheiten
            Gesundheit erhalten und foerdern
            """
        ),
        "focus_cards": [
            {
                "title": "Herausfordernde Situationen",
                "tone": "sky",
                "body": L(
                    """
                    ruhig bleiben, Sicherheit sichern, Verhalten nicht persoenlich nehmen
                    Verhalten verstehen hilft bei der Deeskalation
                    """
                ),
            },
            {
                "title": "Infektionsschutz",
                "tone": "mint",
                "body": L(
                    """
                    Standardmassnahmen sind Basis, Zusatzmassnahmen folgen dem Uebertragungsweg
                    Hygiene ist Teamaufgabe
                    """
                ),
            },
            {
                "title": "Gesundheit erhalten",
                "tone": "sand",
                "body": L(
                    """
                    Praevention, Impfungen, Aufklaerung und Ressourcen foerdern
                    Gesundheit ist mehr als Infektfreiheit
                    """
                ),
            },
            {
                "title": "Pflegehaltung",
                "tone": "rose",
                "body": L(
                    """
                    Schutzmassnahmen klaeren, nicht nur anordnen
                    Anhoerige und Patienten verstaendlich anleiten
                    """
                ),
            },
        ],
        "blocks_title": "Praxisblöcke fuer Hygiene und Deeskalation",
        "blocks_subtitle": "In BHF 17 werden besonders Standardmassnahmen, Uebertragungswege und professioneller Umgang geprueft.",
        "blocks": [
            {
                "title": "Infektionskette",
                "tone": "sky",
                "body": L(
                    """
                    Erreger, Reservoir, Austritt, Uebertragung, Eintritt, empfänglicher Mensch
                    an jeder Stelle kann Pflege unterbrechen
                    """
                ),
            },
            {
                "title": "Schutzmassnahmen",
                "tone": "mint",
                "body": L(
                    """
                    Haendedesinfektion, Handschuhe, Kittel, Masken passend waehlen
                    Massnahme immer am Uebertragungsweg begruenden
                    """
                ),
            },
            {
                "title": "Herausforderndes Verhalten",
                "tone": "sand",
                "body": L(
                    """
                    Ausloeser suchen: Schmerz, Angst, Verwirrung, Ueberforderung
                    mit klarer Sprache, Abstand und Struktur arbeiten
                    """
                ),
            },
            {
                "title": "Gesundheitsfoerderung",
                "tone": "rose",
                "body": L(
                    """
                    Praevention ist Beratung, Motivation und Lebensweltarbeit
                    Schutzfaktoren staerken statt nur Risiken benennen
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Infektionskette erklaeren
            Schutzmassnahmen begruenden
            Herausfordernde Pflegesituationen einordnen
            Gesundheitsfoerderung beschreiben
            """
        ),
        "check_right": L(
            """
            Immer den Uebertragungsweg nennen
            Verhalten hat oft einen Grund
            Deeskalation beginnt mit eigener Ruhe
            """
        ),
    },
    {
        "no": "18",
        "title": "BHF 18 | Unfall, Polytrauma, ITS und Frakturen",
        "subtitle": "Erstversorgung, Intensivstation, SHT, Frakturen und informierte Entscheidung",
        "filename": "BHF_18_Unfall_Polytrauma_ITS_Frakturen.pptx",
        "accent": rgb("B45D44"),
        "accent_dark": rgb("874735"),
        "topics": L(
            """
            Unfallgeschehen und Erstversorgung
            Auf einer ITS arbeiten
            Dyen Pham nach Polytrauma pflegen
            Johanna Soltek nach SHT und Fraktur pflegen
            Informierte Entscheidung treffen
            Erstversorgung am Unfallort
            Frakturen allgemein
            Simon Roeder
            """
        ),
        "focus_cards": [
            {
                "title": "Erstversorgung",
                "tone": "sky",
                "body": L(
                    """
                    Lebensgefahr zuerst, dann weitere Verletzungen
                    klares Teamhandeln und Prioritaeten sind entscheidend
                    """
                ),
            },
            {
                "title": "ITS-Denken",
                "tone": "mint",
                "body": L(
                    """
                    engmaschige Beobachtung, Monitoring und kleine Trends
                    Technik ersetzt nicht klinisches Denken
                    """
                ),
            },
            {
                "title": "Trauma und Fraktur",
                "tone": "sand",
                "body": L(
                    """
                    Schmerz, Blutung, Schwellung, Durchblutung und Neurologie beachten
                    Immobilisation und Sicherheit begruenden
                    """
                ),
            },
            {
                "title": "Entscheidungen",
                "tone": "rose",
                "body": L(
                    """
                    Information, Einwilligung und Beteiligung bleiben wichtig
                    auch auf der ITS braucht es Kommunikation
                    """
                ),
            },
        ],
        "blocks_title": "Trauma-Krankheitslehre im Kurzformat",
        "blocks_subtitle": "Bei BHF 18 zaehlt priorisierte Notfalllogik plus pflegerische Detailbeobachtung.",
        "blocks": [
            {
                "title": "Polytrauma",
                "tone": "sky",
                "body": L(
                    """
                    D: mehrere schwere Verletzungen mit vitaler Bedrohung
                    Risiken: Blutverlust, Atemweg, Schock, Hypothermie
                    Pflege: ABCDE, Monitoring, Waerme, Lagerung, Teamkommunikation
                    """
                ),
            },
            {
                "title": "Schaedel-Hirn-Trauma",
                "tone": "mint",
                "body": L(
                    """
                    D: Verletzung von Schaedelknochen und/oder Gehirn
                    Grade: leicht, mittel, schwer
                    S: Bewusstsein, Pupillen, Erbrechen, Kopfschmerz, Ausfaelle
                    T/P: Ueberwachung, Reizkontrolle, Hirndruckzeichen beachten
                    """
                ),
            },
            {
                "title": "Fraktur",
                "tone": "sand",
                "body": L(
                    """
                    D: Unterbrechung der Knochenkontinuitaet
                    Formen: offen/geschlossen, disloziert/nicht disloziert
                    Pflege: Immobilisation, Schmerzmanagement, DMS-Kontrolle
                    """
                ),
            },
            {
                "title": "ITS und Einwilligung",
                "tone": "rose",
                "body": L(
                    """
                    Auch in Akutsituationen sind Information und Beteiligung relevant
                    Patientenrechte und Teamkommunikation bleiben wichtig
                    """
                ),
            },
        ],
        "check_left": L(
            """
            ABCDE bei Trauma erklaeren
            Polytrauma, SHT und Fraktur knapp darstellen
            ITS-Pflegebeobachtungen nennen
            DMS-Kontrolle erklaeren
            """
        ),
        "check_right": L(
            """
            Erst lebensbedrohlich, dann wichtig
            DMS immer vor und nach Massnahmen pruefen
            Technik unterstuetzt, ersetzt aber nicht Beobachtung
            """
        ),
    },
    {
        "no": "19",
        "title": "BHF 19 | Schlaganfall-Reha und Aphasie",
        "subtitle": "Rehabilitative Pflege, Pflegeprozess in der Reha und Kommunikation mit Aphasikern",
        "filename": "BHF_19_Schlaganfall_Reha_Aphasie.pptx",
        "accent": rgb("5A83B0"),
        "accent_dark": rgb("446387"),
        "topics": L(
            """
            Nach Schlaganfall rehabilitativ pflegen
            Nach Reha den Pflegeprozess steuern
            Mit Aphasikern richtig kommunizieren
            """
        ),
        "focus_cards": [
            {
                "title": "Rehabilitative Haltung",
                "tone": "sky",
                "body": L(
                    """
                    Funktion und Teilhabe sind das Ziel
                    kleine alltagsnahe Schritte bringen Fortschritt
                    """
                ),
            },
            {
                "title": "Pflegeprozess in Reha",
                "tone": "mint",
                "body": L(
                    """
                    Assessment, Ziele und Evaluation muessen sehr konkret sein
                    Fortschritte sichtbar machen und motivierend rueckmelden
                    """
                ),
            },
            {
                "title": "Aphasie",
                "tone": "sand",
                "body": L(
                    """
                    Sprachstoerung ist nicht gleich Denkstoerung
                    Zeit, Blickkontakt und Hilfen schaffen Kommunikation
                    """
                ),
            },
            {
                "title": "Interprofessionell arbeiten",
                "tone": "rose",
                "body": L(
                    """
                    Pflege, Physio, Ergo, Logopaedie und Arztteam muessen verzahnt sein
                    Ziele sollten gemeinsam getragen werden
                    """
                ),
            },
        ],
        "blocks_title": "Reha-Praxis fuer die Pruefung",
        "blocks_subtitle": "Hier ist weniger die Akutdiagnose als die funktionelle Pflege und Kommunikation entscheidend.",
        "blocks": [
            {
                "title": "Reha nach Schlaganfall",
                "tone": "sky",
                "body": L(
                    """
                    Mobilitaet, Selbstversorgung, Schlucken und Sicherheit priorisieren
                    Wiederholung und Anleitung sind zentrale Pflegeaufgaben
                    """
                ),
            },
            {
                "title": "Aphasie kommunizieren",
                "tone": "mint",
                "body": L(
                    """
                    kurze Saetze, eine Frage nach der anderen
                    Zeit lassen, Gestik/Bildkarten nutzen, nicht fuer den Patienten sprechen
                    """
                ),
            },
            {
                "title": "Pflegeprozess steuern",
                "tone": "sand",
                "body": L(
                    """
                    Ziele muessen messbar und alltagsnah sein
                    Evaluation zeigt, ob die Reha-Strategie angepasst werden muss
                    """
                ),
            },
            {
                "title": "Familie und Motivation",
                "tone": "rose",
                "body": L(
                    """
                    Angehoerige brauchen Anleitung und realistische Hoffnung
                    Motivation entsteht durch erkennbare kleine Erfolge
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Reha-Pflege nach Schlaganfall erklaeren
            Kommunikation mit Aphasikern beschreiben
            Ziele in der Reha formulieren
            Teamarbeit begruenden
            """
        ),
        "check_right": L(
            """
            Reha ist alltagsbezogen, nicht nur trainingsbezogen
            Aphasie braucht Geduld und Struktur
            Pflegeprozess bleibt auch in der Reha das Grundgeruest
            """
        ),
    },
    {
        "no": "20",
        "title": "BHF 20 | Urogenitalsystem, Dialyse und Kontinenz",
        "subtitle": "Nieren, Dialyse, Kontinenzprobleme, Qualitaetssicherung und Beratung",
        "filename": "BHF_20_Urogenital_Dialyse_Kontinenz.pptx",
        "accent": rgb("A36A2B"),
        "accent_dark": rgb("7A4E20"),
        "topics": L(
            """
            Anatomie, Physiologie und Erkrankungen des Urogenitalsystems
            Dialysepatient pflegen und beraten
            Menschen mit Kontinenzproblemen pflegen
            Qualitaetssichernd handeln
            Verschiedene Perspektiven entwickeln
            """
        ),
        "focus_cards": [
            {
                "title": "Urogenital denken",
                "tone": "sky",
                "body": L(
                    """
                    Anatomie und Funktion erklaeren viele Symptome
                    Ausscheidung ist ein wichtiges Assessmentfeld
                    """
                ),
            },
            {
                "title": "Dialyse",
                "tone": "mint",
                "body": L(
                    """
                    Betroffene brauchen Beobachtung, Beratung und Sicherheit
                    Gewicht, Bilanz und Allgemeinzustand sind Schluessel
                    """
                ),
            },
            {
                "title": "Kontinenz",
                "tone": "sand",
                "body": L(
                    """
                    Scham, Hautschutz und Selbststaendigkeit mitdenken
                    Inkontinenz ist ein Pflege- und Beratungsthema
                    """
                ),
            },
            {
                "title": "Qualitaet",
                "tone": "rose",
                "body": L(
                    """
                    Standards und Evaluation machen Pflege verlässlich
                    Perspektivenwechsel verbessert Entscheidungen
                    """
                ),
            },
        ],
        "blocks_title": "DURST fuer Niere und Kontinenz",
        "blocks_subtitle": "Bei BHF 20 solltest du Krankheitslehre immer mit Bilanz, Haut und Beratung verbinden.",
        "blocks": [
            {
                "title": "Chronische Niereninsuffizienz",
                "tone": "sky",
                "body": L(
                    """
                    D: fortschreitender Funktionsverlust der Niere
                    Stadien: G1 bis G5 nach GFR
                    S: Oedeme, Muedigkeit, Juckreiz, Hypertonie, Uraemiesymptome
                    T/P: Blutdruck, Bilanz, Ernaehrung, Zugang, Dialysebeobachtung
                    """
                ),
            },
            {
                "title": "Dialyse",
                "tone": "mint",
                "body": L(
                    """
                    Ziel: Ersatz von Ausscheidungs- und Regulationsfunktion
                    Pflege: Gewicht, Volumenstatus, Kreislauf, Zugang und Beratung
                    Risiken: Hypotonie, Infektion, Erschoepfung
                    """
                ),
            },
            {
                "title": "Harninkontinenz",
                "tone": "sand",
                "body": L(
                    """
                    D: unwillkuerlicher Urinverlust
                    Typen: Belastungs-, Drang-, Ueberlauf-, funktionelle Inkontinenz
                    S: Urinverlust, Harndrang, Hautprobleme, Rueckzug
                    T/P: Toilettentraining, Hautschutz, Beobachtung, passende Hilfsmittel
                    """
                ),
            },
            {
                "title": "Qualitaetssicherung",
                "tone": "rose",
                "body": L(
                    """
                    Beobachtung, Standard, Doku und Evaluation greifen zusammen
                    gute Pflege ist nachpruefbar und begruendbar
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Niereninsuffizienz und Dialyse erklaeren
            Formen der Harninkontinenz nennen
            Bilanzierung und Hautschutz begruenden
            Qualitaetssicherung beschreiben
            """
        ),
        "check_right": L(
            """
            Ausscheidung ist immer ein zentrales Assessmentthema
            Bei Dialyse nie nur Technik, sondern auch Allgemeinzustand beachten
            Inkontinenz immer ursachenbezogen denken
            """
        ),
    },
    {
        "no": "21",
        "title": "BHF 21 | Fruehgeborene, Schreibaby und paediatrische Erkrankungen",
        "subtitle": "Neonatologie, kindliche Darmerkrankungen und familienzentrierte Pflege",
        "filename": "BHF_21_Fruehgeborene_Schreibaby_Paediatrie.pptx",
        "accent": rgb("2F8A8B"),
        "accent_dark": rgb("236768"),
        "topics": L(
            """
            Pflege eines Fruehgeborenen
            Schreibaby
            Pflege von Kindern mit Darmerkrankungen
            Pflege von Kindern mit paediatrischen Erkrankungen
            Friedemann
            """
        ),
        "focus_cards": [
            {
                "title": "Fruehgeborene",
                "tone": "sky",
                "body": L(
                    """
                    Temperatur, Atmung, Ernaehrung und Bindung sind zentral
                    Eltern brauchen besonders viel Anleitung und Sicherheit
                    """
                ),
            },
            {
                "title": "Schreibaby",
                "tone": "mint",
                "body": L(
                    """
                    Erschoepfung der Eltern mitsehen
                    Beruhigung, Struktur und Entlastung sind wesentliche Pflegeziele
                    """
                ),
            },
            {
                "title": "Darmerkrankungen",
                "tone": "sand",
                "body": L(
                    """
                    Beobachte Bauch, Nahrung, Ausscheidung, Schmerz und Allgemeinzustand
                    Dehydratation und Verschlechterung frueh erkennen
                    """
                ),
            },
            {
                "title": "Familienzentrierung",
                "tone": "rose",
                "body": L(
                    """
                    Kindliche Erkrankung betrifft immer auch Eltern und Tagesablauf
                    Pflege braucht Erklaerung, Ruhe und Beteiligung
                    """
                ),
            },
        ],
        "blocks_title": "Paediatrischer Pruefungstransfer",
        "blocks_subtitle": "Bei BHF 21 werden haeufig Entwicklungsstand, Elternanleitung und fruehes Erkennen von Verschlechterung abgefragt.",
        "blocks": [
            {
                "title": "Fruehgeborenes",
                "tone": "sky",
                "body": L(
                    """
                    Risiken: Temperaturinstabilitaet, Atemprobleme, Trinkschwaeche
                    Pflege: Waerme, Beobachtung, Schutz vor Ueberreizung, Eltern einbeziehen
                    """
                ),
            },
            {
                "title": "Schreibaby",
                "tone": "mint",
                "body": L(
                    """
                    exzessives Schreien belastet Kind und Eltern stark
                    Reize reduzieren, Struktur geben, Eltern entlasten und beraten
                    """
                ),
            },
            {
                "title": "Darmerkrankungen",
                "tone": "sand",
                "body": L(
                    """
                    auf Bauchumfang, Erbrechen, Stuhl, Schmerzen und Trinkverhalten achten
                    Dehydratation und Kreislauf frueh mitdenken
                    """
                ),
            },
            {
                "title": "Paediatrische Pflege",
                "tone": "rose",
                "body": L(
                    """
                    alter, Entwicklung und Familie bestimmen die Pflegegestaltung
                    kindgerechte Sprache und Erklaerung bleiben entscheidend
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Pflege eines Fruehgeborenen beschreiben
            Schreibaby pflegerisch einordnen
            kindliche Darmerkrankungen beobachten
            Elternarbeit begruenden
            """
        ),
        "check_right": L(
            """
            Bei Kindern veraendert sich der Zustand schnell
            Eltern sind Ressource und muessen mitversorgt werden
            Reizschutz und Beobachtung sind Schluesselelemente
            """
        ),
    },
    {
        "no": "22",
        "title": "BHF 22 | Zukunftswerkstatt und Projektdenken",
        "subtitle": "Ideen entwickeln, priorisieren, praesentieren und reflektieren",
        "filename": "BHF_22_Zukunftswerkstatt_Projektdenken.pptx",
        "accent": rgb("4A7FA1"),
        "accent_dark": rgb("355F79"),
        "topics": L(
            """
            Zukunftswerkstatt
            Ideenentwicklung
            Priorisierung
            Teamarbeit
            Praesentation und Reflexion
            """
        ),
        "focus_cards": [
            {
                "title": "Problem sehen",
                "tone": "sky",
                "body": L(
                    """
                    erst das Problem klar beschreiben, dann Loesungen bauen
                    gute Analyse spart spaet Chaos
                    """
                ),
            },
            {
                "title": "Ideen entwickeln",
                "tone": "mint",
                "body": L(
                    """
                    breit denken, noch nicht zu frueh bewerten
                    kreative Ideen brauchen danach Struktur
                    """
                ),
            },
            {
                "title": "Priorisieren",
                "tone": "sand",
                "body": L(
                    """
                    Nutzen, Umsetzbarkeit und Ressourcen vergleichen
                    nicht alles gleichzeitig anfangen
                    """
                ),
            },
            {
                "title": "Praesentieren",
                "tone": "rose",
                "body": L(
                    """
                    Ziel, Problem, Loesung und Nutzen klar zeigen
                    Reflexion gehoert zum professionellen Arbeiten
                    """
                ),
            },
        ],
        "blocks_title": "Projektlogik fuer muendliche Pruefungen",
        "blocks_subtitle": "BHF 22 wirkt weniger medizinisch, ist aber stark fuer Struktur-, Team- und Praesentationskompetenz.",
        "blocks": [
            {
                "title": "Analyse",
                "tone": "sky",
                "body": L(
                    """
                    Problem, Zielgruppe und Ausgangslage klaeren
                    Ohne Analyse keine tragfaehige Loesung
                    """
                ),
            },
            {
                "title": "Ideenphase",
                "tone": "mint",
                "body": L(
                    """
                    kreativ sammeln, Perspektiven wechseln, offen denken
                    erst danach priorisieren
                    """
                ),
            },
            {
                "title": "Umsetzung",
                "tone": "sand",
                "body": L(
                    """
                    Schritte, Rollen, Zeit und Ressourcen planen
                    kleine realistische Ziele setzen
                    """
                ),
            },
            {
                "title": "Reflexion",
                "tone": "rose",
                "body": L(
                    """
                    Was hat funktioniert, was nicht, was lernen wir?
                    Reflexion verbessert kuenftige Projekte und Teamarbeit
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Zukunftswerkstatt in Phasen erklaeren
            Ideen von Bewertung trennen
            Priorisierung begruenden
            Praesentation strukturiert aufbauen
            """
        ),
        "check_right": L(
            """
            Erst Problem, dann Loesung
            Teamarbeit braucht Rollen und Klarheit
            Reflexion ist kein Extra, sondern Teil des Prozesses
            """
        ),
    },
    {
        "no": "23",
        "title": "BHF 23 | Gewalt, Konflikt und Deeskalation",
        "subtitle": "Gewaltsituationen erkennen, Konflikte analysieren und praeventiv handeln",
        "filename": "BHF_23_Gewalt_Konflikt_Deeskalation.pptx",
        "accent": rgb("A04D5B"),
        "accent_dark": rgb("7A3945"),
        "topics": L(
            """
            Gewaltsituationen in der Pflege erkennen
            Konflikte erkennen
            Loesungen ableiten
            Loesungsstrategien bei Konflikten
            Gewalt in der Pflege praeventiv begegnen
            """
        ),
        "focus_cards": [
            {
                "title": "Gewalt erkennen",
                "tone": "sky",
                "body": L(
                    """
                    Gewalt kann physisch, psychisch, strukturell oder verbal sein
                    fruehe Warnzeichen ernst nehmen
                    """
                ),
            },
            {
                "title": "Konflikte analysieren",
                "tone": "mint",
                "body": L(
                    """
                    Was ist der Ausloeser? Welche Interessen prallen aufeinander?
                    Konflikte muessen erst verstanden werden, bevor sie geloest werden koennen
                    """
                ),
            },
            {
                "title": "Deeskalation",
                "tone": "sand",
                "body": L(
                    """
                    klare Sprache, ruhige Haltung, Abstand und Sicherheit
                    nicht provozieren, sondern strukturieren
                    """
                ),
            },
            {
                "title": "Praevention",
                "tone": "rose",
                "body": L(
                    """
                    Teamkultur, Reflexion und Meldesysteme verhindern Wiederholung
                    Gewaltpraevention ist Qualitaetsarbeit
                    """
                ),
            },
        ],
        "blocks_title": "Pruefungstransfer fuer Konfliktfaelle",
        "blocks_subtitle": "BHF 23 ist ein starkes muendliches Thema, weil hier Haltung, Sprache und Struktur sichtbar werden.",
        "blocks": [
            {
                "title": "Gewaltformen",
                "tone": "sky",
                "body": L(
                    """
                    physisch, psychisch, sexualisiert, strukturell
                    auch Zeitdruck und Demuetigung koennen Gewaltcharakter haben
                    """
                ),
            },
            {
                "title": "Konfliktanalyse",
                "tone": "mint",
                "body": L(
                    """
                    Situation, Beteiligte, Ausloeser und Dynamik erfassen
                    Beobachtung vor Bewertung
                    """
                ),
            },
            {
                "title": "Loesungsstrategien",
                "tone": "sand",
                "body": L(
                    """
                    klaeren, begrenzen, deeskalieren, Hilfe holen
                    klare Absprachen und Dokumentation sind wichtig
                    """
                ),
            },
            {
                "title": "Praeventive Kultur",
                "tone": "rose",
                "body": L(
                    """
                    Reflexion, Teamgespraeche und Schutzkonzepte aufbauen
                    Gewaltpraevention beginnt vor dem Vorfall
                    """
                ),
            },
        ],
        "check_left": L(
            """
            Gewaltformen nennen
            Konfliktanalyse Schritt fuer Schritt beschreiben
            Deeskalation begruenden
            Praeventionsmassnahmen aufzaehlen
            """
        ),
        "check_right": L(
            """
            Eigene Ruhe ist Teil der Deeskalation
            Beobachten, sichern, dokumentieren, reflektieren
            Gewalt ist nie nur ein Individualproblem
            """
        ),
    },
    {
        "no": "24",
        "title": "BHF 24 | SIS und Falltraining fuer die Pruefung",
        "subtitle": "Fallübersicht, SIS und die Faelle Frau Weiss bis Nadine systematisch bearbeiten",
        "filename": "BHF_24_SIS_Falltraining_Pruefung.pptx",
        "accent": rgb("3B7996"),
        "accent_dark": rgb("2B5A71"),
        "topics": L(
            """
            Einfuehrung
            Falluebersicht
            SIS
            Fall 1 Frau Weiss
            Fall 2 Frau Feldmann
            Fall 3 Herr Gaertner
            Fall 4 Lara
            Fall 5 Studie
            Fall 6 Frau Falk
            Fall 7 Leon
            Fall 8 Frau Wunsch
            Fall 9 Lena
            Fall 10 Nadine
            """
        ),
        "focus_cards": [
            {
                "title": "Fallübersicht zuerst",
                "tone": "sky",
                "body": L(
                    """
                    bei jedem Fall zuerst Situation ordnen, nicht sofort Loesungen springen
                    Diagnosen, Symptome, Risiken und Ressourcen trennen
                    """
                ),
            },
            {
                "title": "SIS nutzen",
                "tone": "mint",
                "body": L(
                    """
                    strukturiert Informationen sammeln und auf den Alltag beziehen
                    Pflege wird dadurch klarer und nachvollziehbarer
                    """
                ),
            },
            {
                "title": "Priorisieren",
                "tone": "sand",
                "body": L(
                    """
                    Was ist akut? Was ist gefaehrlich? Was belastet am meisten?
                    erst dann Ziele und Massnahmen formulieren
                    """
                ),
            },
            {
                "title": "Pruefungssicherheit",
                "tone": "rose",
                "body": L(
                    """
                    Struktur schlägt Detailchaos
                    Problem + Ziel + Massnahme + Evaluation fast immer mitnennen
                    """
                ),
            },
        ],
        "blocks_title": "Die Fallstruktur, die fast immer funktioniert",
        "blocks_subtitle": "BHF 24 ist dein direkter Uebergang in die Klausur- und Fallpruefung.",
        "blocks": [
            {
                "title": "SIS",
                "tone": "sky",
                "body": L(
                    """
                    Informationen strukturiert sammeln und mit Lebenssituation verbinden
                    nicht nur Defizite, sondern auch Ressourcen sichtbar machen
                    """
                ),
            },
            {
                "title": "Fallanalyse",
                "tone": "mint",
                "body": L(
                    """
                    1. Lage erfassen
                    2. Risiken priorisieren
                    3. Ziele formulieren
                    4. Massnahmen + Evaluation planen
                    """
                ),
            },
            {
                "title": "Typische Fallthemen",
                "tone": "sand",
                "body": L(
                    """
                    Mobilitaet, Schmerzen, Ausscheidung, Kommunikation, Selbstversorgung
                    psychosoziale Aspekte nie vergessen
                    """
                ),
            },
            {
                "title": "Faelle 1-10",
                "tone": "rose",
                "body": L(
                    """
                    Frau Weiss, Frau Feldmann, Herr Gaertner, Lara, Leon, Nadine u.a.
                    Namen lernen ist weniger wichtig als die systematische Bearbeitung
                    """
                ),
            },
        ],
        "check_left": L(
            """
            SIS in einfachen Worten erklaeren
            Fallanalyse in 4 Schritten darstellen
            Priorisierung in einem Fall begruenden
            Problem, Ziel und Massnahme verknuepfen
            """
        ),
        "check_right": L(
            """
            Im Fall immer zuerst ordnen, dann handeln
            Ressourcen sind genauso wichtig wie Probleme
            Gute Pruefungsantworten sind klar, konkret und begruendet
            """
        ),
    },
]


def set_fill(shape, color, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1.0, transparency: float = 0.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_box(
    slide,
    x,
    y,
    w,
    h,
    text="",
    font_size=18,
    color=TEXT,
    bold=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_lines(
    slide,
    x,
    y,
    w,
    h,
    lines,
    font_size=15,
    color=TEXT,
    bold_first=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    line_spacing=1.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and idx == 0
        run.font.color.rgb = color
    return box


def add_background(slide, accent, accent_dark):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_fill(bg, LIGHT)
    set_line(bg, LIGHT, 0)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.24))
    set_fill(top, accent)
    set_line(top, accent, 0)
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(7.18), Inches(13.333), Inches(0.32))
    set_fill(bottom, rgb("EEF4F7"))
    set_line(bottom, rgb("EEF4F7"), 0)
    for x, y, size, tr in [(11.3, 0.55, 0.62, 0.8), (11.95, 0.95, 0.34, 0.84), (10.78, 1.0, 0.26, 0.86)]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        set_fill(circle, accent_dark, tr)
        set_line(circle, accent_dark, 0.5, 0.9)


def add_footer(slide, text, no, total, accent):
    add_box(slide, 0.55, 7.16, 10.8, 0.15, text, 9, MUTED)
    add_box(slide, 12.0, 7.12, 0.8, 0.18, f"{no}/{total}", 10, accent, True, font_name="Aptos Display", align=PP_ALIGN.RIGHT)


def add_tag(slide, text, accent):
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.48), Inches(1.55), Inches(0.34)
    )
    set_fill(tag, accent)
    set_line(tag, accent, 0.8)
    add_box(slide, 0.68, 0.51, 1.42, 0.24, text, 11, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_card(slide, x, y, w, h, title, body_lines, fill_color, accent, body_size=14):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, fill_color)
    set_line(shape, accent, 1.0, 0.25)
    add_box(slide, x + 0.14, y + 0.12, w - 0.28, 0.34, title, 17, accent, True, font_name="Aptos Display")
    add_lines(slide, x + 0.14, y + 0.5, w - 0.28, h - 0.62, body_lines, body_size, TEXT)


def add_title_slide(prs, deck, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.86), Inches(6.35), Inches(5.65)
    )
    set_fill(panel, deck["accent"])
    set_line(panel, deck["accent"], 0.8)
    add_box(slide, 0.95, 1.0, 2.2, 0.3, deck["no"], 13, WHITE, True)
    add_box(slide, 0.95, 1.4, 5.55, 1.85, deck["title"], 27, WHITE, True, font_name="Aptos Display")
    add_lines(slide, 0.95, 3.55, 5.4, 1.15, [deck["subtitle"]], 16, rgb("E7F3FA"))
    add_card(
        slide,
        7.35,
        1.0,
        5.15,
        1.5,
        "Lernziel",
        [
            "alle OneNote-Themen dieses BHF abdecken",
            "Krankheitslehre im DURST-Format lernen",
            "fuer die Pruefung strukturiert antworten",
        ],
        TONES["sky"],
        deck["accent_dark"],
    )
    add_card(
        slide,
        7.35,
        2.8,
        5.15,
        1.45,
        "Pruefungsstrategie",
        [
            "Definition oder Einordnung zuerst",
            "dann Symptome/Risiken",
            "dann Therapie, Pflege und Begruendung",
        ],
        TONES["sand"],
        deck["accent_dark"],
    )
    add_card(
        slide,
        7.35,
        4.55,
        5.15,
        1.3,
        "Merksatz",
        [
            "Struktur schlaegt Chaos: Problem, Ziel, Massnahme, Evaluation.",
        ],
        TONES["mint"],
        deck["accent_dark"],
    )
    add_footer(slide, "Exam Library | aus den BHF-Themen im lokalen OneNote abgeleitet", no, total, deck["accent"])


def split_in_half(items):
    mid = (len(items) + 1) // 2
    return items[:mid], items[mid:]


def add_topics_slide(prs, deck, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, deck["no"], deck["accent"])
    add_box(slide, 0.72, 0.94, 9.5, 0.42, "Themen aus dem OneNote", 24, INK, True, font_name="Aptos Display")
    add_box(slide, 0.75, 1.3, 10.6, 0.25, "Bereinigte Themenliste fuer dieses BHF. Das ist die inhaltliche Vollstaendigkeit der Deck-Basis.", 12, MUTED)
    left, right = split_in_half(deck["topics"])
    left_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.8), Inches(5.9), Inches(4.95)
    )
    right_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.78), Inches(1.8), Inches(5.85), Inches(4.95)
    )
    set_fill(left_shape, TONES["sky"])
    set_fill(right_shape, TONES["white"])
    set_line(left_shape, deck["accent_dark"], 1.0, 0.25)
    set_line(right_shape, deck["accent_dark"], 1.0, 0.25)
    add_lines(slide, 0.94, 2.02, 5.38, 4.4, [f"• {item}" for item in left], 15, TEXT)
    add_lines(slide, 7.0, 2.02, 5.28, 4.4, [f"• {item}" for item in right], 15, TEXT)
    add_footer(slide, "Diese Themenliste ist die Referenz fuer die folgenden Lern- und DURST-Slides.", no, total, deck["accent"])


def add_focus_slide(prs, deck, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, "Fokus", deck["accent"])
    add_box(slide, 0.72, 0.94, 9.2, 0.42, "Pruefungsfokus fuer dieses BHF", 24, INK, True, font_name="Aptos Display")
    add_box(slide, 0.75, 1.3, 10.0, 0.25, "Wenn du diese vier Felder sicher kannst, bist du in der Pruefung schon sehr stabil.", 12, MUTED)
    positions = [
        (0.72, 1.75, 5.85, 2.05),
        (6.78, 1.75, 5.85, 2.05),
        (0.72, 4.0, 5.85, 2.05),
        (6.78, 4.0, 5.85, 2.05),
    ]
    for pos, card in zip(positions, deck["focus_cards"]):
        add_card(slide, *pos, card["title"], card["body"], TONES[card["tone"]], deck["accent_dark"])
    add_footer(slide, "Fokus-Slides geben dir die schnellste Wiederholung kurz vor der Klausur.", no, total, deck["accent"])


def chunk(items, size):
    return [items[i : i + size] for i in range(0, len(items), size)]


def add_block_slide(prs, deck, group, index, groups_len, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, "DURST / Transfer", deck["accent"])
    title = deck["blocks_title"] if groups_len == 1 else f"{deck['blocks_title']} ({index}/{groups_len})"
    add_box(slide, 0.72, 0.92, 9.5, 0.42, title, 24, INK, True, font_name="Aptos Display")
    add_box(slide, 0.75, 1.28, 10.5, 0.26, deck["blocks_subtitle"], 12, MUTED)
    positions = [
        (0.72, 1.72, 5.85, 2.15),
        (6.78, 1.72, 5.85, 2.15),
        (0.72, 4.05, 5.85, 2.15),
        (6.78, 4.05, 5.85, 2.15),
    ]
    for pos, block in zip(positions, group):
        add_card(slide, *pos, block["title"], block["body"], TONES[block["tone"]], deck["accent_dark"], body_size=13)
    add_footer(slide, "Bei Krankheitslehre immer D, Formen/Typen, Risiken, Symptome und Pflege mitdenken.", no, total, deck["accent"])


def add_checklist_slide(prs, deck, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, "Checklist", deck["accent"])
    add_box(slide, 0.72, 0.94, 9.0, 0.42, "Pruefungscheck fuer dieses BHF", 24, INK, True, font_name="Aptos Display")
    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.78), Inches(5.9), Inches(4.95)
    )
    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.78), Inches(1.78), Inches(5.85), Inches(4.95)
    )
    set_fill(left, TONES["mint"])
    set_fill(right, TONES["sand"])
    set_line(left, deck["accent_dark"], 1.0, 0.25)
    set_line(right, deck["accent_dark"], 1.0, 0.25)
    add_box(slide, 0.95, 1.98, 5.2, 0.28, "Das solltest du frei sagen koennen", 17, deck["accent_dark"], True, font_name="Aptos Display")
    add_lines(slide, 0.95, 2.3, 5.2, 4.1, [f"• {line}" for line in deck["check_left"]], 15, TEXT)
    add_box(slide, 7.0, 1.98, 5.15, 0.28, "Merksaetze fuer die letzten 48 Stunden", 17, deck["accent_dark"], True, font_name="Aptos Display")
    add_lines(slide, 7.0, 2.3, 5.15, 4.1, [f"• {line}" for line in deck["check_right"]], 15, TEXT)
    add_footer(slide, "Wiederhole diesen Slide am Ende jedes Lernblocks laut und in eigenen Worten.", no, total, deck["accent"])


def build_deck(deck):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = deck["title"]
    prs.core_properties.subject = deck["subtitle"]
    prs.core_properties.author = "Codex"
    prs.core_properties.comments = "Vollstaendige BHF-Exam-Library auf Basis des lokalen OneNote"

    block_groups = chunk(deck["blocks"], 4)
    total = 4 + len(block_groups)
    slide_no = 1
    add_title_slide(prs, deck, slide_no, total)
    slide_no += 1
    add_topics_slide(prs, deck, slide_no, total)
    slide_no += 1
    add_focus_slide(prs, deck, slide_no, total)
    slide_no += 1
    for idx, group in enumerate(block_groups, start=1):
        add_block_slide(prs, deck, group, idx, len(block_groups), slide_no, total)
        slide_no += 1
    add_checklist_slide(prs, deck, slide_no, total)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / deck["filename"]
    prs.save(str(path))
    return path


def build_all():
    paths = [build_deck(deck) for deck in DECKS]
    for path in paths:
        print(f"saved {path}")


if __name__ == "__main__":
    build_all()
