from __future__ import annotations

import json
import re
import unicodedata
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Iterable
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("/Users/arielhavana/antigr/porc/mockups/docx_theme_powerpoints")
MANIFEST_PATH = OUT_DIR / "theme_manifest.json"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


WHITE = rgb("FFFFFF")
INK = rgb("223247")
TEXT = rgb("32465B")
MUTED = rgb("61758D")
LIGHT = rgb("F8FBFD")

TONES = {
    "sky": rgb("E7F3FA"),
    "mint": rgb("E7F4EE"),
    "sand": rgb("FCF3DA"),
    "rose": rgb("F9E6EA"),
    "white": WHITE,
}

ACCENTS = [
    ("1E7A8C", "175E6D"),
    ("2F6F9F", "214F72"),
    ("7C5D9A", "604776"),
    ("4F8A53", "37673A"),
    ("AF6B1F", "884F12"),
    ("8B4E8E", "67386A"),
]

OCR_FALLBACKS = {
    "02": """
    BURGSTRASSE
    (Graphische) Lernfelduebersicht
    BHF 2: Menschen in ihrer Orientierung und Beweglichkeit unterstuetzen und Sicherheit mitgestalten
    BHF 2
    Menschen in ihrer Orientierung und Beweglichkeit unterstuetzen und Sicherheit mitgestalten
    Grundzuege der Bewegung und Pflegehandeln bei Bewegungseinschraenkung
    Ich lerne etwas ueber:
    Bedeutung von Bewegung fuer den Menschen
    Anatomie und Physiologie des Bewegungsapparates
    Richtungsbezeichnungen des Koerpers
    Veraenderung der Motorik im Alter
    Erkrankungen des Bewegungsapparates (z.B. Arthrose)
    Risiken bei eingeschraenkter Mobilitaet
    Pflegerische Prophylaxen bei eingeschraenkter Mobilitaet (z.B. Kontrakturenprophylaxe)
    Bewegungsinteraktion gestalten und rechtliche Grundlagen der Bewegungsunterstuetzung
    Ich lerne etwas ueber:
    Foerderung der Sicherheit eines zu pflegenden Menschen
    Massnahmen zur Vermeidung von Stuerzen
    Massnahmen der ersten Hilfe nach einem Sturz
    Haftungsrechtliche Regelungen in der Pflege
    Massnahmen zur Verhuetung von Unfaellen im Pflegealltag
    Persoenliche Gesunderhaltung in Bezug auf Bewegung
    Ich lerne etwas ueber:
    Reflektion des eigenen Bewegungsverhaltens
    Grenzen der Belastbarkeit des eigenen Bewegungsapparates
    Berufliche Gesundheitsrisiken im Hinblick auf den Bewegungsapparat
    Rueckengerechtes Arbeiten im Pflegealltag
    Nutzung von technischen Hilfsmitteln im Pflegealltag
    Pflegeinterventionen zur Bewegungsaktivierung (Kinaesthetik)
    Ich lerne etwas ueber:
    Erfassung und Analyse von menschlichen Bewegungsmustern
    Entwicklung und Grundlagen der Kinaesthetik
    Anwendung von kinaesthetischen Prinzipien im Pflegealltag
    Menschen zum Handeln und zur Bewegung anleiten
    Ich lerne etwas ueber:
    Bedeutung von Anleitung in der Pflege
    Grundprinzipien der Anleitung zu einfachen Handlungen und Bewegungsablaeufen
    Menschen mit Orientierungsstoerungen begleiten
    Ich lerne etwas ueber:
    Orientierung und Orientierungsstoerungen
    Kommunikation und Interaktion mit Menschen, die in ihrer Orientierung eingeschraenkt sind
    Verwendung von Orientierungshilfen im Pflegealltag
    Basiswissen ueber die Erkrankung Demenz
    """,
    "20": """
    BURGSTRASSE
    Uebersicht BHF 20
    BHF 20: Professionelles Handeln bei taburelevanten Themen und bei der Umsetzung von qualitaetsentwickelnden Massnahmen
    BHF 20
    Professionelles Handeln bei taburelevanten Themen und bei der Umsetzung von qualitaetsentwickelnden Massnahmen
    Menschen mit Erkrankungen des Urogenitalsystems pflegen und beraten
    Ich lerne etwas ueber:
    Anatomie und Physiologie des Urogenitalsystems
    ausgewaehlte Erkrankungen des Urogenitalsystems
    Beratung und Pflege von Menschen mit chronischer Niereninsuffizienz und Dialyse
    Menschen mit Kontinenzproblemen beraten und schulen
    Ich lerne etwas ueber:
    Harninkontinenz
    Anzeichen, Risikofaktoren, Formen, Kontinenzprofile
    Massnahmen zur Foerderung der Harnkontinenz
    Beckenbodentraining, Blasen- und Toilettentraining, Hilfsmittel
    Initiierung, Konzeption und Durchfuehrung von Beratungsangeboten
    Qualitaetssicherndes Handeln
    Ich lerne etwas ueber:
    Die Bedeutung des Qualitaetsmanagements
    Grundlegende Begriffe des Qualitaetsmanagements
    Prozesse und Massnahmen der internen Qualitaetssicherung
    Einarbeitung neuer Mitarbeiter und Anleitung von Praktikanten
    Weiblich maennlich divers
    Ich lerne etwas ueber:
    Sexuelle Orientierung und geschlechtliche Vielfalt
    Begriffserklaerung Transsexuell oder Transident
    Geschlechtliche Identitaet
    Transgender Kinder
    Diversitaetssensible Pflege
    """,
}

DOCS = [
    {
        "id": "01B",
        "bhf": "BHF 1B",
        "path": Path("/Users/arielhavana/Downloads/Themenstruktur_Schüler BHF 1B.docx"),
        "mode": "boundary",
        "themes": [
            {"title": "Vertiefung Informationsgespraeche", "match": "Vertiefung Informationsgespräche"},
            {"title": "Verstaendnis des Berufsbildes", "match": "Verständnis des Berufsbildes"},
            {"title": "Auswertung des ersten Praxiseinsatzes", "match": "Auswertung des ersten Praxiseinsatzes"},
            {"title": "Orientierung im System Familie", "match": "Orientierung im System Familie"},
        ],
    },
    {
        "id": "02",
        "bhf": "BHF 2",
        "path": Path("/Users/arielhavana/Downloads/bhf 2.docx"),
        "mode": "boundary",
        "themes": [
            {
                "title": "Grundzuege der Bewegung und Pflegehandeln bei Bewegungseinschraenkung",
                "match": "Grundzuege der Bewegung und Pflegehandeln bei Bewegungseinschraenkung",
                "kind": "disease_bundle",
            },
            {
                "title": "Bewegungsinteraktion und rechtliche Grundlagen der Bewegungsunterstuetzung",
                "match": "Bewegungsinteraktion gestalten und rechtliche Grundlagen der Bewegungsunterstuetzung",
            },
            {
                "title": "Persoenliche Gesunderhaltung in Bezug auf Bewegung",
                "match": "Persoenliche Gesunderhaltung in Bezug auf Bewegung",
            },
            {
                "title": "Pflegeinterventionen zur Bewegungsaktivierung (Kinaesthetik)",
                "match": "Pflegeinterventionen zur Bewegungsaktivierung (Kinaesthetik)",
            },
            {
                "title": "Menschen zum Handeln und zur Bewegung anleiten",
                "match": "Menschen zum Handeln und zur Bewegung anleiten",
            },
            {
                "title": "Menschen mit Orientierungsstoerungen begleiten",
                "match": "Menschen mit Orientierungsstoerungen begleiten",
                "kind": "disease_bundle",
            },
        ],
    },
    {
        "id": "03A",
        "bhf": "BHF 3A",
        "path": Path("/Users/arielhavana/Downloads/Themenstruktur Schüler BHF3A.docx"),
        "mode": "boundary",
        "themes": [
            {"title": "Im Pflegealltag des Einsatzbereiches mitwirken", "match": "Im Pflegealltag des Einsatzbereiches mitwirken"},
            {"title": "Grundlagen des hygienischen Handelns", "match": "Grundlagen des hygienischen Handelns"},
            {"title": "Haut und Hautzustaende", "match": "Haut und Hautzustände", "kind": "disease_bundle"},
            {"title": "Koerperpflege unterstuetzen", "match": "Körperpflege unterstützen"},
            {"title": "Beruehrung und Intimpflege", "match": "Berührung – Interaktion bei der körpernahen Versorgung"},
            {"title": "Erste Grundlagen der Ernaehrung und Ausscheidung", "match": "Erste Grundlagen der Ernährung und Ausscheidung", "kind": "disease_bundle"},
            {"title": "Recht und Pflegeethik", "match": "Recht und Pflegeethik"},
        ],
    },
    {
        "id": "03B",
        "bhf": "BHF 3B",
        "path": Path("/Users/arielhavana/Downloads/Übersiche Schüler 3B.docx"),
        "mode": "boundary",
        "themes": [
            {"title": "Vertiefung des Pflegeprozesses und seiner Bedeutung", "match": "Vertiefung des Pflegeprozesses und seiner Bedeutung"},
            {"title": "Im Pflegealltag des Einsatzbereiches mitwirken und reflektieren", "match": "Im Pflegealltag des Einsatzbereiches mitwirken/ Reflexion"},
            {
                "title": "Unterstuetzung bei der Nahrungsaufnahme und Risiko Mangelernaehrung",
                "match": "Unterstützung bei der Nahrungsaufnahme/ Risiko Mangelernährung",
                "kind": "disease_bundle",
            },
            {"title": "Pflegerische Versorgung von Saeuglingen", "match": "Pflegerische Versorgung von Säuglingen", "kind": "disease_bundle"},
            {"title": "Grundbegriffe der Pflegeethik und pflegerischer Kommunikation", "match": "Grundbegriffe der Pflegeethik und pflegerischer Kommunikation"},
            {"title": "Kultursensibilitaet und Familienorientierung im Pflegeprozess", "match": "Kultursensibilität und Familienorientierung als Bezugspunkte im Pflegeprozess"},
        ],
    },
    {
        "id": "04",
        "bhf": "BHF 4",
        "path": Path("/Users/arielhavana/Downloads/BHF4_Themenstruktur_Schüler .docx"),
        "mode": "boundary",
        "themes": [
            {
                "title": "Theorieblock I: Urogenitalsystem, Fertilisation und Schwangerschaftsdiagnostik",
                "match": "Theorieblock I (10 Std.)",
                "kind": "disease_bundle",
            },
            {
                "title": "Theorieblock II: Schutz von Mutter und Kind, Familienhilfen und Erste Hilfe",
                "match": "Theorieblock II (8 Std.)",
            },
            {
                "title": "Lernsituation I Frau Schnelle: Schwangerschaft, Geburt und Wochenbett",
                "match": "Lernsituation I „Frau Schnelle“ (15 Std.)",
                "kind": "disease_bundle",
            },
            {
                "title": "Lernsituation II Endlich ist Paul da: Versorgung des Neugeborenen",
                "match": "Lernsituation II „Endlich ist Paul da“",
                "kind": "disease_bundle",
            },
        ],
    },
    {
        "id": "06",
        "bhf": "BHF 6",
        "path": Path("/Users/arielhavana/Downloads/BHF Übersicht 6.docx"),
        "mode": "boundary",
        "themes": [
            {"title": "Zu ermittelnde Vitalparameter und ihre Aussagekraft", "match": "Zu ermittelnde Vitalparameter und ihre Aussagekraft"},
            {
                "title": "Erstmassnahmen in Notfallsituationen",
                "match": "Erstmaßnahmen in Notfallsituationen",
                "kind": "disease_bundle",
            },
            {"title": "Arbeiten in der Notaufnahme", "match": "Arbeiten in der Notaufnahme", "kind": "disease_bundle"},
            {"title": "Durchfuehrung einer Reanimation", "match": "Durchführung einer Reanimation"},
            {"title": "Rechtliche Grundlagen der Notfallversorgung", "match": "Rechtliche Grundlagen der Notfallversorgung"},
            {
                "title": "Erstmassnahmen in verschiedenen Notfallsituationen",
                "match": "Erstmaßnahmen in verschiedenen Notfallsituationen",
                "kind": "disease_bundle",
            },
        ],
    },
    {
        "id": "07",
        "bhf": "BHF 7",
        "path": Path("/Users/arielhavana/Downloads/bhf7.docx"),
        "mode": "boundary",
        "themes": [
            {
                "title": "Arbeiten in der ambulanten Pflege",
                "match": "Arbeiten in der ambulanten Pflege",
            },
            {
                "title": "Menschen mit Diabetes mellitus versorgen",
                "match": "Menschen mit Diabetes mellitus versorgen",
                "kind": "disease",
            },
            {
                "title": "Menschen mit Hoer- und Seheinschraenkungen versorgen",
                "match": "Menschen mit Hör- und Seheinschränkungen versorgen",
                "kind": "disease_bundle",
            },
            {
                "title": "Menschen nach einer ambulanten OP versorgen",
                "match": "Menschen nach einer ambulanten OP versorgen",
            },
            {
                "title": "Gesundheitsfoerderung im haeuslichen Umfeld",
                "match": "Die Gesundheitsförderung beachten",
            },
        ],
    },
    {
        "id": "08",
        "bhf": "BHF 8",
        "path": Path("/Users/arielhavana/Downloads/Themenstruktur BHF8.docx"),
        "mode": "boundary",
        "themes": [
            {
                "title": "Henriette Schulz: multimorbide Patienten versorgen",
                "match": "Henriette Schulz – multimorbide Patienten versorgen",
                "kind": "disease_bundle",
            },
            {"title": "Kind mit Asthma", "match": "Kind mit Asthma", "kind": "disease"},
            {"title": "Gesundheitfoerderung, Resilienz und Adhaerenz", "match": "Gesundheitförderung"},
            {"title": "Kollegiale Fallberatung", "match": "Kollegiale Fallberatung"},
            {"title": "Krankenhausfinanzierung", "match": "Krankenhausfinanzierung"},
        ],
    },
    {
        "id": "10",
        "bhf": "BHF 10",
        "path": Path("/Users/arielhavana/Downloads/Themenstruktur BHF10.docx"),
        "mode": "boundary",
        "themes": [
            {"title": "Pflegeprozess in der stationaeren Langzeitpflege", "match": "Pflegeprozess in der stationären Langzeitpflege"},
            {"title": "Organisation in der stationaeren Langzeitpflege", "match": "Organisation in der stationären Langzeitpflege"},
            {"title": "Hygienisches Handeln in der stationaeren Langzeitpflege", "match": "Hygienisches Handeln in der stationären Langzeitpflege"},
            {"title": "Existentielle Erfahrungen und Gewalt", "match": "Existentielle Erfahrungen und Gewalt"},
            {"title": "Eigene Professionalitaet", "match": "Eigene Professionalität"},
        ],
    },
    {
        "id": "13",
        "bhf": "BHF 13",
        "path": Path("/Users/arielhavana/Downloads/Themenstruktur Schüler BHF 13.docx"),
        "mode": "boundary",
        "themes": [
            {"title": "Chronische Erkrankungen", "match": "Chronische Erkrankungen", "kind": "disease_bundle"},
            {"title": "Umgang mit chronischen Erkrankungen", "match": "Umgang mit chronischen Erkrankungen"},
            {"title": "Chronische Schmerzen", "match": "Chronische Schmerzen", "kind": "disease"},
            {"title": "Foerderung der Mobilitaet", "match": "Förderung der Mobilität"},
            {"title": "Belastungen der Pflegenden", "match": "Belastungen der Pflegenden"},
            {"title": "Integrierte Versorgungsprozesse", "match": "Integrierte Versorgungsprozesse"},
        ],
    },
    {
        "id": "16",
        "bhf": "BHF 16",
        "path": Path("/Users/arielhavana/Downloads/bhf16.docx"),
        "mode": "keyword",
        "themes": [
            {
                "title": "Interventionen in der Psychiatrie",
                "keywords": ["Bezugspflege", "Milieutherapie", "Gesprächsführung", "Tagesstruktur", "Deeskalationsstrategien", "Recovery", "Empowerment", "SDM-Modell"],
                "fallback_lines": ["Bezugspflege", "Milieutherapie", "Gesprächsfuehrung", "Tagesstruktur", "Deeskalationsstrategien", "Recovery-Konzept", "Empowerment", "SDM-Modell"],
            },
            {
                "title": "Anorexia nervosa",
                "keywords": ["Anorexia nervosa", "Krankheitsbild Anorexia", "Selbsthilfe", "Rechtliche Grundlagen"],
                "kind": "disease",
                "fallback_lines": ["Krankheitsbild Anorexia nervosa", "Selbsthilfe", "Rechtliche Grundlagen"],
            },
            {
                "title": "Geschichte der Psychiatrie",
                "keywords": ["Geschichte der Psychiatrie", "Erklärungsversuche", "Antike bis heute", "Konzept Hoffnung"],
                "fallback_lines": ["Erklaerungsversuche und Umgang mit psychischen Erkrankungen von der Antike bis heute", "Konzept Hoffnung"],
            },
            {
                "title": "Depression",
                "keywords": ["Depression", "Schweregrade", "Psychoedukation", "Suizidalität", "Co-Abhängigkeit", "Epidemiologische Relevanz"],
                "kind": "disease",
                "fallback_lines": ["Abgrenzung zu negativen Befindlichkeiten", "Epidemiologische Relevanz", "Risikofaktoren, Diagnostik, Verlauf und Schweregrade", "Psychoedukation", "Suizidalitaet", "Auswirkungen auf Kinder chronisch kranker Eltern", "Co-Abhaengigkeit"],
            },
            {
                "title": "Schizophrenie und der Film Das weisse Rauschen",
                "keywords": ["Schizophrenie", "Das weiße Rauschen", "Säulen der psychiatrischen Therapie", "Therapie & Pflegerische Besonderheiten"],
                "kind": "disease",
                "fallback_lines": ["Krankheitsbild Schizophrenie", "Film Das weisse Rauschen", "Saeulen der psychiatrischen Therapie", "Therapie und pflegerische Besonderheiten"],
            },
        ],
    },
    {
        "id": "17",
        "bhf": "BHF 17",
        "path": Path("/Users/arielhavana/Downloads/bhf17.docx"),
        "mode": "boundary",
        "themes": [
            {
                "title": "Herausfordernde Pflegesituationen",
                "match": "Herausfordernde Pflegesituationen",
                "fallback_lines": ["Diversitaetssensible Pflege", "Diversitaetsmerkmale", "Sunrise-Modell nach Leininger", "Partizipative Entscheidungsfindung in sozial herausfordernden Situationen"],
            },
            {
                "title": "Gesundheitsfoerderung und Gesunderhaltung",
                "match": "Gesundheitsförderung- und erhaltung:",
                "fallback_lines": ["Selbst- und Fremdgefaehrdung durch Infektionskrankheiten", "Massnahmen des Selbstschutzes", "Postexpositionsprophylaxe", "Salutogenese", "Gesunderhaltung in sozial schwierigen Situationen"],
            },
            {
                "title": "Infektionskrankheiten und Infektionsschutz",
                "match": "Infektionskrankheiten",
                "split_markers": ["Infektionskrankheiten Ich festige mein Wissen über:"],
                "kind": "disease_bundle",
                "fallback_lines": ["Erregerarten", "Infektionswege", "Infektionsschutzgesetz", "Skabies", "Tuberkulose", "Hepatitis B", "COVID-19", "Impfungen als Massnahme des Infektionsschutzes"],
            },
        ],
    },
    {
        "id": "19",
        "bhf": "BHF 19",
        "path": Path("/Users/arielhavana/Downloads/Themenstruktur Schüler BHF 19.docx"),
        "mode": "boundary",
        "themes": [
            {"title": "Menschen mit Aphasie professionell begleiten", "match": "Menschen mit Aphasie professionell begleiten", "kind": "disease_bundle"},
            {"title": "Pflege eines Jugendlichen mit Cerebralparese und Rehabilitation", "match": "Pflege eines Jugendlichen mit einer Cerebralparese und Rehabilitation", "kind": "disease_bundle"},
            {"title": "Die Pflege eines Menschen nach Hirninfarkt planen und evaluieren", "match": "Die Pflege eines Menschen nach Hirninfarkt planen und evaluieren", "kind": "disease_bundle"},
        ],
    },
    {
        "id": "20",
        "bhf": "BHF 20",
        "path": Path("/Users/arielhavana/Downloads/BHF20.docx"),
        "mode": "boundary",
        "themes": [
            {"title": "Menschen mit Erkrankungen des Urogenitalsystems pflegen und beraten", "match": "Menschen mit Erkrankungen des Urogenitalsystems pflegen und beraten", "kind": "disease_bundle"},
            {"title": "Menschen mit Kontinenzproblemen beraten und schulen", "match": "Menschen mit Kontinenzproblemen beraten und schulen", "kind": "disease_bundle"},
            {"title": "Qualitaetssicherndes Handeln", "match": "Qualitaetssicherndes Handeln"},
            {"title": "Weiblich maennlich divers", "match": "Weiblich maennlich divers"},
        ],
    },
    {
        "id": "21",
        "bhf": "BHF 21",
        "path": Path("/Users/arielhavana/Downloads/BHF 21_Übersicht SuS.docx"),
        "mode": "boundary",
        "themes": [
            {
                "title": "Fallbezogene Erarbeitung und Planung einer komplexen Pflegesituation",
                "match": "Fallbezogene Erarbeitung und Planung einer komplexen Pflegesituation",
            },
            {
                "title": "Betreuung eines Kindes mit einer chronischen Darmerkrankung (CED)",
                "match": "Betreuung eines Kindes mit einer chronischen Darmerkrankung (CED)",
                "kind": "disease_bundle",
            },
            {
                "title": "Umgang mit einer Krisensituation am Beispiel eines Schreibabys",
                "match": "Umgang mit einer Krisensituation am Beispiel eines „Schreibabys“",
            },
            {
                "title": "Pflege eines Fruehgeborenen",
                "match": "Pflege eines Frühgeborenen",
                "kind": "disease_bundle",
            },
            {
                "title": "Perspektivwechsel und Mitverantwortung fuer Arbeitsprozesse",
                "match": "Perspektivwechsel – Mitverantwortung für die Organisation und Gestaltung von Pflegeprozessen",
            },
        ],
    },
]


NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
NOISE_LINE_FRAGMENTS = [
    "zielgruppe",
    "zielgruppen",
    "setting",
    "kompetenzen",
    "zeitlicher umfang",
    "ich kann",
    "ersteller/datum",
    "bearbeitet/datum",
    "seite 1 von 1",
]


def fold_ascii(value: str) -> str:
    value = value.replace("ß", "ss").replace("ẞ", "SS")
    return unicodedata.normalize("NFKD", value).encode("ascii", "ignore").decode("ascii")


def clean_space(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def slugify(value: str) -> str:
    slug = fold_ascii(value)
    slug = re.sub(r"[^A-Za-z0-9]+", "_", slug)
    return slug.strip("_")


def read_doc_text(doc: dict) -> str:
    parts: list[str] = []
    with ZipFile(doc["path"]) as zf:
        xml = zf.read("word/document.xml")
    root = ET.fromstring(xml)
    for paragraph in root.findall(".//w:p", NS):
        paragraph_parts = []
        for node in paragraph.findall(".//w:t", NS):
            if node.text:
                paragraph_parts.append(node.text)
        text = clean_space("".join(paragraph_parts))
        if text:
            parts.append(text)
    for elem in root.iter():
        tag = elem.tag.split("}")[-1]
        if tag in {"docPr", "cNvPr"}:
            descr = clean_space(elem.attrib.get("descr", ""))
            if descr:
                parts.append(descr)
    text = "\n".join(parts)
    if len(fold_ascii(text)) < 200 and doc["id"] in OCR_FALLBACKS:
        text = f"{text}\n{OCR_FALLBACKS[doc['id']]}"
    return text


def clean_lines(text: str) -> list[str]:
    normalized = text.replace("Computergenerierter Alternativtext:", "")
    normalized = normalized.replace("Þ", "\n• ")
    normalized = normalized.replace("•", "\n• ")
    normalized = normalized.replace("→", "\n• ")
    normalized = normalized.replace("›", "\n• ")
    normalized = normalized.replace("\xa0", " ")
    normalized = normalized.replace("\r", "\n")
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    lines = [clean_space(line) for line in normalized.split("\n")]
    lines = [line for line in lines if line]
    deduped: list[str] = []
    seen_recent: list[str] = []
    for line in lines:
        if deduped and deduped[-1] == line:
            continue
        if line in seen_recent[-8:]:
            continue
        deduped.append(line)
        seen_recent.append(line)
    return deduped


INLINE_LEARNING_MARKERS = [
    "Ich lerne etwas über:",
    "Ich festige mein Wissen über:",
    "Ich denke nach über:",
    "Ich reflektiere:",
    "Ich vertiefe meine Kenntnisse über:",
]


def split_on_markers(line: str, markers: list[str]) -> list[str]:
    segments = [clean_space(line)]
    while True:
        changed = False
        next_segments: list[str] = []
        for segment in segments:
            split_index: int | None = None
            for marker in markers:
                idx = segment.find(marker)
                if idx > 0 and (split_index is None or idx < split_index):
                    split_index = idx
            if split_index is None:
                next_segments.append(segment)
                continue
            head = clean_space(segment[:split_index])
            tail = clean_space(segment[split_index:])
            if head:
                next_segments.append(head)
            if tail:
                next_segments.append(tail)
            changed = True
        segments = next_segments
        if not changed:
            break
    return [segment for segment in segments if segment]


def expand_doc_lines(lines: list[str], doc: dict) -> list[str]:
    markers = list(INLINE_LEARNING_MARKERS)
    for theme in doc["themes"]:
        markers.extend(theme.get("split_markers", []))
        if theme.get("match") and not theme.get("split_markers"):
            markers.append(theme["match"])
    markers = dedupe_keep_order(markers)
    expanded: list[str] = []
    for line in lines:
        expanded.extend(split_on_markers(line, markers))
    return dedupe_keep_order(expanded)


def is_noise(line: str) -> bool:
    lowered = fold_ascii(line).lower()
    stripped = lowered.strip(" •.:-…")
    if not lowered:
        return True
    if any(
        stripped.startswith(prefix)
        for prefix in {
            "ich lerne etwas uber",
            "ich denke nach uber",
            "ich reflektiere",
            "ich vertiefe meine kenntnisse uber",
            "ich festige mein wissen uber",
        }
    ):
        return True
    if any(fragment in lowered for fragment in NOISE_LINE_FRAGMENTS):
        return True
    if re.fullmatch(r"\(?tk\s*\d+\)?", stripped):
        return True
    if stripped.startswith("tk"):
        return True
    return False


def dedupe_keep_order(items: Iterable[str]) -> list[str]:
    result: list[str] = []
    seen: set[str] = set()
    for item in items:
        key = fold_ascii(item).lower()
        key = re.sub(r"^[^a-z0-9]+", "", key)
        key = re.sub(r"[^a-z0-9]+$", "", key)
        key = re.sub(r"\s+", " ", key).strip()
        if not key or key in seen:
            continue
        seen.add(key)
        result.append(item)
    return result


def trim_line(value: str, limit: int = 104) -> str:
    value = clean_space(value)
    if len(value) <= limit:
        return value
    return f"{value[:limit - 3].rstrip()}..."


def tokens(value: str) -> set[str]:
    lowered = fold_ascii(value).lower()
    found = re.findall(r"[a-z0-9]+", lowered)
    stop = {
        "und",
        "der",
        "die",
        "das",
        "mit",
        "den",
        "dem",
        "des",
        "eine",
        "einer",
        "eines",
        "fuer",
        "zum",
        "zur",
        "nach",
        "auf",
        "von",
        "bei",
        "bzw",
        "ich",
        "lerne",
        "uber",
        "menschen",
        "pflege",
        "bfh",
        "bhf",
    }
    return {token for token in found if len(token) > 2 and token not in stop}


def find_match_index(lines: list[str], match: str, start: int = 0) -> int | None:
    needle = fold_ascii(match).lower()
    best_index: int | None = None
    best_score = -1
    for index in range(start, len(lines)):
        hay = fold_ascii(lines[index]).lower()
        score = -1
        if hay == needle:
            score = 3
        elif hay.startswith(needle):
            score = 2
        elif needle in hay:
            score = 1
        if score > best_score:
            best_index = index
            best_score = score
            if score == 3:
                break
    return best_index if best_score >= 0 else None


def extract_boundary_block(lines: list[str], theme: dict, next_theme: dict | None) -> list[str]:
    start = find_match_index(lines, theme["match"])
    if start is None:
        return theme.get("fallback_lines", [])
    end = len(lines)
    if next_theme:
        next_start = find_match_index(lines, next_theme["match"], start + 1)
        if next_start is not None:
            end = next_start
    block = lines[start:end]
    cleaned: list[str] = []
    for line in block:
        if is_noise(line):
            continue
        if clean_space(line) == clean_space(theme["match"]):
            continue
        cleaned.append(trim_line(line))
    cleaned = dedupe_keep_order(cleaned)
    return cleaned or theme.get("fallback_lines", [])


def boundary_theme_starts(lines: list[str], doc: dict) -> list[int | None]:
    starts: list[int | None] = []
    search_start = 0
    for theme in doc["themes"]:
        start = find_match_index(lines, theme["match"], search_start)
        if start is None:
            start = find_match_index(lines, theme["match"], 0)
        starts.append(start)
        if start is not None:
            search_start = start + 1
    return starts


def extract_keyword_block(lines: list[str], theme: dict) -> list[str]:
    keywords = theme.get("keywords", [])
    folded_keywords = [fold_ascii(keyword).lower() for keyword in keywords]
    collected: list[str] = []
    for index, line in enumerate(lines):
        folded_line = fold_ascii(line).lower()
        if any(keyword in folded_line for keyword in folded_keywords):
            window = lines[max(0, index - 1) : min(len(lines), index + 8)]
            for candidate in window:
                if is_noise(candidate):
                    continue
                collected.append(trim_line(candidate))
    cleaned = dedupe_keep_order(collected)
    if theme.get("fallback_lines"):
        cleaned = dedupe_keep_order(cleaned + theme["fallback_lines"])
    return cleaned[:8]


def extract_theme_lines(lines: list[str], doc: dict, theme_index: int) -> list[str]:
    theme = doc["themes"][theme_index]
    if doc["mode"] == "boundary":
        starts = boundary_theme_starts(lines, doc)
        start = starts[theme_index]
        if start is None:
            extracted = theme.get("fallback_lines", [])
        else:
            end = len(lines)
            for next_start in starts[theme_index + 1 :]:
                if next_start is not None and next_start > start:
                    end = next_start
                    break
            block = lines[start:end]
            cleaned: list[str] = []
            for line in block:
                if is_noise(line):
                    continue
                if clean_space(line) == clean_space(theme["match"]):
                    continue
                cleaned.append(trim_line(line))
            extracted = dedupe_keep_order(cleaned) or theme.get("fallback_lines", [])
    else:
        extracted = extract_keyword_block(lines, theme)
    return extracted[:8]


def infer_kind(theme: dict, lines: list[str]) -> str:
    if theme.get("kind"):
        return theme["kind"]
    corpus = fold_ascii(" ".join(lines + [theme["title"]])).lower()
    disease_signals = [
        "krankheitsbild",
        "diabetes",
        "asthma",
        "copd",
        "parkinson",
        "arthritis",
        "schmerzen",
        "schlaganfall",
        "aphasie",
        "demenz",
        "pneumonie",
        "inkontinenz",
        "hepatitis",
        "covid",
        "tuberkulose",
    ]
    return "disease_bundle" if any(signal in corpus for signal in disease_signals) else "generic"


def exam_focus(kind: str, title: str, lines: list[str]) -> list[str]:
    if kind in {"disease", "disease_bundle"}:
        focus = [
            "DURST: Definition, Ursachen/Risiken, Symptome und Therapie/Pflege sauber ordnen.",
            "Formen, Typen, Grade oder Verlauf extra nennen, wenn das Thema das hergibt.",
        ]
        if kind == "disease_bundle":
            focus.append("Wenn mehrere Krankheitsbilder im Thema vorkommen, jedes kurz getrennt wiederholen.")
        if any("notfall" in fold_ascii(line).lower() for line in lines + [title]):
            focus.append("Akutzeichen, Erstmassnahmen und Beobachtung priorisieren.")
        return focus[:3]
    return [
        "Definition in 1-2 Saetzen geben und direkt ein klares Praxisbeispiel nennen.",
        "Pflege, Kommunikation, Anleitung, Sicherheit oder Recht mit dem Thema verknuepfen.",
        "Bei Fallbezug: Problem, Ziel, Massnahme und Evaluation mitdenken.",
    ]


def learning_questions(kind: str, title: str, lines: list[str]) -> list[str]:
    bullet_preview = [line for line in lines[:3] if line]
    questions = [
        f"Was ist der Kern von '{title}' in 2 Saetzen?",
        f"Welche 3 Unterpunkte musst du zu '{title}' frei nennen koennen?",
    ]
    if kind in {"disease", "disease_bundle"}:
        questions.append("Welche Formen, Komplikationen oder Warnzeichen musst du unterscheiden?")
    else:
        questions.append("Welches konkrete Praxisbeispiel passt zu diesem Thema?")
    if bullet_preview:
        questions.append(f"Welche Aussage aus dem Dokument bleibt dir besonders wichtig: {bullet_preview[0]}?")
    return [trim_line(question, 115) for question in questions[:4]]


def set_fill(shape, color, transparency=0.0):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    fill.transparency = transparency


def set_line(shape, color, width=1.0, transparency=0.0):
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width)
    line.transparency = transparency


def add_box(
    slide,
    x,
    y,
    w,
    h,
    text="",
    font_size=18,
    color=TEXT,
    bold=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_lines(
    slide,
    x,
    y,
    w,
    h,
    lines,
    font_size=15,
    color=TEXT,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.line_spacing = 1.08
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def add_background(slide, accent, accent_dark):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_fill(bg, LIGHT)
    set_line(bg, LIGHT, 0)
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.24))
    set_fill(top, accent)
    set_line(top, accent, 0)
    for x, y, size, tr in [(11.35, 0.55, 0.62, 0.8), (12.0, 0.94, 0.34, 0.84), (10.8, 1.0, 0.26, 0.86)]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        set_fill(circle, accent_dark, tr)
        set_line(circle, accent_dark, 0.5, 0.9)


def add_footer(slide, text, no, total, accent):
    add_box(slide, 0.55, 7.12, 10.8, 0.18, text, 9, MUTED)
    add_box(slide, 12.0, 7.1, 0.9, 0.2, f"{no}/{total}", 10, accent, True, font_name="Aptos Display", align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body_lines, fill_color, accent, body_size=14):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill_color)
    set_line(shape, accent, 1.0, 0.22)
    add_box(slide, x + 0.14, y + 0.12, w - 0.28, 0.3, title, 17, accent, True, font_name="Aptos Display")
    add_lines(slide, x + 0.14, y + 0.48, w - 0.28, h - 0.58, body_lines, body_size, TEXT)


def add_title_slide(prs, payload, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, payload["accent"], payload["accent_dark"])
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.86), Inches(6.35), Inches(5.7))
    set_fill(panel, payload["accent"])
    set_line(panel, payload["accent"], 0.8)
    add_box(slide, 0.95, 1.0, 2.8, 0.28, payload["bhf"], 13, WHITE, True)
    add_box(slide, 0.95, 1.42, 5.55, 1.55, payload["title"], 25, WHITE, True, font_name="Aptos Display")
    add_lines(slide, 0.95, 3.2, 5.4, 1.4, [payload["subtitle"]], 15, rgb("E7F3FA"))
    add_card(
        slide,
        7.35,
        1.0,
        5.12,
        1.5,
        "Quelle",
        [f"• {payload['source_name']}", "• direkt aus dem uebergebenen DOCX extrahiert"],
        TONES["sky"],
        payload["accent_dark"],
    )
    add_card(
        slide,
        7.35,
        2.78,
        5.12,
        1.65,
        "Wie lernen?",
        [
            "erst Thema definieren",
            "dann Unterpunkte aus dem Dokument sortieren",
            "zum Schluss in eigenen Worten pruefungsreif sagen",
        ],
        TONES["sand"],
        payload["accent_dark"],
    )
    add_card(
        slide,
        7.35,
        4.73,
        5.12,
        1.05,
        "Ziel",
        ["ein Deck = ein klares Thema aus dem Dokument"],
        TONES["mint"],
        payload["accent_dark"],
    )
    add_footer(slide, "Thema-fuer-Thema aus dem uebergebenen DOCX erstellt", no, total, payload["accent"])


def add_content_slide(prs, payload, chunk_lines, slide_title, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, payload["accent"], payload["accent_dark"])
    add_box(slide, 0.72, 0.92, 10.5, 0.42, slide_title, 23, INK, True, font_name="Aptos Display")
    add_box(slide, 0.75, 1.28, 10.8, 0.24, "Diese Punkte kommen direkt aus dem uebergebenen Dokument.", 12, MUTED)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.76), Inches(12.0), Inches(4.98))
    set_fill(left, TONES["white"])
    set_line(left, payload["accent_dark"], 1.0, 0.22)
    add_lines(slide, 0.92, 2.0, 11.6, 4.45, [f"• {line}" for line in chunk_lines], 16, TEXT)
    add_footer(slide, "Lies diese Punkte laut als Kurzvortrag und notiere, was du noch nicht frei kannst.", no, total, payload["accent"])


def add_exam_slide(prs, payload, no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, payload["accent"], payload["accent_dark"])
    add_box(slide, 0.72, 0.92, 9.8, 0.42, "Pruefungsfokus", 23, INK, True, font_name="Aptos Display")
    add_card(slide, 0.72, 1.65, 5.85, 2.15, "So antwortest du", payload["focus_lines"], TONES["mint"], payload["accent_dark"], body_size=14)
    add_card(slide, 6.78, 1.65, 5.85, 2.15, "Das solltest du koennen", payload["question_lines"], TONES["sand"], payload["accent_dark"], body_size=14)
    add_card(slide, 0.72, 4.05, 11.91, 1.75, "Merksatz", [payload["memory_line"]], TONES["rose"], payload["accent_dark"], body_size=16)
    add_footer(slide, "Wenn du diese Slide frei sagen kannst, sitzt das Thema schon deutlich besser.", no, total, payload["accent"])


def build_presentation(payload: dict) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = payload["title"]
    prs.core_properties.subject = payload["subtitle"]
    prs.core_properties.author = "Codex"
    chunks = [payload["document_lines"][index : index + 6] for index in range(0, len(payload["document_lines"]), 6)] or [[]]
    total = 2 + len(chunks)
    slide_no = 1
    add_title_slide(prs, payload, slide_no, total)
    slide_no += 1
    for index, chunk in enumerate(chunks, start=1):
        slide_title = "Themenpunkte aus dem Dokument"
        if len(chunks) > 1:
            slide_title = f"{slide_title} ({index}/{len(chunks)})"
        add_content_slide(prs, payload, chunk, slide_title, slide_no, total)
        slide_no += 1
    add_exam_slide(prs, payload, slide_no, total)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / payload["filename"]
    prs.save(str(path))
    return path


def accent_for_index(index: int) -> tuple[RGBColor, RGBColor]:
    accent, accent_dark = ACCENTS[index % len(ACCENTS)]
    return rgb(accent), rgb(accent_dark)


def build_payload(doc: dict, theme: dict, theme_index: int, lines: list[str], deck_index: int) -> dict:
    extracted_lines = extract_theme_lines(lines, doc, theme_index)
    extracted_lines = [trim_line(line) for line in extracted_lines[:10]]
    kind = infer_kind(theme, extracted_lines)
    focus_lines = exam_focus(kind, theme["title"], extracted_lines)
    question_lines = learning_questions(kind, theme["title"], extracted_lines)
    accent, accent_dark = accent_for_index(deck_index)
    subtitle = f"{doc['bhf']} | {theme['title']}"
    memory_line = f"{theme['title']}: zuerst Thema einordnen, dann Unterpunkte aus dem Dokument logisch verbinden."
    filename = f"{doc['id']}_{theme_index + 1:02d}_{slugify(theme['title'])}.pptx"
    return {
        "bhf": doc["bhf"],
        "title": theme["title"],
        "subtitle": subtitle,
        "source_name": doc["path"].name,
        "document_lines": extracted_lines or [theme["title"]],
        "focus_lines": focus_lines,
        "question_lines": question_lines,
        "memory_line": trim_line(memory_line, 120),
        "accent": accent,
        "accent_dark": accent_dark,
        "filename": filename,
    }


def build_all():
    manifest: list[dict] = []
    deck_index = 0
    for doc in DOCS:
        raw_text = read_doc_text(doc)
        lines = expand_doc_lines(clean_lines(raw_text), doc)
        for theme_index, theme in enumerate(doc["themes"]):
            payload = build_payload(doc, theme, theme_index, lines, deck_index)
            path = build_presentation(payload)
            manifest.append(
                {
                    "bhf": doc["bhf"],
                    "source_doc": str(doc["path"]),
                    "theme": theme["title"],
                    "deck": str(path),
                    "points": payload["document_lines"],
                }
            )
            print(f"saved {path}")
            deck_index += 1
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    MANIFEST_PATH.write_text(json.dumps(manifest, indent=2, ensure_ascii=True), encoding="utf-8")
    print(f"saved {MANIFEST_PATH}")


if __name__ == "__main__":
    build_all()
