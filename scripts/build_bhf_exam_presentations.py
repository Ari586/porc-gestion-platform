from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path("/Users/arielhavana/antigr/porc/mockups/exam_prep")


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


WHITE = rgb("FFFFFF")
INK = rgb("223247")
TEXT = rgb("32465B")
MUTED = rgb("61758D")
LINE = rgb("D5E1E9")
SOFT_BLUE = rgb("E8F3F8")
SOFT_MINT = rgb("E9F5EF")
SOFT_SAND = rgb("FBF1DB")
SOFT_ROSE = rgb("F9E6E7")

CARD_TONES = {
    "sky": rgb("E7F3FA"),
    "mint": rgb("E6F4EE"),
    "sand": rgb("FCF2D8"),
    "rose": rgb("F9E5EA"),
    "white": WHITE,
}


DECKS = [
    {
        "filename": "01_BHF_1_4_Grundlagen_Pflege_und_Berufsrolle.pptx",
        "title": "BHF 1-4 | Grundlagen der Pflege",
        "subtitle": "Berufsrolle, Pflegeprozess, Hygiene, Kommunikation und Start in Schwangerschaft/Neugeborene",
        "bhf": "BHF 1a-1b-2-3a-3b-4",
        "accent": rgb("1D7A8C"),
        "accent_dark": rgb("175E6C"),
        "hero": "Fundament für die Prüfung: sauber definieren, beobachten, priorisieren und verständlich begründen.",
        "slides": [
            {
                "type": "cards",
                "title": "Was du sicher können musst",
                "subtitle": "Diese vier Blöcke tauchen in fast jeder Pflegeprüfung wieder auf.",
                "cards": [
                    {
                        "title": "Berufsrolle",
                        "tone": "sky",
                        "body": [
                            "Pflege ist geplant, beobachtend und begründet.",
                            "Schweigepflicht, Würde, Selbstbestimmung und Sicherheit sind Kernprinzipien.",
                            "Rollenklärung verhindert Grenzverletzungen.",
                        ],
                    },
                    {
                        "title": "Pflegeprozess",
                        "tone": "mint",
                        "body": [
                            "Informationssammlung",
                            "Probleme/Ressourcen",
                            "Ziele, Maßnahmen, Durchführung, Evaluation",
                        ],
                    },
                    {
                        "title": "Beobachtung",
                        "tone": "sand",
                        "body": [
                            "Wahrnehmen ist nicht gleich interpretieren.",
                            "Beschreibe objektiv, bevor du bewertest.",
                            "Abweichungen sofort einordnen und weitergeben.",
                        ],
                    },
                    {
                        "title": "Hygiene und Grundpflege",
                        "tone": "rose",
                        "body": [
                            "Händehygiene, Intimschutz, Mobilisation, Hautbeobachtung",
                            "Ernährung, Ausscheidung und Prophylaxen immer mitdenken",
                        ],
                    },
                ],
            },
            {
                "type": "cards",
                "title": "Pflegeprozess klar und prüfungssicher",
                "subtitle": "Formuliere immer von Ressource + Problem + Ziel + Maßnahme + Evaluation aus.",
                "cards": [
                    {
                        "title": "1. Assessment",
                        "tone": "sky",
                        "body": [
                            "Biografie, Vitalwerte, Schmerzen, Mobilität, Ernährung, Ausscheidung",
                            "Subjektive + objektive Daten trennen",
                        ],
                    },
                    {
                        "title": "2. Pflegeproblem",
                        "tone": "mint",
                        "body": [
                            "Aktuell oder potenziell?",
                            "Beispiel: Sturzrisiko bei Schwindel und unsicherem Gang",
                        ],
                    },
                    {
                        "title": "3. Pflegeziel",
                        "tone": "sand",
                        "body": [
                            "SMART formulieren",
                            "Kurzfristig: heute sicher mobilisieren",
                            "Langfristig: Selbstständigkeit steigern",
                        ],
                    },
                    {
                        "title": "4. Evaluation",
                        "tone": "rose",
                        "body": [
                            "Hat die Maßnahme gewirkt?",
                            "Wenn nein: Ursache neu prüfen und Plan anpassen",
                        ],
                    },
                ],
            },
            {
                "type": "cards",
                "title": "Kommunikation, Team und Recht",
                "subtitle": "Hier werden oft Definition + Beispiel + Konsequenz verlangt.",
                "cards": [
                    {
                        "title": "Empathie",
                        "tone": "mint",
                        "body": [
                            "Gefühl wahrnehmen, ohne Grenzen zu verlieren",
                            "Beispiel: Angst benennen und strukturierend begleiten",
                        ],
                    },
                    {
                        "title": "Kollegiale Zusammenarbeit",
                        "tone": "sky",
                        "body": [
                            "klar, wertschätzend, konkret",
                            "Feedback beschreibt Verhalten, nicht Persönlichkeit",
                        ],
                    },
                    {
                        "title": "Rechte und Pflichten",
                        "tone": "sand",
                        "body": [
                            "Aufklärung, Einwilligung, Datenschutz, Dokumentation",
                            "Unterlassene Beobachtung kann haftungsrelevant werden",
                        ],
                    },
                    {
                        "title": "Kultursensibilität",
                        "tone": "rose",
                        "body": [
                            "Nicht stereotype Annahmen treffen",
                            "Fragen stellen, Ressourcen erfragen, Würde schützen",
                        ],
                    },
                ],
            },
            {
                "type": "cards",
                "title": "Grundpflege mit Blick auf Risiken",
                "subtitle": "Körperpflege ist nie nur Waschen, sondern immer Assessment in Aktion.",
                "cards": [
                    {
                        "title": "Haut und Dekubitus",
                        "tone": "sky",
                        "body": [
                            "Rötung, Wärme, Druckstellen, trockene Haut erfassen",
                            "Lagerung, Mobilisierung, Hautschutz begründen",
                        ],
                    },
                    {
                        "title": "Ernährung",
                        "tone": "mint",
                        "body": [
                            "Essmenge, Trinkmenge, Kau-/Schluckprobleme beobachten",
                            "Mangelernährung und Dehydratation früh erkennen",
                        ],
                    },
                    {
                        "title": "Ausscheidung",
                        "tone": "sand",
                        "body": [
                            "Menge, Farbe, Geruch, Schmerzen, Obstipation, Diarrhö",
                            "Scham berücksichtigen, Intimsphäre sichern",
                        ],
                    },
                    {
                        "title": "Mobilität",
                        "tone": "rose",
                        "body": [
                            "Transfer, Sturzrisiko, Hilfsmittel, Schmerzbeobachtung",
                            "Sicherheit vor Tempo",
                        ],
                    },
                ],
            },
            {
                "type": "cards",
                "title": "Schwangerschaft, Geburt und Neugeborene",
                "subtitle": "Hier zählt besonders das sichere Benennen normaler vs. kritischer Zeichen.",
                "cards": [
                    {
                        "title": "Schwangerschaft",
                        "tone": "sky",
                        "body": [
                            "normale Veränderungen vs. Warnzeichen",
                            "Red Flags: Blutung, starke Schmerzen, Ödeme + RR-Anstieg",
                        ],
                    },
                    {
                        "title": "Geburt",
                        "tone": "mint",
                        "body": [
                            "Beobachtung von Wehen, Mutter und Kind",
                            "ruhig strukturieren, Sicherheit vermitteln",
                        ],
                    },
                    {
                        "title": "Neugeborenes",
                        "tone": "sand",
                        "body": [
                            "Temperatur, Atmung, Hautfarbe, Trinkverhalten",
                            "Bindung und Elternanleitung mitdenken",
                        ],
                    },
                    {
                        "title": "Prüfungssatz",
                        "tone": "rose",
                        "body": [
                            "Definition + Beispiel + pflegerische Konsequenz",
                            "So antwortest du fast immer vollständig",
                        ],
                    },
                ],
            },
            {
                "type": "checklist",
                "title": "Prüfungscheck BHF 1-4",
                "left_title": "Das solltest du frei erklären können",
                "left_lines": [
                    "Pflegeprozess in 6 Schritten",
                    "Unterschied Beobachtung vs. Interpretation",
                    "wichtige Rechte und Pflichten in der Pflege",
                    "Grundpflege als Risiko-Assessment",
                    "Red Flags in Schwangerschaft und beim Neugeborenen",
                ],
                "right_title": "Merkhilfen für die Klausur",
                "right_lines": [
                    "Erst objektiv beschreiben, dann bewerten",
                    "Immer Problem + Ziel + Maßnahme + Evaluation nennen",
                    "Bei jeder Situation Sicherheit, Würde und Kommunikation mitdenken",
                    "Bei Unsicherheit: Was ist akut? Was ist gefährlich? Was muss ich weitergeben?",
                ],
            },
        ],
    },
    {
        "filename": "02_BHF_5_6_Akutpflege_Blutzirkulation_und_Notfaelle.pptx",
        "title": "BHF 5-6 | Akutpflege, Blut und Notfälle",
        "subtitle": "Wunden, Kreislauf, Dehydratation, Notfallsituationen, Blut und Intensivpflege",
        "bhf": "BHF 05-06",
        "accent": rgb("C66A4C"),
        "accent_dark": rgb("9B4E35"),
        "hero": "Prüfungsziel: in Akutsituationen schnell priorisieren, Symptome erkennen und pflegerisch sicher handeln.",
        "slides": [
            {
                "type": "cards",
                "title": "Prüfungsfokus im Akutblock",
                "subtitle": "Die Frage lautet fast immer: Was ist lebenswichtig, was ist sofort zu tun?",
                "cards": [
                    {
                        "title": "ABCDE",
                        "tone": "sky",
                        "body": [
                            "Airway, Breathing, Circulation, Disability, Exposure",
                            "Lebensbedrohung zuerst beheben",
                        ],
                    },
                    {
                        "title": "Vitalzeichen",
                        "tone": "mint",
                        "body": [
                            "RR, Puls, AF, SpO2, Temperatur, Bewusstsein",
                            "Trends sind wichtiger als Einzelwerte",
                        ],
                    },
                    {
                        "title": "Kreislauf und Volumen",
                        "tone": "sand",
                        "body": [
                            "Dehydratation, Blutverlust, Schock früh erkennen",
                            "Bilanz und Haut-/Schleimhautbeobachtung zählen",
                        ],
                    },
                    {
                        "title": "Sicherheit",
                        "tone": "rose",
                        "body": [
                            "Hilfe holen, Monitoring sichern, Ruhe ausstrahlen",
                            "Dokumentation und Übergabe klar halten",
                        ],
                    },
                ],
            },
            {
                "type": "durst",
                "title": "DURST-Schema | Dehydratation",
                "term": "Dehydratation",
                "hook": "Typisch in Prüfung und Praxis: Flüssigkeitsmangel erkennen, einordnen und rechtzeitig gegensteuern.",
                "cards": [
                    ("Definition", ["Mangel an Körperwasser, oft auch mit Elektrolytstörung", "führt zu Kreislauf-, Nieren- und Bewusstseinsproblemen"]),
                    ("Formen / Typen", ["isoton", "hyperton", "hypoton", "leicht, mittel, schwer je nach Klinik"]),
                    ("Ursachen", ["zu geringe Trinkmenge", "Fieber, Erbrechen, Diarrhö", "Diuretika, Schluckstörung"]),
                    ("Risikofaktoren", ["Säuglinge und alte Menschen", "Demenz, Immobilität, Dysphagie", "Hitze, Infekt, Pflegeabhängigkeit"]),
                    ("Leitsymptome", ["Durst, trockene Schleimhäute", "Oligurie, dunkler Urin", "Schwindel, Delir, Tachykardie"]),
                    ("Therapie / Pflege", ["Ursache behandeln", "oral oder i.v. Flüssigkeit nach Standard", "Bilanz, Gewicht, Hautturgor, Mundpflege, Sturzprophylaxe"]),
                ],
                "exam": "Prüfungspunkt: Schockzeichen, Verwirrtheit und geringe Urinmenge immer als Warnsignale nennen.",
            },
            {
                "type": "durst",
                "title": "DURST-Schema | Anaphylaxie",
                "term": "Anaphylaxie",
                "hook": "Merke: plötzlich, systemisch, potenziell tödlich. Handlungssicherheit zählt mehr als Spezialwissen.",
                "cards": [
                    ("Definition", ["akute schwere allergische Sofortreaktion", "betrifft meist Haut, Atmung und Kreislauf gleichzeitig"]),
                    ("Grade / Schwere", ["Grad I: Hautsymptome", "Grad II-III: Atmung/Kreislauf", "Grad IV: Reanimation"]),
                    ("Ursachen", ["Nahrungsmittel", "Medikamente", "Insektengifte", "seltener andere Trigger"]),
                    ("Risikofaktoren", ["bekannte Allergien", "Asthma", "frühere Anaphylaxie", "schneller Kontakt mit Allergen"]),
                    ("Leitsymptome", ["Urtikaria, Juckreiz", "Atemnot, Stridor, Schwellung", "Hypotonie, Tachykardie, Angst"]),
                    ("Therapie / Pflege", ["Notruf, Lagerung, O2", "Adrenalin i.m. nach Standard", "Monitoring, venöser Zugang, beruhigen, Auslöser stoppen"]),
                ],
                "exam": "Prüfungspunkt: bei Atemnot + Kreislaufzeichen immer an Anaphylaxie denken und sofort eskalieren.",
            },
            {
                "type": "cards",
                "title": "Blut, Gerinnung und Intensivlogik",
                "subtitle": "Einfach erklärt, aber pflegerisch sehr relevant.",
                "cards": [
                    {
                        "title": "Erythrozyten",
                        "tone": "sky",
                        "body": [
                            "Sauerstofftransport",
                            "bei Mangel: Blässe, Müdigkeit, Belastungsdyspnoe",
                        ],
                    },
                    {
                        "title": "Leukozyten",
                        "tone": "mint",
                        "body": [
                            "Abwehr von Erregern",
                            "bei Infektion oft erhöht, bei Immunschwäche Risiko beachten",
                        ],
                    },
                    {
                        "title": "Thrombozyten",
                        "tone": "sand",
                        "body": [
                            "wichtig für Gerinnung",
                            "bei Mangel: Blutungsneigung, Petechien",
                        ],
                    },
                    {
                        "title": "ITS-Denken",
                        "tone": "rose",
                        "body": [
                            "engmaschig beobachten",
                            "kleine Änderungen können große Bedeutung haben",
                        ],
                    },
                ],
            },
            {
                "type": "checklist",
                "title": "Prüfungscheck BHF 5-6",
                "left_title": "Das solltest du frei erklären können",
                "left_lines": [
                    "ABCDE-Schema mit Beispiel",
                    "Dehydratation sicher erkennen",
                    "Anaphylaxie als Notfall benennen",
                    "Rolle von Erythrozyten, Leukozyten, Thrombozyten",
                    "wichtige Beobachtungen bei akuten Zustandsänderungen",
                ],
                "right_title": "Merkhilfen für die Klausur",
                "right_lines": [
                    "Akut geht vor vollständig",
                    "Immer zuerst Airway, Atmung, Kreislauf prüfen",
                    "Bei Volumenmangel: Mund, Haut, Urin, Puls, RR, Bewusstsein",
                    "Bei Allergie + Atemnot + Hypotonie: sofort Anaphylaxie denken",
                ],
            },
        ],
    },
    {
        "filename": "03_BHF_7_12_Sinne_Neurologie_Langzeitpflege_und_Paediatrie.pptx",
        "title": "BHF 7-12 | Sinne, Neurologie, Langzeitpflege und Pädiatrie",
        "subtitle": "Sinnesbeeinträchtigung, Fieber, Neurologie, Langzeitpflege, Diabetes, ADHS und Medikamente",
        "bhf": "BHF 07-12",
        "accent": rgb("3F7E5E"),
        "accent_dark": rgb("2E5D44"),
        "hero": "Hier wird oft Transfer geprüft: Symptome beobachten, Alltag anpassen, Angehörige anleiten.",
        "slides": [
            {
                "type": "cards",
                "title": "Prüfungsfokus im Mittelblock",
                "subtitle": "Viele Themen sind alltagsnah. Saubere Pflegebegründung bringt hier viele Punkte.",
                "cards": [
                    {
                        "title": "Sehen / Hören",
                        "tone": "sky",
                        "body": [
                            "Orientierung, Sicherheit und Kommunikation anpassen",
                            "Umgebung strukturieren, Hilfsmittel einbeziehen",
                        ],
                    },
                    {
                        "title": "Fieber",
                        "tone": "mint",
                        "body": [
                            "Beobachtung: Temperatur, Kreislauf, Flüssigkeit, Allgemeinzustand",
                            "Fieber ist Symptom, nicht die Diagnose",
                        ],
                    },
                    {
                        "title": "Langzeitpflege",
                        "tone": "sand",
                        "body": [
                            "Ressourcen fördern, Routinen nutzen, Risiken erkennen",
                            "Pflegeprozess und Biografiearbeit sind zentral",
                        ],
                    },
                    {
                        "title": "Pädiatrie",
                        "tone": "rose",
                        "body": [
                            "Kind + Eltern gleichzeitig im Blick",
                            "Erklärung altersgerecht und ruhig",
                        ],
                    },
                ],
            },
            {
                "type": "durst",
                "title": "DURST-Schema | Schlaganfall",
                "term": "Schlaganfall",
                "hook": "Merke: Zeit ist Gehirn. In Prüfungen sind Symptome, FAST und Pflegeprioritäten entscheidend.",
                "cards": [
                    ("Definition", ["akute Durchblutungsstörung oder Blutung im Gehirn", "führt zu plötzlichen neurologischen Ausfällen"]),
                    ("Formen / Typen", ["ischämisch", "hämorrhagisch", "TIA als Warnereignis ohne bleibenden Infarkt"]),
                    ("Ursachen", ["Thrombus, Embolus", "Gefäßruptur", "Vorhofflimmern, Hypertonie als wichtige Hintergründe"]),
                    ("Risikofaktoren", ["Hypertonie, Rauchen, Diabetes", "Adipositas, Bewegungsmangel", "Vorhofflimmern, Dyslipidämie"]),
                    ("Leitsymptome", ["FAST: Face, Arm, Speech, Time", "Hemiparese, Sprachstörung, Sehstörung", "Schluckstörung, Bewusstseinsveränderung"]),
                    ("Therapie / Pflege", ["Stroke Unit, ggf. Lyse/Thrombektomie", "Aspirationsprophylaxe, Lagerung, Monitoring", "Frühmobilisation und Reha einleiten"]),
                ],
                "exam": "Prüfungspunkt: FAST nennen, Aspirationsgefahr mitdenken und jede plötzliche Veränderung als Notfall einordnen.",
            },
            {
                "type": "durst",
                "title": "DURST-Schema | Diabetes mellitus Typ 1",
                "term": "Diabetes Typ 1",
                "hook": "Bei Kindern besonders wichtig: Symptome, Hypo-/Hyperglykämie und Elternanleitung sicher erklären.",
                "cards": [
                    ("Definition", ["Autoimmunerkrankung mit absolutem Insulinmangel", "beginnt oft im Kindes- oder Jugendalter"]),
                    ("Formen / Typen", ["Typ 1 vs. Typ 2 klar unterscheiden", "akute Entgleisung: Ketoazidose"]),
                    ("Ursachen", ["Autoimmunreaktion gegen Betazellen", "genetische und Umweltfaktoren spielen mit"]),
                    ("Risikofaktoren", ["familiäre Belastung", "Autoimmunneigung", "bei Typ 1 weniger lifestyle-gebunden"]),
                    ("Leitsymptome", ["Polydipsie, Polyurie, Gewichtsverlust", "Müdigkeit, Bauchschmerzen", "bei Ketoazidose: Übelkeit, Azetongeruch, tiefe Atmung"]),
                    ("Therapie / Pflege", ["Insulin, BZ-Kontrolle, Ernährung, Bewegung", "Hypo-/Hyperglykämie erkennen", "Kind und Eltern schulen, Spritzen/CGM erklären"]),
                ],
                "exam": "Prüfungspunkt: Hypoglykämie immer als akut gefährlich benennen und schnell handeln.",
            },
            {
                "type": "cards",
                "title": "ADHS, Medikamente und Pflegealltag",
                "subtitle": "Nicht nur Krankheit nennen, sondern immer Alltag und Beziehungsgestaltung mitdenken.",
                "cards": [
                    {
                        "title": "ADHS kurz erklärt",
                        "tone": "sky",
                        "body": [
                            "Aufmerksamkeit, Impulsivität, Hyperaktivität",
                            "Beeinträchtigt Schule, Beziehungen und Selbststeuerung",
                        ],
                    },
                    {
                        "title": "Pflege / Umgang",
                        "tone": "mint",
                        "body": [
                            "kurze Anweisungen, klare Struktur, wenig Reizüberflutung",
                            "Lob konkret und zeitnah geben",
                        ],
                    },
                    {
                        "title": "Medikamente",
                        "tone": "sand",
                        "body": [
                            "6-R-Regel, Wirkung und Nebenwirkung beobachten",
                            "Eltern informieren, Einnahmefehler vermeiden",
                        ],
                    },
                    {
                        "title": "Langzeitpflege",
                        "tone": "rose",
                        "body": [
                            "Pflege im Heim ist beziehungs- und ressourcenorientiert",
                            "Orientierung, Tagesstruktur und Sicherheit sind Schlüssel",
                        ],
                    },
                ],
            },
            {
                "type": "checklist",
                "title": "Prüfungscheck BHF 7-12",
                "left_title": "Das solltest du frei erklären können",
                "left_lines": [
                    "Pflege bei Seh- und Höreinschränkung",
                    "Fieber als Symptom richtig beobachten",
                    "Schlaganfall inkl. FAST und Aspirationsrisiko",
                    "Diabetes Typ 1 inkl. Hypo-/Hyperglykämie",
                    "ADHS + Strukturierung im Alltag",
                ],
                "right_title": "Merkhilfen für die Klausur",
                "right_lines": [
                    "Bei Sinnesverlust: Orientierung und Kommunikation zuerst",
                    "Bei Neurologie: Zeitfaktor und Schluckstörung nennen",
                    "Bei Kindern: Elternarbeit ist Teil der Pflege",
                    "Bei Medikamenten: Wirkung, Nebenwirkung, Sicherheit",
                ],
            },
        ],
    },
    {
        "filename": "04_BHF_13_15_Chronische_Erkrankungen_Expertenstandard_und_Gesundheitsfoerderung.pptx",
        "title": "BHF 13-15 | Chronische Erkrankungen, Expertenstandard und Gesundheitsförderung",
        "subtitle": "Parkinson, Dialyse/NI, Expertenstandards, Berufskrankheiten, Burnout und Prävention",
        "bhf": "BHF 13-15",
        "accent": rgb("9B7A2E"),
        "accent_dark": rgb("745C22"),
        "hero": "Prüfungsziel: chronische Verläufe systematisch erfassen und Pflege als langfristige Begleitung beschreiben.",
        "slides": [
            {
                "type": "cards",
                "title": "Prüfungsfokus im Chronik-Block",
                "subtitle": "Hier zählt Verstehen statt Auswendiglernen: Verlauf, Ressourcen, Beratung, Prävention.",
                "cards": [
                    {
                        "title": "Chronisch krank",
                        "tone": "sky",
                        "body": [
                            "langfristiger Verlauf mit Anpassungsbedarf",
                            "Autonomie, Edukation und Alltag stehen im Vordergrund",
                        ],
                    },
                    {
                        "title": "Expertenstandard",
                        "tone": "mint",
                        "body": [
                            "evidenzbasiertes Vorgehen",
                            "hilft, Pflege zu begründen und Qualität zu sichern",
                        ],
                    },
                    {
                        "title": "Gesundheitsförderung",
                        "tone": "sand",
                        "body": [
                            "Risikofaktoren senken, Ressourcen stärken",
                            "Beratung muss alltagstauglich sein",
                        ],
                    },
                    {
                        "title": "Arbeitsgesundheit",
                        "tone": "rose",
                        "body": [
                            "Burnout, Mobbing, Berufskrankheiten kennen",
                            "Selbstschutz ist professionell, nicht egoistisch",
                        ],
                    },
                ],
            },
            {
                "type": "durst",
                "title": "DURST-Schema | Parkinson-Syndrom",
                "term": "Parkinson",
                "hook": "Typische Prüfungsfrage: Leitsymptome + Medikamentenbesonderheiten + pflegerische Unterstützung im Alltag.",
                "cards": [
                    ("Definition", ["neurodegeneratives Syndrom mit Dopaminmangel", "führt zu motorischen und nicht-motorischen Symptomen"]),
                    ("Formen / Typen", ["idiopathisch", "atypische Parkinson-Syndrome", "sekundäre Formen"]),
                    ("Ursachen", ["Untergang dopaminerger Nervenzellen", "genaue Ursache oft multifaktoriell"]),
                    ("Risikofaktoren", ["höheres Alter", "familiäre Faktoren", "bestimmte toxische Einflüsse möglich"]),
                    ("Leitsymptome", ["Bradykinese", "Rigor", "Tremor", "posturale Instabilität", "Hypomimie, kleinschrittiger Gang"]),
                    ("Therapie / Pflege", ["L-Dopa und weitere dopaminerge Therapie", "Sturzprophylaxe, Bewegungsförderung", "genügend Zeit geben, Schlucken und Obstipation beachten"]),
                ],
                "exam": "Prüfungspunkt: Bradykinese ist Leitsymptom, nicht nur Tremor. Medikamentenzeitpunkte sind wichtig.",
            },
            {
                "type": "durst",
                "title": "DURST-Schema | Chronische Niereninsuffizienz",
                "term": "Chronische Niereninsuffizienz",
                "hook": "Wichtig für Prüfung und Praxis: Stadien, Flüssigkeit, Ernährung und Dialysebezug.",
                "cards": [
                    ("Definition", ["langsam fortschreitender Funktionsverlust der Nieren", "führt zu Retention von harnpflichtigen Substanzen"]),
                    ("Grade / Stadien", ["G1 bis G5 nach GFR", "Endstadium oft dialysepflichtig"]),
                    ("Ursachen", ["Diabetes mellitus", "Hypertonie", "Glomerulopathien", "langjährige Nierenschäden"]),
                    ("Risikofaktoren", ["Diabetes, Hypertonie, Alter", "Nephrotoxine, kardiovaskuläre Erkrankungen"]),
                    ("Leitsymptome", ["Müdigkeit, Ödeme, Juckreiz", "Anämie, Übelkeit, Hypertonie", "später Urämiesymptome"]),
                    ("Therapie / Pflege", ["RR- und Diabeteskontrolle", "Ernährungs- und Trinkanpassung", "Gewicht, Bilanz, Ödeme, Haut, Dialysezugang beobachten"]),
                ],
                "exam": "Prüfungspunkt: Bilanzierung, Elektrolyte und Dialysebeobachtung immer mit nennen.",
            },
            {
                "type": "cards",
                "title": "Expertenstandard, Medikation und Burnout",
                "subtitle": "In Prüfungen oft als Transferfrage formuliert.",
                "cards": [
                    {
                        "title": "Expertenstandard",
                        "tone": "sky",
                        "body": [
                            "z. B. Schmerz, Dekubitus, Ernährung, Sturz",
                            "liefert Struktur für Assessment und Maßnahmen",
                        ],
                    },
                    {
                        "title": "Medikationssicherheit",
                        "tone": "mint",
                        "body": [
                            "6-R-Regel, Beobachtung von Wirkung/Nebenwirkung",
                            "Polypharmazie bei chronischer Krankheit kritisch prüfen",
                        ],
                    },
                    {
                        "title": "Burnout",
                        "tone": "sand",
                        "body": [
                            "Erschöpfung, Distanzierung, Leistungsabfall",
                            "Frühzeichen ernst nehmen, Ressourcen und Pausen stärken",
                        ],
                    },
                    {
                        "title": "Gesundheitsförderung",
                        "tone": "rose",
                        "body": [
                            "realistische Ziele vereinbaren",
                            "nicht nur informieren, sondern Verhalten anbahnen",
                        ],
                    },
                ],
            },
            {
                "type": "checklist",
                "title": "Prüfungscheck BHF 13-15",
                "left_title": "Das solltest du frei erklären können",
                "left_lines": [
                    "Merkmale chronischer Krankheit",
                    "Parkinson mit typischen Symptomen",
                    "Chronische Niereninsuffizienz + Dialysebezug",
                    "Wozu Expertenstandards dienen",
                    "Burnout von normaler Belastung abgrenzen",
                ],
                "right_title": "Merkhilfen für die Klausur",
                "right_lines": [
                    "Bei Chronik immer Beratung + Ressourcen + Adhärenz nennen",
                    "Bei Parkinson Medikamentenzeitpunkt und Sturzrisiko",
                    "Bei Niereninsuffizienz Bilanz und Ödeme nicht vergessen",
                    "Bei Burnout: Prävention ist Team- und Selbstmanagement",
                ],
            },
        ],
    },
    {
        "filename": "05_BHF_16_19_Psychiatrie_Infektion_Trauma_und_Rehabilitation.pptx",
        "title": "BHF 16-19 | Psychiatrie, Infektion, Trauma und Rehabilitation",
        "subtitle": "Psychiatrische Pflege, Infektionsschutz, Polytrauma, Schlaganfall-Reha und Aphasie",
        "bhf": "BHF 16-19",
        "accent": rgb("7C5C9B"),
        "accent_dark": rgb("604575"),
        "hero": "Prüfungsziel: Sicherheit, Beziehungsgestaltung und funktionelle Ziele miteinander verbinden.",
        "slides": [
            {
                "type": "cards",
                "title": "Prüfungsfokus im 3. Ausbildungsjahr",
                "subtitle": "Hier werden häufig komplexe Situationen mit mehreren Risiken gleichzeitig geprüft.",
                "cards": [
                    {
                        "title": "Psychiatrie",
                        "tone": "sky",
                        "body": [
                            "Beziehung, Sicherheit und Struktur sind zentral",
                            "ruhig, klar, nicht konfrontativ kommunizieren",
                        ],
                    },
                    {
                        "title": "Infektion",
                        "tone": "mint",
                        "body": [
                            "Übertragungswege kennen",
                            "Standard- und Zusatzmaßnahmen korrekt begründen",
                        ],
                    },
                    {
                        "title": "Trauma / ITS",
                        "tone": "sand",
                        "body": [
                            "Prioritäten setzen, Monitoring verstehen",
                            "Frühkomplikationen erkennen",
                        ],
                    },
                    {
                        "title": "Rehabilitation",
                        "tone": "rose",
                        "body": [
                            "Funktion erhalten oder wieder aufbauen",
                            "kleine Teilziele und interprofessionelle Zusammenarbeit",
                        ],
                    },
                ],
            },
            {
                "type": "durst",
                "title": "DURST-Schema | Depression",
                "term": "Depression",
                "hook": "In Prüfungen wichtig: Symptome ernst nehmen, Suizidalität aktiv mitdenken und Beziehung professionell gestalten.",
                "cards": [
                    ("Definition", ["affektive Störung mit gedrückter Stimmung, Interessenverlust und Antriebsmangel", "betrifft Denken, Fühlen und Handeln"]),
                    ("Schweregrade / Formen", ["leicht, mittel, schwer", "episodisch oder rezidivierend", "mit/ohne psychotische Symptome"]),
                    ("Ursachen", ["multifaktoriell: biologisch, psychisch, sozial", "Belastungen und Vorerkrankungen spielen mit"]),
                    ("Risikofaktoren", ["frühere Episoden", "chronische Krankheit", "soziale Isolation", "Belastungen und Verluste"]),
                    ("Leitsymptome", ["Traurigkeit, Interessenverlust, Erschöpfung", "Schlafstörung, Schuldgefühle", "Suizidgedanken als Warnsignal"]),
                    ("Therapie / Pflege", ["Psychotherapie, Antidepressiva, Aktivierung", "wertschätzend, klar, nicht bagatellisieren", "Suizidalität ansprechen und Sicherheit herstellen"]),
                ],
                "exam": "Prüfungspunkt: Suizidalität direkt erfragen ist professionell und wichtig.",
            },
            {
                "type": "durst",
                "title": "DURST-Schema | Anorexia nervosa",
                "term": "Anorexia nervosa",
                "hook": "Wichtig: nicht nur Gewicht, sondern Körperbild, Kontrolle und somatische Risiken verstehen.",
                "cards": [
                    ("Definition", ["Essstörung mit absichtlich herbeigeführtem Untergewicht", "starke Angst vor Gewichtszunahme und Körperschemastörung"]),
                    ("Formen / Typen", ["restriktiver Typ", "Binge-Eating/Purging-Typ"]),
                    ("Ursachen", ["multifaktoriell: psychisch, familiär, sozial", "Leistungsdruck und Kontrollthemen können relevant sein"]),
                    ("Risikofaktoren", ["Jugendalter", "Perfektionismus", "geringes Selbstwertgefühl", "sozialer Druck und Ideale"]),
                    ("Leitsymptome", ["Untergewicht, Essvermeidung", "Amenorrhö möglich", "Bradykardie, Hypotonie, Frieren, Elektrolytstörung"]),
                    ("Therapie / Pflege", ["multimodal: Psychotherapie, Ernährungstherapie, medizinische Überwachung", "Gewicht, Kreislauf und Essverhalten beobachten", "ruhig, konsequent, nicht moralisieren"]),
                ],
                "exam": "Prüfungspunkt: somatische Gefahren wie Bradykardie und Elektrolytstörung immer mit nennen.",
            },
            {
                "type": "cards",
                "title": "Infektion, Trauma und Reha praxisnah",
                "subtitle": "Die Stärke in der Prüfung ist die Verknüpfung von Beobachtung, Schutz und Förderung.",
                "cards": [
                    {
                        "title": "Infektionsschutz",
                        "tone": "sky",
                        "body": [
                            "Händehygiene, Schutzkleidung, Isolationsform begründen",
                            "Patienten und Angehörige verständlich anleiten",
                        ],
                    },
                    {
                        "title": "Polytrauma",
                        "tone": "mint",
                        "body": [
                            "Atemweg, Atmung, Kreislauf zuerst",
                            "Blutung, Schmerz, Hypothermie und Schock im Blick",
                        ],
                    },
                    {
                        "title": "Aphasie",
                        "tone": "sand",
                        "body": [
                            "Sprache ist gestört, Intelligenz nicht automatisch",
                            "kurze Sätze, Zeit lassen, Hilfsmittel nutzen",
                        ],
                    },
                    {
                        "title": "Rehabilitation",
                        "tone": "rose",
                        "body": [
                            "Ziele alltagsnah formulieren",
                            "kleine Fortschritte sichtbar machen und loben",
                        ],
                    },
                ],
            },
            {
                "type": "checklist",
                "title": "Prüfungscheck BHF 16-19",
                "left_title": "Das solltest du frei erklären können",
                "left_lines": [
                    "Grundprinzipien psychiatrischer Pflege",
                    "Depression inkl. Suizidalität",
                    "Anorexie mit somatischen Risiken",
                    "Aphasie in der Kommunikation",
                    "Infektionsschutz und Prioritäten bei Trauma",
                ],
                "right_title": "Merkhilfen für die Klausur",
                "right_lines": [
                    "In Psychiatrie zuerst Sicherheit + Beziehung",
                    "Bei Infektion Übertragungsweg + Maßnahme nennen",
                    "Bei Trauma immer ABCDE denken",
                    "Bei Reha: realistische Teilziele und Motivation benennen",
                ],
            },
        ],
    },
    {
        "filename": "06_BHF_20_24_Qualitaet_Paediatrie_Konflikt_und_Falltraining.pptx",
        "title": "BHF 20-24 | Qualität, Pädiatrie, Konflikt und Falltraining",
        "subtitle": "Qualitätssicherung, Kontinenz, Frühgeborene, Gewaltprävention, SIS und Falldenken",
        "bhf": "BHF 20-24",
        "accent": rgb("A0475E"),
        "accent_dark": rgb("7D3648"),
        "hero": "Prüfungsziel: Fälle strukturiert lösen, Konflikte professionell deeskalieren und Pflegeplanung sauber aufbauen.",
        "slides": [
            {
                "type": "cards",
                "title": "Prüfungsfokus im Abschlussblock",
                "subtitle": "Hier zeigt sich, ob du Wissen in echte Pflegesituationen übersetzen kannst.",
                "cards": [
                    {
                        "title": "Qualität",
                        "tone": "sky",
                        "body": [
                            "Standards, Dokumentation, Evaluation",
                            "Fehlerprävention ist Teil guter Pflege",
                        ],
                    },
                    {
                        "title": "Pädiatrie / Neonatologie",
                        "tone": "mint",
                        "body": [
                            "Frühgeborene, Schreibaby, Darmerkrankungen",
                            "Familienzentrierung ist Pflicht",
                        ],
                    },
                    {
                        "title": "Konflikt / Gewalt",
                        "tone": "sand",
                        "body": [
                            "Gewalt erkennen, benennen, vorbeugen",
                            "Deeskalation braucht Haltung und Technik",
                        ],
                    },
                    {
                        "title": "Falltraining",
                        "tone": "rose",
                        "body": [
                            "BHF 24 ist pure Prüfungspraxis: Fälle, SIS, Prioritäten",
                            "Struktur schlägt Detailchaos",
                        ],
                    },
                ],
            },
            {
                "type": "durst",
                "title": "DURST-Schema | Harninkontinenz",
                "term": "Harninkontinenz",
                "hook": "Typisch prüfungsrelevant: Formen unterscheiden und passende pflegerische Maßnahmen ableiten.",
                "cards": [
                    ("Definition", ["unwillkürlicher Urinverlust", "beeinträchtigt Lebensqualität und Hautgesundheit"]),
                    ("Formen / Typen", ["Belastungsinkontinenz", "Dranginkontinenz", "Überlaufinkontinenz", "reflektorische Formen"]),
                    ("Ursachen", ["Beckenbodenschwäche", "neurologische Störung", "Obstruktion, Medikamente", "Infekt oder funktionelle Faktoren"]),
                    ("Risikofaktoren", ["Alter, Geburten, Immobilität", "Demenz, Schlaganfall, Diabetes", "Adipositas"]),
                    ("Leitsymptome", ["Urinverlust, Harndrang", "nächtliches Einnässen, Restharn", "Hautprobleme, Scham, Rückzug"]),
                    ("Therapie / Pflege", ["Beckenbodentraining, Toilettentraining, Ursachenbehandlung", "Hautschutz und Trinkmanagement", "Miktion beobachten, Inkontinenzmaterial passend wählen"]),
                ],
                "exam": "Prüfungspunkt: Form benennen und Maßnahme daran anpassen, nicht alles gleich behandeln.",
            },
            {
                "type": "cards",
                "title": "Konflikt, Gewalt und Deeskalation",
                "subtitle": "BHF 23 ist hochrelevant für mündliche Prüfungen.",
                "cards": [
                    {
                        "title": "Gewalt erkennen",
                        "tone": "sky",
                        "body": [
                            "physisch, psychisch, strukturell, sexualisiert",
                            "auch Überforderung und Demütigung mitdenken",
                        ],
                    },
                    {
                        "title": "Konflikte analysieren",
                        "tone": "mint",
                        "body": [
                            "Worum geht es wirklich?",
                            "Interessen, Gefühle, Auslöser, Rollen klären",
                        ],
                    },
                    {
                        "title": "Deeskalation",
                        "tone": "sand",
                        "body": [
                            "ruhige Stimme, Abstand, klare Sätze",
                            "nicht provozieren, Sicherheit sichern, Hilfe holen",
                        ],
                    },
                    {
                        "title": "Prävention",
                        "tone": "rose",
                        "body": [
                            "Teamkultur, Reflexion, Fallbesprechung, Selbstschutz",
                            "Dokumentation und Meldesysteme nutzen",
                        ],
                    },
                ],
            },
            {
                "type": "cards",
                "title": "SIS und Falltraining aus BHF 24",
                "subtitle": "Die gleiche Struktur hilft fast in jedem Prüfungsfall.",
                "cards": [
                    {
                        "title": "1. Lage erfassen",
                        "tone": "sky",
                        "body": [
                            "Was ist akut? Was ist gefährlich? Was belastet?",
                            "Diagnosen, Symptome, Ressourcen, Umfeld sammeln",
                        ],
                    },
                    {
                        "title": "2. Priorisieren",
                        "tone": "mint",
                        "body": [
                            "Atemweg, Kreislauf, Schmerz, Sicherheit",
                            "danach Selbstversorgung, Mobilität, Ausscheidung, Psyche",
                        ],
                    },
                    {
                        "title": "3. Pflege planen",
                        "tone": "sand",
                        "body": [
                            "Problem + Ziel + Maßnahme + Evaluation",
                            "interprofessionelle Zusammenarbeit nennen",
                        ],
                    },
                    {
                        "title": "4. Beispiele aus BHF 24",
                        "tone": "rose",
                        "body": [
                            "Frau Weiß, Frau Feldmann, Herr Gärtner, Lara, Leon",
                            "SIS und Fallübersicht als Prüfungsroutine nutzen",
                        ],
                    },
                ],
            },
            {
                "type": "checklist",
                "title": "Prüfungscheck BHF 20-24",
                "left_title": "Das solltest du frei erklären können",
                "left_lines": [
                    "Formen der Harninkontinenz",
                    "Qualitätssicherung einfach und praxisnah",
                    "Frühgeborenes / kindliche Pflege mit Elternfokus",
                    "Gewalt und Deeskalation in Pflegesituationen",
                    "SIS-Logik und Fallanalyse in 4 Schritten",
                ],
                "right_title": "Merkhilfen für die Klausur",
                "right_lines": [
                    "Qualität = Standard + Beobachtung + Evaluation",
                    "Bei Kindern immer Elternarbeit einbauen",
                    "Bei Konflikt zuerst Sicherheit, dann Kommunikation",
                    "Bei jedem Fall: akut, Risiko, Ressource, Ziel, Maßnahme",
                ],
            },
        ],
    },
]


def set_fill(shape, color, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1.1, transparency: float = 0.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_box(
    slide,
    x,
    y,
    w,
    h,
    text="",
    font_size=18,
    color=TEXT,
    bold=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_lines(
    slide,
    x,
    y,
    w,
    h,
    lines,
    font_size=15,
    color=TEXT,
    bold_first=False,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
    line_spacing=1.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(5)
    tf.margin_right = Pt(5)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold_first and idx == 0
        run.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, title, body_lines, fill_color, accent):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(card, fill_color)
    set_line(card, accent, 1.0, 0.2)
    add_box(
        slide,
        x + 0.14,
        y + 0.12,
        w - 0.28,
        0.36,
        title,
        18,
        accent,
        True,
        font_name="Aptos Display",
    )
    add_lines(slide, x + 0.14, y + 0.52, w - 0.28, h - 0.66, body_lines, 14, TEXT)


def add_background(slide, accent, accent_dark):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_fill(bg, WHITE)
    set_line(bg, WHITE, 0)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.22))
    set_fill(bar, accent)
    set_line(bar, accent, 0)
    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(7.18), Inches(13.333), Inches(0.32))
    set_fill(ribbon, SOFT_BLUE)
    set_line(ribbon, SOFT_BLUE, 0)
    for x, y, size, color, tr in [
        (11.55, 0.45, 0.62, accent, 0.78),
        (12.1, 0.92, 0.42, accent_dark, 0.84),
        (10.85, 1.05, 0.32, accent, 0.86),
    ]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        set_fill(circle, color, tr)
        set_line(circle, color, 0)


def add_footer(slide, left_text, slide_no, total, accent):
    add_box(slide, 0.55, 7.18, 10.7, 0.16, left_text, 9, MUTED)
    add_box(
        slide,
        12.0,
        7.14,
        0.8,
        0.18,
        f"{slide_no}/{total}",
        10,
        accent,
        True,
        font_name="Aptos Display",
        align=PP_ALIGN.RIGHT,
    )


def add_tag(slide, text, accent):
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.58),
        Inches(0.46),
        Inches(1.9),
        Inches(0.34),
    )
    set_fill(tag, accent)
    set_line(tag, accent, 0.8)
    add_box(
        slide,
        0.64,
        0.49,
        1.78,
        0.26,
        text,
        11,
        WHITE,
        True,
        font_name="Aptos",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_title_slide(prs, deck, slide_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.78), Inches(6.3), Inches(5.7)
    )
    set_fill(panel, deck["accent"])
    set_line(panel, deck["accent"], 0.8)
    halo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(4.6), Inches(1.8), Inches(1.8))
    set_fill(halo, deck["accent_dark"], 0.72)
    set_line(halo, deck["accent_dark"], 0.8, 0.85)
    add_box(slide, 0.95, 1.02, 5.35, 0.34, deck["bhf"], 13, WHITE, True, font_name="Aptos")
    add_box(
        slide,
        0.95,
        1.42,
        5.45,
        1.8,
        deck["title"],
        28,
        WHITE,
        True,
        font_name="Aptos Display",
    )
    add_lines(slide, 0.95, 3.45, 5.4, 1.35, [deck["subtitle"], deck["hero"]], 16, SOFT_BLUE)
    add_card(
        slide,
        7.35,
        1.0,
        5.2,
        1.45,
        "Lernziel in 2 Wochen",
        [
            "grobe Themen beherrschen",
            "Krankheitslehre mit DURST strukturieren",
            "Pflege immer mit Beobachtung + Begründung erklären",
        ],
        CARD_TONES["sky"],
        deck["accent_dark"],
    )
    add_card(
        slide,
        7.35,
        2.7,
        5.2,
        1.45,
        "Prüfungsstrategie",
        [
            "Definition zuerst",
            "dann Symptome / Risiken",
            "dann konkrete Pflege und Prioritäten",
        ],
        CARD_TONES["sand"],
        deck["accent_dark"],
    )
    add_card(
        slide,
        7.35,
        4.4,
        5.2,
        1.45,
        "Merksatz",
        [
            "Akut geht vor Detail.",
            "Pflege ist immer Beobachtung, Sicherheit, Kommunikation und Evaluation.",
        ],
        CARD_TONES["mint"],
        deck["accent_dark"],
    )
    add_footer(slide, "Exam Prep Deck | erstellt aus den BHF-Themen des OneNote", slide_no, total, deck["accent"])


def layout_positions(count: int):
    if count == 3:
        return [
            (0.72, 1.85, 3.95, 2.0),
            (4.88, 1.85, 3.95, 2.0),
            (9.04, 1.85, 3.55, 2.0),
        ]
    return [
        (0.72, 1.72, 5.85, 2.15),
        (6.78, 1.72, 5.85, 2.15),
        (0.72, 4.02, 5.85, 2.15),
        (6.78, 4.02, 5.85, 2.15),
    ]


def add_cards_slide(prs, deck, payload, slide_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, deck["bhf"], deck["accent"])
    add_box(slide, 0.72, 0.95, 9.2, 0.48, payload["title"], 24, INK, True, font_name="Aptos Display")
    add_box(slide, 0.75, 1.34, 9.8, 0.28, payload["subtitle"], 12, MUTED)
    for (x, y, w, h), card in zip(layout_positions(len(payload["cards"])), payload["cards"]):
        add_card(slide, x, y, w, h, card["title"], card["body"], CARD_TONES[card["tone"]], deck["accent_dark"])
    add_footer(slide, "Große Themen zuerst lernen, Details danach verankern", slide_no, total, deck["accent"])


def add_durst_slide(prs, deck, payload, slide_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, "DURST", deck["accent"])
    add_box(slide, 0.72, 0.88, 8.6, 0.42, payload["title"], 24, INK, True, font_name="Aptos Display")
    add_box(slide, 0.75, 1.28, 10.8, 0.24, payload["hook"], 12, MUTED)
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.3), Inches(0.82), Inches(2.2), Inches(0.46)
    )
    set_fill(pill, CARD_TONES["rose"])
    set_line(pill, deck["accent_dark"], 0.9, 0.4)
    add_box(
        slide,
        10.42,
        0.89,
        1.95,
        0.3,
        payload["term"],
        14,
        deck["accent_dark"],
        True,
        font_name="Aptos Display",
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    positions = [
        (0.72, 1.75, 4.0, 1.55),
        (4.88, 1.75, 4.0, 1.55),
        (9.04, 1.75, 3.58, 1.55),
        (0.72, 3.52, 4.0, 1.55),
        (4.88, 3.52, 4.0, 1.55),
        (9.04, 3.52, 3.58, 1.55),
    ]
    tones = ["sky", "mint", "sand", "rose", "white", "mint"]
    for (x, y, w, h), tone, (title, body) in zip(positions, tones, payload["cards"]):
        add_card(slide, x, y, w, h, title, body, CARD_TONES[tone], deck["accent_dark"])
    exam_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(5.45), Inches(11.9), Inches(1.08)
    )
    set_fill(exam_box, SOFT_MINT)
    set_line(exam_box, deck["accent"], 1.0, 0.35)
    add_box(slide, 0.9, 5.6, 1.4, 0.2, "Prüfungsfokus", 13, deck["accent_dark"], True, font_name="Aptos Display")
    add_lines(slide, 0.9, 5.86, 11.3, 0.42, [payload["exam"]], 14, TEXT)
    add_footer(slide, "DURST = Definition, Ursachen, Risiken, Symptome, Therapie, Pflege", slide_no, total, deck["accent"])


def add_checklist_slide(prs, deck, payload, slide_no, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, deck["accent"], deck["accent_dark"])
    add_tag(slide, "Prüfungscheck", deck["accent"])
    add_box(slide, 0.72, 0.92, 9.0, 0.46, payload["title"], 24, INK, True, font_name="Aptos Display")
    left = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.72), Inches(5.92), Inches(4.92)
    )
    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.82), Inches(1.72), Inches(5.8), Inches(4.92)
    )
    set_fill(left, CARD_TONES["sky"])
    set_fill(right, CARD_TONES["sand"])
    set_line(left, deck["accent_dark"], 1.0, 0.25)
    set_line(right, deck["accent_dark"], 1.0, 0.25)
    add_box(slide, 0.94, 1.92, 5.4, 0.28, payload["left_title"], 17, deck["accent_dark"], True, font_name="Aptos Display")
    add_lines(slide, 0.94, 2.28, 5.4, 3.95, [f"• {line}" for line in payload["left_lines"]], 15, TEXT)
    add_box(slide, 7.04, 1.92, 5.24, 0.28, payload["right_title"], 17, deck["accent_dark"], True, font_name="Aptos Display")
    add_lines(slide, 7.04, 2.28, 5.24, 3.95, [f"• {line}" for line in payload["right_lines"]], 15, TEXT)
    add_footer(slide, "Die sicherste Prüfungsantwort ist immer klar, priorisiert und pflegebezogen.", slide_no, total, deck["accent"])


def build_deck(deck):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = deck["title"]
    prs.core_properties.subject = deck["subtitle"]
    prs.core_properties.author = "Codex"
    prs.core_properties.comments = "Exam prep Deck aus den OneNote-BHF-Themen"
    total = len(deck["slides"]) + 1
    add_title_slide(prs, deck, 1, total)
    for idx, slide in enumerate(deck["slides"], start=2):
        if slide["type"] == "cards":
            add_cards_slide(prs, deck, slide, idx, total)
        elif slide["type"] == "durst":
            add_durst_slide(prs, deck, slide, idx, total)
        elif slide["type"] == "checklist":
            add_checklist_slide(prs, deck, slide, idx, total)
        else:
            raise ValueError(f"Unknown slide type: {slide['type']}")
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUT_DIR / deck["filename"]
    prs.save(str(out_path))
    return out_path


def build_all():
    paths = [build_deck(deck) for deck in DECKS]
    for path in paths:
        print(f"saved {path}")


if __name__ == "__main__":
    build_all()
